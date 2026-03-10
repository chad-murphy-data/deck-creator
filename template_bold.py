"""
Bold template — confident, geometric, artistic-but-corporate.
Signature: thick color strips, oversized numbers, right-edge color column motif,
dark title slides, warm gray content backgrounds.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

W = Inches(10)
H = Inches(5.625)

# ── Palette ──────────────────────────────────────────────
CHARCOAL   = RGBColor(0x2A, 0x2A, 0x2A)   # title/closer bg
DK_GREEN   = RGBColor(0x04, 0x40, 0x14)   # section divider bg, primary brand
MINT       = RGBColor(0x1D, 0xE4, 0xCA)   # bright accent on dark slides
WARM_GRAY  = RGBColor(0xF2, 0xF0, 0xEC)   # content slide bg
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)    # card backgrounds on gray
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x1A)   # primary text
MID_TEXT    = RGBColor(0x58, 0x58, 0x58)   # secondary text
MUTED       = RGBColor(0x8A, 0x8A, 0x8A)   # captions
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_RULE  = RGBColor(0xDD, 0xDB, 0xD7)   # subtle separators

# Confident accent palette
GREEN   = RGBColor(0x36, 0x87, 0x27)
BLUE    = RGBColor(0x38, 0x80, 0xF3)
PURPLE  = RGBColor(0x5B, 0x2C, 0x8F)
COBALT  = RGBColor(0x04, 0x54, 0x7C)
GOLD    = RGBColor(0xD4, 0xA8, 0x43)
TEAL    = RGBColor(0x1B, 0x7A, 0x6E)

ACCENTS = [DK_GREEN, COBALT, PURPLE, GOLD, TEAL]

SECTION_COLORS = {
    "green": DK_GREEN, "blue": BLUE, "purple": PURPLE,
    "cobalt": COBALT, "gold": GOLD,
    "red": RGBColor(0xC2, 0x3B, 0x22), "teal": TEAL,
    "ochre": RGBColor(0xCC, 0x7A, 0x2E), "slate": RGBColor(0x5A, 0x6A, 0x7A),
    "plum": RGBColor(0x8E, 0x45, 0x85),
}

def _resolve_section_color(c, default=None):
    sc = c.get("sectionColor", "")
    return SECTION_COLORS.get(sc, default or DK_GREEN)

TITLE_FONT = "Fidelity Slab"
TITLE_FONT_FALLBACK = "Trebuchet MS"
BODY_FONT  = "Fidelity Sans"
BODY_FONT_FALLBACK = "Calibri"

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "logo.png")

LM = Inches(0.7)
CW = Inches(8.6)
EDGE_W = Inches(0.35)   # right-edge color column width


# ── Helpers ──────────────────────────────────────────────

def _font_fallback(font_name):
    if font_name == TITLE_FONT: return TITLE_FONT, TITLE_FONT_FALLBACK
    if font_name == BODY_FONT: return BODY_FONT, BODY_FONT_FALLBACK
    return font_name, None

def _set_font_with_fallback(font_obj, font_name):
    primary, fallback = _font_fallback(font_name)
    font_obj.name = primary
    if fallback:
        el = font_obj._element
        latin = el.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(el, qn('a:latin'))
            latin.set('typeface', primary)
        cs = el.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(el, qn('a:cs'))
        cs.set('typeface', fallback)

def _rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

def _oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

def _tb(slide, x, y, w, h, text, font=None, size=12, color=None,
        bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
        line_spacing=None):
    font = font or BODY_FONT; color = color or DARK_TEXT
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    _set_font_with_fallback(p.font, font); p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.italic = italic
    p.alignment = align
    if line_spacing: p.line_spacing = Pt(line_spacing)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', {MSO_ANCHOR.TOP:'t', MSO_ANCHOR.MIDDLE:'ctr', MSO_ANCHOR.BOTTOM:'b'}[valign])
    return box

def _warm_bg(slide):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = WARM_GRAY

def _edge_column(slide, color=DK_GREEN):
    """Right-edge color column — the bold motif."""
    _rect(slide, Emu(W - EDGE_W), Inches(0), EDGE_W, H, color)

def _content_title(slide, title, size=28, edge_color=DK_GREEN):
    """Bold geometric title: thick color bar above, then title text."""
    _tb(slide, LM, Inches(0.42), Inches(8.0), Inches(0.65), title,
        TITLE_FONT, size, DARK_TEXT, bold=True, valign=MSO_ANCHOR.BOTTOM)

CONTENT_TOP = Inches(1.25)


# ── Builders ─────────────────────────────────────────────

def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = CHARCOAL
    # Bold green block — left side
    _rect(slide, Inches(0), Inches(0), Inches(0.5), H, DK_GREEN)
    # Mint accent strip
    _rect(slide, Inches(0.5), Inches(0), Inches(0.08), H, MINT)
    # Title
    _tb(slide, Inches(0.9), Inches(0.8), Inches(8.5), Inches(2.0),
        c.get("title", "Title"), TITLE_FONT, 44, WHITE, bold=True,
        valign=MSO_ANCHOR.BOTTOM)
    # Thick green bar under title
    _rect(slide, Inches(0.9), Inches(3.0), Inches(3.0), Inches(0.08), MINT)
    if c.get("subtitle"):
        _tb(slide, Inches(0.9), Inches(3.25), Inches(8.5), Inches(0.5),
            c["subtitle"], BODY_FONT, 16, RGBColor(0xBB, 0xBB, 0xBB))
    meta = []
    if c.get("author"): meta.append(c["author"])
    if c.get("date"): meta.append(c["date"])
    if meta:
        _tb(slide, Inches(0.9), Inches(4.3), Inches(5), Inches(0.7),
            "\n".join(meta), BODY_FONT, 12, RGBColor(0x99, 0x99, 0x99))
    # Logo in upper-right
    logo = c.get("logo_path", LOGO_PATH)
    if logo and os.path.isfile(logo):
        slide.shapes.add_picture(logo, Inches(7.6), Inches(0.3), Inches(2.0))


def build_in_brief(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "In Brief"))
    bullets = c.get("bullets", [])
    n = len(bullets)
    avail = H - CONTENT_TOP - Inches(0.35)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.95))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, b in enumerate(bullets):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # White card
        _rect(slide, LM, y, Inches(8.2), rowH, CARD_BG)
        # Thick left color strip
        _rect(slide, LM, y, Inches(0.12), rowH, col)
        # Number
        _tb(slide, Emu(LM + Inches(0.25)), y, Inches(0.45), rowH,
            str(i + 1), TITLE_FONT, 22, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        # Text
        _tb(slide, Emu(LM + Inches(0.7)), y, Inches(7.2), rowH,
            b, BODY_FONT, 13, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=18)


def build_section_divider(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    # Mint accent strip at left
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, MINT)
    if c.get("sectionNumber"):
        _tb(slide, Inches(0.6), Inches(0.9), Inches(2), Inches(1.0),
            f"0{c['sectionNumber']}", TITLE_FONT, 56, MINT, bold=True,
            valign=MSO_ANCHOR.BOTTOM)
    _tb(slide, Inches(0.6), Inches(2.0), Inches(8.5), Inches(1.2),
        c.get("title", "Section"), TITLE_FONT, 38, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(0.6), Inches(3.25), Inches(2.5), Inches(0.06), MINT)
    if c.get("subtitle"):
        _tb(slide, Inches(0.6), Inches(3.5), Inches(8.5), Inches(0.5),
            c["subtitle"], BODY_FONT, 14, RGBColor(0xC8, 0xD8, 0xC8))


def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide)
    sc = _resolve_section_color(c)
    _edge_column(slide, sc)
    _content_title(slide, c.get("title", "Key Metric"), size=22, edge_color=sc)
    # Big number on a color block
    _rect(slide, Inches(1.5), Inches(1.4), Inches(6.0), Inches(2.0), sc)
    _tb(slide, Inches(1.5), Inches(1.4), Inches(6.0), Inches(2.0),
        c.get("stat", "—"), TITLE_FONT, 80, WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if c.get("headline"):
        _tb(slide, Inches(1.2), Inches(3.6), Inches(6.6), Inches(0.6),
            c["headline"], BODY_FONT, 16, DARK_TEXT, bold=True,
            align=PP_ALIGN.CENTER)
    if c.get("detail"):
        _tb(slide, Inches(1.5), Inches(4.2), Inches(6.0), Inches(0.7),
            c["detail"], BODY_FONT, 11, MID_TEXT, align=PP_ALIGN.CENTER,
            line_spacing=16)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(5.0), Inches(8.5), Inches(0.3),
            c["source"], BODY_FONT, 9, MUTED, italic=True, align=PP_ALIGN.CENTER)


def build_quote(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide, PURPLE)
    _content_title(slide, c.get("title", "In Their Words"), size=22, edge_color=PURPLE)
    # Bold quote block with thick left border
    _rect(slide, Inches(0.7), Inches(1.3), Inches(0.1), Inches(2.5), PURPLE)
    _tb(slide, Inches(1.1), Inches(1.4), Inches(7.5), Inches(2.2),
        c.get("quote", ""), BODY_FONT, 20, DARK_TEXT, italic=True,
        valign=MSO_ANCHOR.MIDDLE, line_spacing=28)
    # Attribution
    _rect(slide, Inches(1.1), Inches(4.0), Inches(1.5), Inches(0.06), PURPLE)
    if c.get("attribution"):
        _tb(slide, Inches(1.1), Inches(4.15), Inches(7), Inches(0.35),
            c["attribution"], BODY_FONT, 12, MID_TEXT, bold=True)
    if c.get("context"):
        _tb(slide, Inches(1.1), Inches(4.5), Inches(7), Inches(0.35),
            c["context"], BODY_FONT, 10, MUTED, italic=True)


def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Comparison"), size=24)

    cardH = Inches(3.6)
    # Left card
    _rect(slide, LM, CONTENT_TOP, Inches(3.95), cardH, CARD_BG)
    _rect(slide, LM, CONTENT_TOP, Inches(3.95), Inches(0.5), COBALT)
    _tb(slide, Emu(LM + Inches(0.15)), CONTENT_TOP, Inches(3.6), Inches(0.5),
        c.get("leftLabel", "Before"), TITLE_FONT, 14, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(c.get("leftItems", [])):
        _tb(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.65) + i * Inches(0.6)),
            Inches(3.5), Inches(0.55), item, BODY_FONT, 12, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)

    # Right card
    _rect(slide, Inches(4.9), CONTENT_TOP, Inches(4.1), cardH, CARD_BG)
    _rect(slide, Inches(4.9), CONTENT_TOP, Inches(4.1), Inches(0.5), DK_GREEN)
    _tb(slide, Inches(5.05), CONTENT_TOP, Inches(3.8), Inches(0.5),
        c.get("rightLabel", "After"), TITLE_FONT, 14, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(c.get("rightItems", [])):
        _tb(slide, Inches(5.1), Emu(CONTENT_TOP + Inches(0.65) + i * Inches(0.6)),
            Inches(3.7), Inches(0.55), item, BODY_FONT, 12, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)


def build_text_graph(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Title"), size=24)
    texts = c.get("text", [])
    if not isinstance(texts, list): texts = [texts]
    for i, t in enumerate(texts):
        _tb(slide, LM, Emu(CONTENT_TOP + i * Inches(1.2)),
            Inches(3.8), Inches(1.1), t, BODY_FONT, 12, DARK_TEXT,
            line_spacing=17)
    # Thick vertical divider
    _rect(slide, Inches(4.8), CONTENT_TOP, Inches(0.06), Inches(3.5), DK_GREEN)

    from pptx.chart.data import CategoryChartData
    raw = c.get("chartData", [{"name": "S1", "labels": ["A","B","C"], "values": [25,45,30]}])
    cd_obj = CategoryChartData()
    if raw:
        cd = raw[0]
        cd_obj.categories = cd.get("labels", ["A","B","C"])
        cd_obj.add_series(cd.get("name","Series 1"), cd.get("values",[25,45,30]))
    ct = {"line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}.get(
        c.get("chartType","bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(ct, Inches(5.05), CONTENT_TOP, Inches(3.95), Inches(3.6), cd_obj)


def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Process"), size=24)
    steps = c.get("steps", [])
    count = min(len(steps), 5)
    if count == 0: return

    total_w = 8.2
    gap_w = 0.25
    step_w = (total_w - (count - 1) * gap_w) / count
    startX = 0.7

    for i, step in enumerate(steps[:count]):
        x = Inches(startX + i * (step_w + gap_w))
        sw = Inches(step_w)
        col = ACCENTS[i % len(ACCENTS)]

        # White card with thick color top
        _rect(slide, x, CONTENT_TOP, sw, Inches(3.9), CARD_BG)
        _rect(slide, x, CONTENT_TOP, sw, Inches(0.08), col)
        # Oversized number
        _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.15)),
            Inches(0.6), Inches(0.7), str(i + 1),
            TITLE_FONT, 36, col, bold=True)
        # Step title
        _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.9)),
            Emu(sw - Inches(0.24)), Inches(0.6),
            step.get("title", ""), BODY_FONT, 12, DARK_TEXT, bold=True)
        if step.get("detail"):
            _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(1.55)),
                Emu(sw - Inches(0.24)), Inches(2.1),
                step["detail"], BODY_FONT, 10, MID_TEXT, line_spacing=14)
        # Arrow between steps
        if i < count - 1:
            _tb(slide, Emu(x + sw), Emu(CONTENT_TOP + Inches(1.2)),
                Inches(gap_w), Inches(0.5), "\u2192",
                BODY_FONT, 18, MUTED, align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE)


def build_matrix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Framework"), size=24)
    quads = c.get("quadrants", [{}, {}, {}, {}])
    qW = Inches(3.95)
    qH = Inches(1.7)
    gap = Inches(0.25)
    sX = LM
    sY = CONTENT_TOP

    for i, q in enumerate(quads[:4]):
        col_idx = i % 2; row = i // 2
        x = Emu(sX + col_idx * (qW + gap))
        y = Emu(sY + row * (qH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # White card
        _rect(slide, x, y, qW, qH, CARD_BG)
        # Thick left color strip
        _rect(slide, x, y, Inches(0.1), qH, col)
        _tb(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.1)),
            Emu(qW - Inches(0.3)), Inches(0.35),
            q.get("label", ""), TITLE_FONT, 13, col, bold=True)
        _tb(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.5)),
            Emu(qW - Inches(0.3)), Emu(qH - Inches(0.6)),
            q.get("detail", ""), BODY_FONT, 11, MID_TEXT, line_spacing=15)


def build_methods(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Approach"), size=24)
    fields = c.get("fields", [])
    n = len(fields)
    avail = H - CONTENT_TOP - Inches(0.35)
    rowH = min(int((avail - (n - 1) * Inches(0.06)) / n), Inches(0.8))
    gap = Inches(0.06)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, f in enumerate(fields):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        bg = CARD_BG if i % 2 == 0 else WARM_GRAY
        _rect(slide, LM, y, Inches(8.2), rowH, bg)
        _rect(slide, LM, y, Inches(0.08), rowH, col)
        # Label
        _tb(slide, Emu(LM + Inches(0.2)), y, Inches(1.8), rowH,
            f.get("label", ""), TITLE_FONT, 11, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        # Value
        _tb(slide, Emu(LM + Inches(2.1)), y, Inches(5.9), rowH,
            f.get("value", ""), BODY_FONT, 12, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


def build_hypotheses(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Hypotheses"), size=24)
    hyps = c.get("hypotheses", [])
    n = len(hyps)
    avail = H - CONTENT_TOP - Inches(0.35)
    rowH = min(int((avail - (n - 1) * Inches(0.08)) / n), Inches(0.75))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, h in enumerate(hyps):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, Inches(8.2), rowH, CARD_BG)
        _rect(slide, LM, y, Inches(0.1), rowH, col)
        # H-number in circle
        cs = Inches(0.36)
        cx = Emu(LM + Inches(0.2)); cy = Emu(y + (rowH - cs) // 2)
        _oval(slide, cx, cy, cs, cs, col)
        _tb(slide, cx, cy, cs, cs, f"H{i+1}", TITLE_FONT, 11, WHITE,
            bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(LM + Inches(0.7)), y, Inches(5.8), rowH,
            h.get("text", ""), BODY_FONT, 12, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)
        if h.get("status"):
            sc = DK_GREEN if h["status"] == "Confirmed" else (RGBColor(0xC2,0x3B,0x22) if h["status"] == "Rejected" else MUTED)
            _rect(slide, Inches(7.5), Emu(y + (rowH - Inches(0.3)) // 2),
                  Inches(1.2), Inches(0.3), sc)
            _tb(slide, Inches(7.5), Emu(y + (rowH - Inches(0.3)) // 2),
                Inches(1.2), Inches(0.3), h["status"],
                BODY_FONT, 9, WHITE, bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE)


def build_wsn_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide, PURPLE)
    _content_title(slide, c.get("title", "Key Finding"), size=24)
    cols_data = [
        ("What", DK_GREEN, c.get("what", {})),
        ("So What", COBALT, c.get("soWhat", {})),
        ("Now What", PURPLE, c.get("nowWhat", {})),
    ]
    colW = Inches(2.6)
    gap_w = Inches(0.2)
    cardH = Inches(3.8)

    for i, (label, color, data) in enumerate(cols_data):
        x = Emu(LM + i * (colW + gap_w))
        _rect(slide, x, CONTENT_TOP, colW, cardH, CARD_BG)
        # Color header bar
        _rect(slide, x, CONTENT_TOP, colW, Inches(0.45), color)
        _tb(slide, Emu(x + Inches(0.12)), CONTENT_TOP, Emu(colW - Inches(0.24)), Inches(0.45),
            label, TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.55)),
            Emu(colW - Inches(0.24)), Inches(1.1),
            data.get("headline", ""), BODY_FONT, 12, DARK_TEXT, bold=True,
            line_spacing=17)
        if data.get("detail"):
            _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(1.7)),
                Emu(colW - Inches(0.24)), Inches(1.9),
                data["detail"], BODY_FONT, 10, MID_TEXT, line_spacing=14)


def build_wsn_reveal(prs, c):
    """Progressive 3-slide reveal."""
    WSN_COLORS = [DK_GREEN, COBALT, PURPLE]
    WSN_LABELS = ["What We Found", "So What", "Now What"]

    def _hdr(slide):
        _warm_bg(slide); _edge_column(slide, PURPLE)
        _rect(slide, LM, Inches(0.3), Inches(2.0), Inches(0.08), DK_GREEN)
        _tb(slide, LM, Inches(0.42), Inches(8.0), Inches(0.55),
            c.get("title", "Key Finding"), TITLE_FONT, 26, DARK_TEXT, bold=True,
            valign=MSO_ANCHOR.BOTTOM)

    def _zone(slide, x, w, label, color, data, condensed=False):
        y = Inches(1.2)
        h = Inches(2.0) if condensed else Inches(3.7)
        _rect(slide, x, y, w, h, CARD_BG)
        _rect(slide, x, y, w, Inches(0.4), color)
        _tb(slide, Emu(x + Inches(0.12)), y, Emu(w - Inches(0.24)), Inches(0.4),
            label, TITLE_FONT, 11 if condensed else 12, WHITE, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        fs = 11 if condensed else 13
        _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.5)),
            Emu(w - Inches(0.24)), Inches(0.6) if condensed else Inches(0.9),
            data.get("headline", ""), BODY_FONT, fs, DARK_TEXT, bold=True,
            line_spacing=16)
        if data.get("detail"):
            _tb(slide, Emu(x + Inches(0.12)),
                Emu(y + Inches(1.1)) if condensed else Emu(y + Inches(1.45)),
                Emu(w - Inches(0.24)), Inches(0.65) if condensed else Inches(1.9),
                data["detail"], BODY_FONT, 9 if condensed else 11, MID_TEXT,
                line_spacing=13 if condensed else 15)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6]); _hdr(s1)
    _zone(s1, LM, Inches(5.5), "What We Found", DK_GREEN, c.get("what", {}))

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); _hdr(s2)
    _zone(s2, LM, Inches(3.95), "What We Found", DK_GREEN, c.get("what", {}))
    _zone(s2, Inches(4.9), Inches(4.1), "So What", COBALT, c.get("soWhat", {}))

    # Slide 3: condensed + Now What
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); _hdr(s3)
    _zone(s3, LM, Inches(3.95), "What We Found", DK_GREEN, c.get("what", {}), True)
    _zone(s3, Inches(4.9), Inches(4.1), "So What", COBALT, c.get("soWhat", {}), True)
    nwY = Inches(3.4); nwH = Inches(1.85)
    _rect(s3, LM, nwY, Inches(8.2), nwH, CARD_BG)
    _rect(s3, LM, nwY, Inches(8.2), Inches(0.45), PURPLE)
    _tb(s3, Emu(LM + Inches(0.15)), nwY, Inches(3), Inches(0.45),
        "Now What", TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    nw = c.get("nowWhat", {})
    _tb(s3, Emu(LM + Inches(0.15)), Emu(nwY + Inches(0.55)),
        Inches(7.8), Inches(0.55),
        nw.get("headline", ""), BODY_FONT, 15, DARK_TEXT, bold=True)
    if nw.get("detail"):
        _tb(s3, Emu(LM + Inches(0.15)), Emu(nwY + Inches(1.15)),
            Inches(7.8), Inches(0.55), nw["detail"], BODY_FONT, 11, MID_TEXT)


def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Findings & Recommendations"), size=22)
    items = c.get("items", [])
    n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.82))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, item in enumerate(items[:5]):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Finding card
        _rect(slide, LM, y, Inches(3.6), rowH, CARD_BG)
        _rect(slide, LM, y, Inches(0.08), rowH, col)
        _tb(slide, Emu(LM + Inches(0.18)), y, Inches(3.3), rowH,
            item.get("finding", ""), BODY_FONT, 11, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)
        # Arrow circle
        acs = Inches(0.32)
        acx = Emu(Inches(4.52)); acy = Emu(y + (rowH - acs) // 2)
        _oval(slide, acx, acy, acs, acs, col)
        _tb(slide, acx, acy, acs, acs, "\u2192",
            BODY_FONT, 14, WHITE, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE)
        # Rec card
        _rect(slide, Inches(5.0), y, Inches(3.9), rowH, CARD_BG)
        _rect(slide, Inches(5.0), y, Inches(0.08), rowH, col)
        _tb(slide, Inches(5.2), y, Inches(3.5), rowH,
            item.get("recommendation", ""), BODY_FONT, 11, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Complete Findings"), size=22)
    items = c.get("items", [])
    n = min(len(items), 8)
    avail = H - CONTENT_TOP - Inches(0.2)
    rowH = min(int((avail - (n - 1) * Inches(0.04)) / n), Inches(0.48))
    gap = Inches(0.04)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, item in enumerate(items[:8]):
        y = Emu(startY + i * (rowH + gap))
        bg = CARD_BG if i % 2 == 0 else WARM_GRAY
        _rect(slide, LM, y, Inches(3.9), rowH, bg)
        _rect(slide, LM, y, Inches(0.06), rowH, DK_GREEN)
        _tb(slide, Emu(LM + Inches(0.15)), y, Inches(3.65), rowH,
            item.get("finding", ""), BODY_FONT, 9, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)
        _rect(slide, Inches(4.75), y, Inches(4.15), rowH, bg)
        _rect(slide, Inches(4.75), y, Inches(0.06), rowH, COBALT)
        _tb(slide, Inches(4.9), y, Inches(3.85), rowH,
            item.get("recommendation", ""), BODY_FONT, 9, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE)


def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide, GOLD)
    _content_title(slide, c.get("title", "Open Questions"), size=24, edge_color=GOLD)
    questions = c.get("questions", [])
    cardW = Inches(3.95)
    cardH = Inches(1.7)
    gX = Inches(0.25)
    gY = Inches(0.25)

    for i, question in enumerate(questions[:4]):
        col_idx = i % 2; row = i // 2
        x = Emu(LM + col_idx * (cardW + gX))
        y = Emu(CONTENT_TOP + row * (cardH + gY))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, y, cardW, cardH, CARD_BG)
        _rect(slide, x, y, Inches(0.1), cardH, col)
        # Bold number
        _tb(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.1)),
            Inches(0.5), Inches(0.5), str(i + 1),
            TITLE_FONT, 26, col, bold=True)
        _tb(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.6)),
            Emu(cardW - Inches(0.4)), Inches(0.95),
            question, BODY_FONT, 12, DARK_TEXT, line_spacing=17)


def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Agenda"), size=26)
    items = c.get("items", [])
    n = len(items)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int(avail / n), Inches(0.72))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, item in enumerate(items):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, Inches(8.2), rowH, CARD_BG)
        _rect(slide, LM, y, Inches(0.1), rowH, col)
        # Number circle
        cs = Inches(0.38)
        cx = Emu(LM + Inches(0.2)); cy = Emu(y + (rowH - cs) // 2)
        _oval(slide, cx, cy, cs, cs, col)
        _tb(slide, cx, cy, cs, cs, str(i + 1), TITLE_FONT, 14, WHITE,
            bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        title_text = item if isinstance(item, str) else item.get("title", "")
        _tb(slide, Emu(LM + Inches(0.75)), y, Inches(5.2), rowH,
            title_text, BODY_FONT, 14, DARK_TEXT, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        if isinstance(item, dict) and item.get("detail"):
            _tb(slide, Inches(7.0), y, Inches(1.7), rowH,
                item["detail"], BODY_FONT, 11, MUTED,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_progressive_reveal(prs, c):
    takeaways = c.get("takeaways", [])
    for n in range(min(len(takeaways), 5)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _warm_bg(slide); _edge_column(slide)
        _rect(slide, LM, Inches(0.3), Inches(2.0), Inches(0.08), DK_GREEN)
        _tb(slide, LM, Inches(0.42), Inches(8.0), Inches(0.55),
            c.get("title", "Building the Picture"), TITLE_FONT, 24, DARK_TEXT,
            bold=True, valign=MSO_ANCHOR.BOTTOM)

        cur = takeaways[n]
        col = ACCENTS[n % len(ACCENTS)]
        # Main content card
        _rect(slide, LM, Inches(1.1), Inches(8.2), Inches(2.3), CARD_BG)
        _rect(slide, LM, Inches(1.1), Inches(0.08), Inches(2.3), col)
        _tb(slide, Emu(LM + Inches(0.2)), Inches(1.2), Inches(7.8), Inches(0.55),
            cur.get("headline", ""), BODY_FONT, 15, DARK_TEXT, bold=True)
        if cur.get("detail"):
            _tb(slide, Emu(LM + Inches(0.2)), Inches(1.8), Inches(7.8), Inches(1.3),
                cur["detail"], BODY_FONT, 11, MID_TEXT, line_spacing=15)

        # Separator
        _rect(slide, Inches(0), Inches(3.6), Emu(W - EDGE_W), Inches(0.06), DK_GREEN)
        _tb(slide, LM, Inches(3.72), Inches(3), Inches(0.22),
            "Running Takeaways", TITLE_FONT, 9, DK_GREEN, bold=True)

        for j in range(n + 1):
            ty = Emu(Inches(4.0) + j * Inches(0.38))
            active = j == n
            jcol = ACCENTS[j % len(ACCENTS)]
            _rect(slide, LM, Emu(ty + Inches(0.06)), Inches(0.12), Inches(0.12), jcol)
            _tb(slide, Emu(LM + Inches(0.22)), ty, Inches(7.8), Inches(0.32),
                takeaways[j].get("summary", takeaways[j].get("headline", "")),
                BODY_FONT, 10, DARK_TEXT if active else MUTED,
                bold=active, valign=MSO_ANCHOR.MIDDLE)


def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    _rect(slide, Inches(0), Inches(0), Inches(0.5), H, CHARCOAL)
    _rect(slide, Inches(0.5), Inches(0), Inches(0.08), H, MINT)
    _tb(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.3),
        c.get("title", "Thank You"), TITLE_FONT, 46, WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
    _rect(slide, Inches(3.5), Inches(2.75), Inches(3.0), Inches(0.06), MINT)
    if c.get("subtitle"):
        _tb(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(0.5),
            c["subtitle"], BODY_FONT, 16, WHITE, align=PP_ALIGN.CENTER)
    if c.get("contact"):
        _tb(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
            c["contact"], BODY_FONT, 12, RGBColor(0xA0, 0xC0, 0xA0),
            align=PP_ALIGN.CENTER)




def build_timeline(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Timeline"))
    milestones = c.get("milestones", [])
    n = min(len(milestones), 6)
    if n == 0: return
    lineY = Inches(2.8)
    startX = 0.7; endX = 8.6
    step = (endX - startX) / max(n - 1, 1)
    _rect(slide, Inches(startX), lineY, Inches(endX - startX), Inches(0.04), LIGHT_RULE)
    STATUS_CLR = {"complete": DK_GREEN, "current": GOLD, "upcoming": MUTED}
    for i, m in enumerate(milestones[:n]):
        cx = Inches(startX + i * step)
        col = STATUS_CLR.get(m.get("status", "upcoming"), MUTED)
        ds = Inches(0.22)
        _oval(slide, Emu(cx - ds // 2), Emu(lineY - ds // 2), ds, ds, col)
        _tb(slide, Emu(cx - Inches(0.6)), Emu(lineY - Inches(0.65)),
            Inches(1.2), Inches(0.3), m.get("date", ""),
            BODY_FONT, 10, col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.25)),
            Inches(1.4), Inches(0.4), m.get("title", ""),
            BODY_FONT, 11, DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        if m.get("detail"):
            _tb(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.65)),
                Inches(1.4), Inches(0.7), m["detail"],
                BODY_FONT, 9, MID_TEXT, align=PP_ALIGN.CENTER, line_spacing=12)


def build_data_table(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Data"), size=24)
    headers = c.get("headers", [])
    rows = c.get("rows", [])
    hc = c.get("highlightCol", None)
    nCols = len(headers); nRows = min(len(rows), 10)
    if nCols == 0: return
    tW = 8.2; tX = 0.7; colW = tW / nCols; rowH = min(0.42, 3.8 / max(nRows + 1, 1)); hdrY = 1.25
    _rect(slide, Inches(tX), Inches(hdrY), Inches(tW), Inches(rowH), DK_GREEN)
    for j, h in enumerate(headers):
        _tb(slide, Inches(tX + j * colW + 0.08), Inches(hdrY),
            Inches(colW - 0.16), Inches(rowH), h,
            BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, row in enumerate(rows[:nRows]):
        ry = hdrY + (i + 1) * rowH
        bg = CARD_BG if i % 2 == 0 else WARM_GRAY
        _rect(slide, Inches(tX), Inches(ry), Inches(tW), Inches(rowH), bg)
        for j, cell in enumerate(row[:nCols]):
            is_hl = hc is not None and j == hc
            _tb(slide, Inches(tX + j * colW + 0.08), Inches(ry),
                Inches(colW - 0.16), Inches(rowH), str(cell),
                BODY_FONT, 10, DK_GREEN if is_hl else DARK_TEXT,
                bold=is_hl, valign=MSO_ANCHOR.MIDDLE)
    if c.get("note"):
        _tb(slide, Inches(tX), Inches(hdrY + (nRows + 1) * rowH + 0.1),
            Inches(tW), Inches(0.3), c["note"], BODY_FONT, 8, MUTED, italic=True)


def build_multi_stat(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Key Metrics"), size=24)
    stats = c.get("stats", []); n = min(len(stats), 4)
    if n == 0: return
    totalW = 8.2; gap = 0.25; statW = (totalW - (n - 1) * gap) / n
    for i, s in enumerate(stats[:n]):
        x = Inches(0.7 + i * (statW + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, CONTENT_TOP, Inches(statW), Inches(3.5), CARD_BG)
        _rect(slide, x, CONTENT_TOP, Inches(statW), Inches(0.08), col)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(0.2)), Inches(statW), Inches(1.4),
            s.get("value", "—"), TITLE_FONT, 48, col, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(1.7)), Inches(statW), Inches(0.45),
            s.get("label", ""), BODY_FONT, 12, DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        if s.get("detail"):
            _tb(slide, x, Emu(CONTENT_TOP + Inches(2.2)), Inches(statW), Inches(0.9),
                s["detail"], BODY_FONT, 10, MID_TEXT, align=PP_ALIGN.CENTER, line_spacing=14)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(5.0), Inches(8.5), Inches(0.3),
            c["source"], BODY_FONT, 9, MUTED, italic=True, align=PP_ALIGN.CENTER)


def build_persona(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Persona"), size=22)
    _rect(slide, LM, CONTENT_TOP, Inches(3.8), Inches(3.9), CARD_BG)
    _rect(slide, LM, CONTENT_TOP, Inches(0.1), Inches(3.9), DK_GREEN)
    _tb(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)), Inches(3.4), Inches(0.5),
        c.get("name", ""), TITLE_FONT, 24, DK_GREEN, bold=True)
    _tb(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.6)), Inches(3.4), Inches(0.3),
        c.get("archetype", ""), BODY_FONT, 12, GOLD, bold=True)
    traits = c.get("traits", [])
    for i, t in enumerate(traits[:5]):
        ty = Emu(CONTENT_TOP + Inches(1.1) + i * Inches(0.4))
        _oval(slide, Emu(LM + Inches(0.2)), Emu(ty + Inches(0.04)), Inches(0.12), Inches(0.12), ACCENTS[i % len(ACCENTS)])
        _tb(slide, Emu(LM + Inches(0.45)), Emu(ty), Inches(3.2), Inches(0.35), t, BODY_FONT, 11, DARK_TEXT)
    _rect(slide, Inches(4.75), CONTENT_TOP, Inches(4.15), Inches(3.9), CARD_BG)
    _rect(slide, Inches(4.75), CONTENT_TOP, Inches(0.1), Inches(3.9), GOLD)
    _tb(slide, Inches(4.95), Emu(CONTENT_TOP + Inches(0.1)), Inches(3.8), Inches(0.25),
        "STRATEGY", BODY_FONT, 9, MUTED, bold=True)
    _tb(slide, Inches(4.95), Emu(CONTENT_TOP + Inches(0.4)), Inches(3.8), Inches(1.5),
        c.get("strategy", ""), BODY_FONT, 12, DARK_TEXT, line_spacing=17)
    if c.get("detail"):
        _tb(slide, Inches(4.95), Emu(CONTENT_TOP + Inches(2.0)), Inches(3.8), Inches(1.5),
            c["detail"], BODY_FONT, 10, MID_TEXT, line_spacing=14)


def build_risk_tradeoff(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Risk & Reward"), size=24)
    risks = c.get("risks", []); rewards = c.get("rewards", [])
    SEV = {"high": RGBColor(0xC2, 0x3B, 0x22), "medium": GOLD, "low": DK_GREEN}
    _rect(slide, LM, CONTENT_TOP, Inches(3.8), Inches(0.4), RGBColor(0xC2, 0x3B, 0x22))
    _tb(slide, Emu(LM + Inches(0.1)), CONTENT_TOP, Inches(3.6), Inches(0.4),
        "RISKS", BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(risks[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.5) + i * Inches(0.65))
        sc = SEV.get(r.get("severity", "medium"), GOLD)
        _rect(slide, LM, ry, Inches(3.8), Inches(0.55), CARD_BG)
        _rect(slide, LM, ry, Inches(0.08), Inches(0.55), sc)
        _tb(slide, Emu(LM + Inches(0.18)), ry, Inches(3.4), Inches(0.55),
            r.get("label", "") + (" — " + r.get("detail", "") if r.get("detail") else ""),
            BODY_FONT, 10, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(4.75), CONTENT_TOP, Inches(4.15), Inches(0.4), DK_GREEN)
    _tb(slide, Inches(4.85), CONTENT_TOP, Inches(3.95), Inches(0.4),
        "REWARDS", BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(rewards[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.5) + i * Inches(0.65))
        _rect(slide, Inches(4.75), ry, Inches(4.15), Inches(0.55), CARD_BG)
        _rect(slide, Inches(4.75), ry, Inches(0.08), Inches(0.55), DK_GREEN)
        _tb(slide, Inches(4.95), ry, Inches(3.8), Inches(0.55),
            r.get("label", "") + (" — " + r.get("detail", "") if r.get("detail") else ""),
            BODY_FONT, 10, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE)


def build_appendix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide)
    _rect(slide, LM, Inches(0.25), Inches(2.0), Inches(0.06), MUTED)
    _tb(slide, LM, Inches(0.35), CW, Inches(0.4),
        c.get("title", "Appendix"), TITLE_FONT, 16, MUTED, bold=True)
    sections = c.get("sections", []); y_cursor = 0.85
    for s in sections:
        if y_cursor > 4.8: break
        _tb(slide, LM, Inches(y_cursor), CW, Inches(0.22),
            s.get("label", ""), BODY_FONT, 9, DK_GREEN, bold=True)
        _tb(slide, LM, Inches(y_cursor + 0.22), CW, Inches(0.75),
            s.get("content", ""), BODY_FONT, 8, MID_TEXT, line_spacing=11)
        y_cursor += 1.05


def build_before_after(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Transformation"), size=24)
    bef = c.get("before", {}); aft = c.get("after", {}); interv = c.get("intervention", "")
    _rect(slide, LM, CONTENT_TOP, Inches(3.3), Inches(3.5), CARD_BG)
    _rect(slide, LM, CONTENT_TOP, Inches(3.3), Inches(0.45), COBALT)
    _tb(slide, Emu(LM + Inches(0.12)), CONTENT_TOP, Inches(3.0), Inches(0.45),
        bef.get("label", "Before"), TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Emu(LM + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.55)), Inches(3.0), Inches(2.7),
        bef.get("detail", ""), BODY_FONT, 11, DARK_TEXT, line_spacing=16)
    _tb(slide, Inches(4.15), Emu(CONTENT_TOP + Inches(1.0)), Inches(0.7), Inches(0.5),
        "\u2192", TITLE_FONT, 28, GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(4.0), Emu(CONTENT_TOP + Inches(1.6)), Inches(1.0), Inches(1.5),
        interv, BODY_FONT, 9, MID_TEXT, align=PP_ALIGN.CENTER, line_spacing=13)
    _rect(slide, Inches(5.15), CONTENT_TOP, Inches(3.75), Inches(3.5), CARD_BG)
    _rect(slide, Inches(5.15), CONTENT_TOP, Inches(3.75), Inches(0.45), DK_GREEN)
    _tb(slide, Inches(5.27), CONTENT_TOP, Inches(3.5), Inches(0.45),
        aft.get("label", "After"), TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(5.27), Emu(CONTENT_TOP + Inches(0.55)), Inches(3.5), Inches(2.7),
        aft.get("detail", ""), BODY_FONT, 11, DARK_TEXT, line_spacing=16)


def build_summary(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Summary"), size=26)
    sections = c.get("sections", []); n = len(sections)
    if n == 0: return
    colW = (8.2 - (n - 1) * 0.2) / n
    for i, sec in enumerate(sections[:4]):
        x = Inches(0.7 + i * (colW + 0.2))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, CONTENT_TOP, Inches(colW), Inches(3.8), CARD_BG)
        _rect(slide, x, CONTENT_TOP, Inches(colW), Inches(0.08), col)
        _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.12)), Inches(colW - 0.24), Inches(0.3),
            sec.get("heading", ""), TITLE_FONT, 11, col, bold=True)
        points = sec.get("points", [])
        for j, p in enumerate(points[:5]):
            py = Emu(CONTENT_TOP + Inches(0.5) + j * Inches(0.6))
            _rect(slide, Emu(x + Inches(0.12)), Emu(py + Inches(0.06)), Inches(0.08), Inches(0.08), col)
            _tb(slide, Emu(x + Inches(0.28)), Emu(py), Inches(colW - 0.44), Inches(0.55),
                p, BODY_FONT, 10, DARK_TEXT, line_spacing=13)




def build_quote_full(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = CHARCOAL
    _rect(slide, Inches(0), Inches(0), Inches(0.5), H, DK_GREEN)
    _rect(slide, Inches(0.5), Inches(0), Inches(0.08), H, MINT)
    _tb(slide, Inches(1.2), Inches(1.0), Inches(7.8), Inches(2.8),
        c.get("quote", ""), BODY_FONT, 24, WHITE, italic=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=34)
    _rect(slide, Inches(3.5), Inches(4.0), Inches(3.0), Inches(0.06), MINT)
    if c.get("attribution"):
        _tb(slide, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.4),
            c["attribution"], BODY_FONT, 13, MINT, bold=True, align=PP_ALIGN.CENTER)
    if c.get("context"):
        _tb(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.35),
            c["context"], BODY_FONT, 10, RGBColor(0x99, 0x99, 0x99), italic=True,
            align=PP_ALIGN.CENTER)


def build_stat_hero(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Key Metric"), size=22)
    hero = c.get("hero", {})
    col = _resolve_section_color(c)
    _rect(slide, LM, Inches(1.2), Inches(8.2), Inches(1.8), CARD_BG)
    _rect(slide, LM, Inches(1.2), Inches(0.1), Inches(1.8), col)
    _tb(slide, LM, Inches(1.2), Inches(8.2), Inches(1.3),
        hero.get("value", "—"), TITLE_FONT, 64, col, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if hero.get("label"):
        _tb(slide, Inches(1.5), Inches(2.5), Inches(6.0), Inches(0.45),
            hero["label"], BODY_FONT, 14, DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    supporting = c.get("supporting", [])
    n = min(len(supporting), 4)
    if n > 0:
        sw = (8.2 - (n - 1) * 0.2) / n
        for i, s in enumerate(supporting[:n]):
            x = Inches(0.7 + i * (sw + 0.2))
            sc = ACCENTS[i % len(ACCENTS)]
            _rect(slide, x, Inches(3.25), Inches(sw), Inches(1.6), CARD_BG)
            _rect(slide, x, Inches(3.25), Inches(sw), Inches(0.06), sc)
            _tb(slide, x, Inches(3.35), Inches(sw), Inches(0.7),
                s.get("value", ""), TITLE_FONT, 26, sc, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _tb(slide, x, Inches(4.05), Inches(sw), Inches(0.35),
                s.get("label", ""), BODY_FONT, 10, MID_TEXT, align=PP_ALIGN.CENTER)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(5.0), Inches(8.5), Inches(0.3),
            c["source"], BODY_FONT, 9, MUTED, italic=True, align=PP_ALIGN.CENTER)


def build_in_brief_featured(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "In Brief"), size=24)
    featured = c.get("featured", "")
    supporting = c.get("supporting", [])
    _rect(slide, LM, CONTENT_TOP, Inches(8.2), Inches(1.3), CARD_BG)
    _rect(slide, LM, CONTENT_TOP, Inches(0.12), Inches(1.3), GOLD)
    _tb(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Inches(7.7), Inches(1.3),
        featured, BODY_FONT, 16, DARK_TEXT, bold=True, valign=MSO_ANCHOR.MIDDLE, line_spacing=22)
    for i, s in enumerate(supporting[:4]):
        y = Emu(Inches(2.7) + i * Inches(0.65))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, Inches(8.2), Inches(0.55), CARD_BG)
        _rect(slide, LM, y, Inches(0.08), Inches(0.55), col)
        _tb(slide, Emu(LM + Inches(0.2)), Emu(y), Inches(7.8), Inches(0.55),
            s, BODY_FONT, 12, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE)


def build_persona_duo(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Archetype Comparison"), size=22)
    personas = c.get("personas", [{}, {}])
    for idx, p in enumerate(personas[:2]):
        x = LM if idx == 0 else Inches(4.75)
        w = Inches(3.85) if idx == 0 else Inches(4.15)
        col = ACCENTS[idx % len(ACCENTS)]
        _rect(slide, x, CONTENT_TOP, w, Inches(3.9), CARD_BG)
        _rect(slide, x, CONTENT_TOP, Inches(0.1), Inches(3.9), col)
        _tb(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)), Emu(w - Inches(0.3)), Inches(0.4),
            p.get("name", ""), TITLE_FONT, 18, col, bold=True)
        _tb(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.5)), Emu(w - Inches(0.3)), Inches(0.25),
            p.get("archetype", ""), BODY_FONT, 10, GOLD, bold=True)
        traits = p.get("traits", [])
        for j, t in enumerate(traits[:4]):
            ty = Emu(CONTENT_TOP + Inches(0.9) + j * Inches(0.38))
            _oval(slide, Emu(x + Inches(0.2)), Emu(ty + Inches(0.04)), Inches(0.1), Inches(0.1), col)
            _tb(slide, Emu(x + Inches(0.4)), Emu(ty), Emu(w - Inches(0.55)), Inches(0.33),
                t, BODY_FONT, 10, DARK_TEXT)
        if p.get("strategy"):
            _rect(slide, Emu(x + Inches(0.1)), Emu(CONTENT_TOP + Inches(2.5)), Emu(w - Inches(0.2)), Inches(0.04), LIGHT_RULE)
            _tb(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(2.65)), Emu(w - Inches(0.3)), Inches(1.1),
                p["strategy"], BODY_FONT, 10, MID_TEXT, line_spacing=14)


def build_process_flow_vertical(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Process"), size=24)
    steps = c.get("steps", [])
    count = min(len(steps), 3)
    if count == 0: return
    avail = 3.8; stepH = (avail - (count - 1) * 0.4) / count
    for i, step in enumerate(steps[:count]):
        y = Inches(1.25 + i * (stepH + 0.4))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, Emu(y), Inches(8.2), Inches(stepH), CARD_BG)
        _rect(slide, LM, Emu(y), Inches(0.1), Inches(stepH), col)
        _tb(slide, Emu(LM + Inches(0.2)), Emu(y + Inches(0.05)), Inches(0.5), Inches(0.5),
            str(i + 1), TITLE_FONT, 24, col, bold=True)
        _tb(slide, Emu(LM + Inches(0.75)), Emu(y + Inches(0.05)), Inches(3.0), Inches(0.4),
            step.get("title", ""), BODY_FONT, 13, DARK_TEXT, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if step.get("detail"):
            _tb(slide, Emu(LM + Inches(0.75)), Emu(y + Inches(0.45)), Inches(7.0), Emu(Inches(stepH) - Inches(0.55)),
                step["detail"], BODY_FONT, 11, MID_TEXT, line_spacing=15)
        if i < count - 1:
            _tb(slide, Emu(LM + Inches(0.3)), Emu(y + Inches(stepH + 0.05)), Inches(0.5), Inches(0.25),
                "\u2193", TITLE_FONT, 18, MUTED, align=PP_ALIGN.CENTER)


def build_text_cards(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Key Points"), size=24)
    items = c.get("items", []); n = min(len(items), 6)
    if n == 0: return
    cols = 2 if n <= 4 else 3
    rows_count = (n + cols - 1) // cols
    cardW = (8.2 - (cols - 1) * 0.2) / cols
    cardH = (3.8 - (rows_count - 1) * 0.2) / rows_count
    for i, item in enumerate(items[:n]):
        col_idx = i % cols; row = i // cols
        x = Inches(0.7 + col_idx * (cardW + 0.2))
        y = Emu(CONTENT_TOP + row * (Inches(cardH) + Inches(0.2)))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, y, Inches(cardW), Inches(cardH), CARD_BG)
        _rect(slide, x, y, Inches(cardW), Inches(0.06), col)
        title_text = item if isinstance(item, str) else item.get("title", "")
        detail_text = "" if isinstance(item, str) else item.get("detail", "")
        if detail_text:
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.1)), Inches(cardW - 0.24), Inches(0.35),
                title_text, BODY_FONT, 12, DARK_TEXT, bold=True)
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.5)), Inches(cardW - 0.24), Emu(Inches(cardH) - Inches(0.6)),
                detail_text, BODY_FONT, 10, MID_TEXT, line_spacing=14)
        else:
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.1)), Inches(cardW - 0.24), Emu(Inches(cardH) - Inches(0.2)),
                title_text, BODY_FONT, 12, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE, line_spacing=17)


def build_text_columns(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Overview"), size=24)
    columns = c.get("columns", []); n = min(len(columns), 3)
    if n == 0: return
    colW = (8.2 - (n - 1) * 0.25) / n
    for i, col_data in enumerate(columns[:n]):
        x = Inches(0.7 + i * (colW + 0.25))
        col = ACCENTS[i % len(ACCENTS)]
        heading = "" if isinstance(col_data, str) else col_data.get("heading", "")
        body = col_data if isinstance(col_data, str) else col_data.get("body", "")
        if heading:
            _tb(slide, x, CONTENT_TOP, Inches(colW), Inches(0.3),
                heading, TITLE_FONT, 11, col, bold=True)
            _rect(slide, x, Emu(CONTENT_TOP + Inches(0.35)), Inches(0.8), Inches(0.04), col)
            _tb(slide, x, Emu(CONTENT_TOP + Inches(0.5)), Inches(colW), Inches(3.3),
                body, BODY_FONT, 11, DARK_TEXT, line_spacing=16)
        else:
            _tb(slide, x, CONTENT_TOP, Inches(colW), Inches(3.8),
                body, BODY_FONT, 11, DARK_TEXT, line_spacing=16)
        if i < n - 1:
            _rect(slide, Inches(0.7 + (i + 1) * colW + i * 0.25 + 0.12), CONTENT_TOP,
                  Inches(0.02), Inches(3.5), LIGHT_RULE)




def build_text_narrative(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Overview"), size=24)
    lede = c.get("lede", "")
    body = c.get("body", "")
    if isinstance(body, list): body = "\n\n".join(body)
    _rect(slide, LM, CONTENT_TOP, Inches(8.2), Inches(1.2), CARD_BG)
    _rect(slide, LM, CONTENT_TOP, Inches(0.1), Inches(1.2), GOLD)
    _tb(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Inches(7.7), Inches(1.2),
        lede, BODY_FONT, 15, DARK_TEXT, bold=True, line_spacing=22,
        valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, LM, Inches(2.6), Inches(8.2), Inches(2.5),
        body, BODY_FONT, 12, MID_TEXT, line_spacing=18)


def build_text_nested(prs, c):
    """Section cards — top-level items as colored label blocks, children as flowing text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Detail"), size=24)
    items = c.get("items", [])
    n = min(len(items), 4)
    if n == 0: return
    avail = H - CONTENT_TOP - Inches(0.3)
    gap = Inches(0.12)
    cardH = (avail - (n - 1) * gap) / n

    for i, item in enumerate(items[:n]):
        y = Emu(CONTENT_TOP + i * (cardH + gap))
        text = item if isinstance(item, str) else item.get("text", "")
        children = [] if isinstance(item, str) else item.get("children", [])
        col = ACCENTS[i % len(ACCENTS)]
        # Colored label block
        _rect(slide, LM, y, Inches(2.2), Emu(cardH), col)
        _tb(slide, Emu(LM + Inches(0.12)), y, Inches(1.95), Emu(cardH),
            text, BODY_FONT, 11, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=15)
        # Children card
        _rect(slide, Inches(3.05), y, Inches(5.85), Emu(cardH), CARD_BG)
        child_texts = []
        for ch in children:
            ch_text = ch if isinstance(ch, str) else ch.get("text", "")
            grandchildren = [] if isinstance(ch, str) else ch.get("children", [])
            child_texts.append(ch_text)
            for gc in grandchildren:
                gc_text = gc if isinstance(gc, str) else gc.get("text", "")
                child_texts.append("  \u2022 " + gc_text)
        body = "\n".join(child_texts)
        _tb(slide, Inches(3.18), y, Inches(5.55), Emu(cardH),
            body, BODY_FONT, 10, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=14)
def build_text_split(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Key Point"), size=24)
    headline = c.get("headline", "")
    detail = c.get("detail", "")
    points = c.get("points", [])
    _rect(slide, LM, CONTENT_TOP, Inches(4.0), Inches(3.8), CARD_BG)
    _rect(slide, LM, CONTENT_TOP, Inches(0.1), Inches(3.8), DK_GREEN)
    _tb(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Inches(3.5), Inches(2.0),
        headline, TITLE_FONT, 18, DARK_TEXT, bold=True, line_spacing=26,
        valign=MSO_ANCHOR.MIDDLE)
    if detail:
        _tb(slide, Emu(LM + Inches(0.25)), Inches(3.0), Inches(3.5), Inches(1.5),
            detail, BODY_FONT, 11, MID_TEXT, line_spacing=16)
    for i, p in enumerate(points[:6]):
        py = Emu(CONTENT_TOP + Inches(0.05) + i * Inches(0.6))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, Inches(4.95), py, Inches(4.0), Inches(0.5), CARD_BG)
        _rect(slide, Inches(4.95), py, Inches(0.08), Inches(0.5), col)
        _tb(slide, Inches(5.15), Emu(py), Inches(3.6), Inches(0.5),
            p, BODY_FONT, 11, DARK_TEXT, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)


def build_text_annotated(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _warm_bg(slide); _edge_column(slide)
    _content_title(slide, c.get("title", "Analysis"), size=24)
    items = c.get("items", [])
    n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.85))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))
    for i, item in enumerate(items[:n]):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, Inches(8.2), Emu(rowH), CARD_BG)
        _rect(slide, LM, y, Inches(1.5), Emu(rowH), col)
        _tb(slide, LM, y, Inches(1.5), Emu(rowH),
            item.get("label", ""), BODY_FONT, 10, WHITE, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(LM + Inches(1.7)), y, Inches(6.3), Emu(rowH),
            item.get("text", ""), BODY_FONT, 11, DARK_TEXT,
            valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


# ── Dispatch ─────────────────────────────────────────────
BUILDERS = {
    "title": build_title, "in_brief": build_in_brief,
    "section_divider": build_section_divider, "stat_callout": build_stat_callout,
    "quote": build_quote, "comparison": build_comparison,
    "text_graph": build_text_graph, "process_flow": build_process_flow,
    "matrix": build_matrix, "methods": build_methods,
    "hypotheses": build_hypotheses, "wsn_dense": build_wsn_dense,
    "wsn_reveal": build_wsn_reveal, "findings_recs": build_findings_recs,
    "findings_recs_dense": build_findings_recs_dense,
    "open_questions": build_open_questions, "agenda": build_agenda,
    "progressive_reveal": build_progressive_reveal, "closer": build_closer,
    "timeline": build_timeline, "data_table": build_data_table,
    "multi_stat": build_multi_stat, "persona": build_persona,
    "risk_tradeoff": build_risk_tradeoff, "appendix": build_appendix,
    "before_after": build_before_after, "summary": build_summary,
    "quote_full": build_quote_full, "stat_hero": build_stat_hero,
    "in_brief_featured": build_in_brief_featured, "persona_duo": build_persona_duo,
    "process_flow_vertical": build_process_flow_vertical,
    "text_cards": build_text_cards, "text_columns": build_text_columns,
    "text_narrative": build_text_narrative, "text_nested": build_text_nested,
    "text_split": build_text_split, "text_annotated": build_text_annotated,
}

def build_deck(slide_configs, output_path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    for slide_type, data in slide_configs:
        if slide_type == "skip": continue
        builder = BUILDERS.get(slide_type)
        if builder: builder(prs, data)
    prs.save(output_path); return output_path
