# -*- coding: utf-8 -*-
"""
Slick Minimal template v2 - with Zilla Slab / Source Sans 3 proxies,
increased text sizes, adaptive layouts, bold lead-ins, rounded cards.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

W = Inches(10)
H = Inches(5.625)

GREEN = RGBColor(0x36, 0x87, 0x27)
GREEN_LIGHT = RGBColor(0xF7, 0xF4, 0xE4)
GREEN_MID = RGBColor(0x1D, 0xE4, 0xCA)
BLUE = RGBColor(0x38, 0x80, 0xF3)
PURPLE = RGBColor(0x5B, 0x2C, 0x8F)
COBALT = RGBColor(0x04, 0x54, 0x7C)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
DARK = RGBColor(0x40, 0x3F, 0x3E)
MID = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xE5, 0xE5, 0xE5)
OFF_WHITE = RGBColor(0xF9, 0xF7, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLORS = [GREEN, BLUE, PURPLE, COBALT, GOLD]

SECTION_COLORS = {
    "green": GREEN, "blue": BLUE, "purple": PURPLE,
    "cobalt": COBALT, "gold": GOLD,
    "red": RGBColor(0xC2, 0x3B, 0x22), "teal": RGBColor(0x1B, 0x7A, 0x6E),
    "ochre": RGBColor(0xCC, 0x7A, 0x2E), "slate": RGBColor(0x5A, 0x6A, 0x7A),
    "plum": RGBColor(0x8E, 0x45, 0x85),
}

def _resolve_section_color(c, default=None):
    """Resolve sectionColor from slide data dict."""
    sc = c.get("sectionColor", "")
    return SECTION_COLORS.get(sc, default or GREEN)

TITLE_FONT = "Fidelity Slab"
BODY_FONT = "Fidelity Sans"

LM = Inches(0.9)
ACC_W = Inches(0.25)
CW = Inches(8.6)


# -- helpers --

def _add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); return shape

def _add_rounded_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    pg = shape._element.find('.//' + qn('a:prstGeom'))
    if pg is not None:
        av = pg.find(qn('a:avLst'))
        if av is None: av = etree.SubElement(pg, qn('a:avLst'))
        for g in av.findall(qn('a:gd')): av.remove(g)
        g = etree.SubElement(av, qn('a:gd')); g.set('name','adj'); g.set('fmla','val 5000')
    return shape

def _add_oval(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); return shape

DEFAULT_ICON_SHAPES = [MSO_SHAPE.OVAL, MSO_SHAPE.RECTANGLE, MSO_SHAPE.ISOSCELES_TRIANGLE]

def _add_icon_or_image(slide, img_path, x, y, w, h, shape_type, fill_color):
    """Add a user-supplied image or a coloured placeholder shape."""
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, w, h)
    else:
        s = slide.shapes.add_shape(shape_type, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill_color
        s.line.fill.background()

def _add_text_box(slide, x, y, w, h, text, font_name=None, font_size=12,
                  color=None, bold=False, italic=False, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, line_spacing=None):
    font_name = font_name or BODY_FONT; color = color or DARK
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.name = font_name
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.italic = italic; p.alignment = align
    if line_spacing: p.line_spacing = Pt(line_spacing)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', {MSO_ANCHOR.TOP:'t', MSO_ANCHOR.MIDDLE:'ctr', MSO_ANCHOR.BOTTOM:'b'}.get(valign,'t'))
    return txBox

def _accent(slide):
    _add_rect(slide, Inches(0), Inches(0), ACC_W, H, GREEN)

def _slide_title(slide, title, y=None, size=28):
    """Title at top with green rule. Moved UP for more content space."""
    _add_text_box(slide, LM, Inches(0.15), CW, Inches(0.65), title,
                  TITLE_FONT, size, DARK, bold=True, valign=MSO_ANCHOR.BOTTOM)

# Content starts at 0.95" (below title + breathing room)
CONTENT_TOP = Inches(0.95)

def _find_split(text):
    """Find where to split text into bold lead-in + regular continuation."""
    for delim in [' — ', ' – ']:
        idx = text.find(delim)
        if idx > 0: return text[:idx + len(delim)], text[idx + len(delim):]
    idx = text.find(': ')
    if idx > 0: return text[:idx + 1], text[idx + 2:]
    words = text.split(); char_count = 0
    for i, w in enumerate(words):
        char_count += len(w) + 1
        if i >= 3:
            ci = text.find(',', char_count - 1)
            if ci > 0 and ci < len(text) * 0.55: return text[:ci + 1], text[ci + 2:]
    if len(words) > 10:
        n = min(8, len(words) // 2)
        pos = sum(len(words[i]) + 1 for i in range(n))
        return text[:pos].rstrip(), text[pos:]
    return text, ""

def _add_split_text(slide, x, y, w, h, text, font_name, font_size, color, line_spacing=None):
    """Text with bold lead-in clause."""
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    bold_part, rest_part = _find_split(text)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = bold_part + (" " if rest_part else "")
    r1.font.name = font_name; r1.font.size = Pt(font_size)
    r1.font.color.rgb = color; r1.font.bold = True
    if rest_part:
        r2 = p.add_run(); r2.text = rest_part
        r2.font.name = font_name; r2.font.size = Pt(font_size)
        r2.font.color.rgb = color; r2.font.bold = False
    p.alignment = PP_ALIGN.LEFT
    if line_spacing: p.line_spacing = Pt(line_spacing)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None: bp.set('anchor', 'ctr')
    return tb

def _adaptive_layout(n, content_top, slide_h, bottom_margin=Inches(0.35)):
    """Calculate card positions - fills available space, biased slightly up."""
    avail = slide_h - content_top - bottom_margin
    gap = Inches(0.12)
    rowH = int((avail - (n - 1) * gap) / n)
    max_row = Inches(1.15)
    if rowH > max_row: rowH = max_row
    total = n * rowH + (n - 1) * gap
    remaining = avail - total
    startY = Emu(content_top + int(remaining * 0.3))
    fs = 16 if n <= 3 else 15
    ls = 20 if n <= 3 else 19
    return startY, rowH, gap, fs, ls


# -- slide builders --

def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), ACC_W, H, GREEN)
    _add_text_box(slide, LM, Inches(1.0), CW, Inches(1.5), c.get("title", "Title"),
                  TITLE_FONT, 34, DARK, bold=True, valign=MSO_ANCHOR.BOTTOM)
    _add_rect(slide, LM, Inches(2.65), Inches(2.5), Inches(0.04), GREEN)
    if c.get("subtitle"):
        _add_text_box(slide, LM, Inches(2.85), CW, Inches(0.6), c["subtitle"],
                      BODY_FONT, 16, MID)
    if c.get("author"):
        _add_text_box(slide, LM, Inches(4.1), CW, Inches(0.35), c["author"],
                      BODY_FONT, 12, MID, bold=True)
    if c.get("date"):
        _add_text_box(slide, LM, Inches(4.45), CW, Inches(0.3), c["date"],
                      BODY_FONT, 11, MID)
    img = c.get("imagePath", "")
    if img and os.path.exists(img):
        slide.shapes.add_picture(img, Inches(5.4), Inches(0.3), Inches(4.3), Inches(5.0))


def build_in_brief(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "In Brief"))
    bullets = c.get("bullets", [])
    n = len(bullets)
    startY, rowH, gap, fs, ls = _adaptive_layout(n, CONTENT_TOP, H)
    
    for i, b in enumerate(bullets):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.10), rowH, col)
        cs = Inches(0.40); cx = Emu(LM + Inches(0.22)); cy = Emu(y + (rowH - cs) // 2)
        _add_oval(slide, cx, cy, cs, cs, col)
        _add_text_box(slide, cx, cy, cs, cs, str(i+1), TITLE_FONT, 14, WHITE,
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_split_text(slide, Emu(LM + Inches(0.8)), y, Emu(CW - Inches(0.95)), rowH,
                        b, BODY_FONT, fs, DARK, line_spacing=ls)


def build_section_divider(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    if c.get("sectionNumber"):
        _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.8),
                      f"0{c['sectionNumber']}", TITLE_FONT, 48, GREEN_MID, bold=True,
                      valign=MSO_ANCHOR.BOTTOM)
    _add_text_box(slide, Inches(0.8), Inches(2.1), Inches(8), Inches(1.0),
                  c.get("title", "Section"), TITLE_FONT, 36, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, Inches(0.8), Inches(3.2), Inches(2.0), Inches(0.04), GREEN_MID)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.8), Inches(3.4), Inches(8), Inches(0.5),
                      c["subtitle"], BODY_FONT, 15, GREEN_LIGHT)


def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    _accent(slide); _slide_title(slide, c.get("title", "Key Metric"), size=22)
    # Visual anchor circle behind stat
    cs = Inches(2.2)
    cx = Emu((W - cs) // 2)
    cy = Emu(Inches(1.5))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, cs, cs)
    oval_fill = oval.fill
    oval_fill.solid()
    r, g, b = sc[0], sc[1], sc[2]
    light_r = min(255, r + int((255 - r) * 0.85))
    light_g = min(255, g + int((255 - g) * 0.85))
    light_b = min(255, b + int((255 - b) * 0.85))
    oval_fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
    oval.line.fill.background()
    _add_text_box(slide, LM, Inches(1.3), CW, Inches(1.8), c.get("stat", "—"),
                  TITLE_FONT, 96, sc, bold=True, align=PP_ALIGN.CENTER,
                  valign=MSO_ANCHOR.MIDDLE)
    if c.get("headline"):
        _add_text_box(slide, Inches(1.5), Inches(3.1), Inches(7), Inches(0.6),
                      c["headline"], BODY_FONT, 17, DARK, bold=True, align=PP_ALIGN.CENTER)
    if c.get("detail"):
        _add_text_box(slide, Inches(1.5), Inches(3.75), Inches(7), Inches(0.8),
                      c["detail"], BODY_FONT, 13, MID, align=PP_ALIGN.CENTER, line_spacing=17)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(4.85), Inches(9), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True, align=PP_ALIGN.CENTER)


def build_quote(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "In Their Words"), size=22)
    _add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.8),
                  "\u201C", "Georgia", 64, GREEN_LIGHT, bold=True,
                  align=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.0), Inches(1.5), Inches(8.0), Inches(2.0),
                  c.get("quote", ""), BODY_FONT, 18, DARK, italic=True,
                  valign=MSO_ANCHOR.MIDDLE, line_spacing=25,
                  align=PP_ALIGN.CENTER)
    _add_rect(slide, Inches(4.25), Inches(3.7), Inches(1.5), Inches(0.04), GREEN)
    if c.get("attribution"):
        _add_text_box(slide, Inches(1.0), Inches(3.85), Inches(8.0), Inches(0.35),
                      c["attribution"], BODY_FONT, 13, MID,
                      align=PP_ALIGN.CENTER)
    if c.get("context"):
        _add_text_box(slide, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.35),
                      c["context"], BODY_FONT, 11, MID, italic=True,
                      align=PP_ALIGN.CENTER)


def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Comparison"))
    sc = _resolve_section_color(c)

    left_color = sc if sc != GREEN else COBALT
    right_color = GREEN

    left_label = c.get("leftLabel", "Before")
    right_label = c.get("rightLabel", "After")
    left_items = c.get("leftItems", [])
    right_items = c.get("rightItems", [])

    card_w_left = Inches(4.1)
    card_w_right = Inches(4.1)
    card_h = Inches(3.8)
    gap = Inches(0.4)
    header_h = Inches(0.5)

    # Left card
    left_x = LM
    _add_rounded_rect(slide, left_x, CONTENT_TOP, card_w_left, card_h, WHITE)
    _add_rect(slide, left_x, CONTENT_TOP, card_w_left, header_h, left_color)
    _add_text_box(slide, Emu(left_x + Inches(0.2)), CONTENT_TOP,
                  Emu(card_w_left - Inches(0.4)), header_h,
                  left_label, BODY_FONT, 14, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(left_items[:6]):
        iy = Emu(CONTENT_TOP + header_h + Inches(0.15) + i * Inches(0.5))
        text = item if isinstance(item, str) else str(item)
        _add_text_box(slide, Emu(left_x + Inches(0.25)), iy,
                      Emu(card_w_left - Inches(0.5)), Inches(0.45),
                      text, BODY_FONT, 13, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)

    # Right card
    right_x = Emu(LM + card_w_left + gap)
    _add_rounded_rect(slide, right_x, CONTENT_TOP, card_w_right, card_h, WHITE)
    _add_rect(slide, right_x, CONTENT_TOP, card_w_right, header_h, right_color)
    _add_text_box(slide, Emu(right_x + Inches(0.2)), CONTENT_TOP,
                  Emu(card_w_right - Inches(0.4)), header_h,
                  right_label, BODY_FONT, 14, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(right_items[:6]):
        iy = Emu(CONTENT_TOP + header_h + Inches(0.15) + i * Inches(0.5))
        text = item if isinstance(item, str) else str(item)
        _add_text_box(slide, Emu(right_x + Inches(0.25)), iy,
                      Emu(card_w_right - Inches(0.5)), Inches(0.45),
                      text, BODY_FONT, 13, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)


def build_text_graph(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Title"), size=24)
    texts = c.get("text", []); texts = [texts] if not isinstance(texts, list) else texts
    for i, t in enumerate(texts):
        _add_text_box(slide, LM, Emu(CONTENT_TOP + i * Inches(1.2)),
                      Inches(4.0), Inches(1.1), t, BODY_FONT, 13, DARK, line_spacing=18)
    _add_rect(slide, Inches(5.1), CONTENT_TOP, Inches(0.03), Inches(3.5), LIGHT)
    chart_data_raw = c.get("chartData", [{"name": "S1", "labels": ["A","B","C"], "values": [25,45,30]}])
    from pptx.chart.data import CategoryChartData
    chart_data = CategoryChartData()
    if chart_data_raw:
        cd = chart_data_raw[0]
        chart_data.categories = cd.get("labels", ["A","B","C"])
        chart_data.add_series(cd.get("name","Series 1"), cd.get("values",[25,45,30]))
    ct = {"line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}.get(c.get("chartType","bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(ct, Inches(5.3), CONTENT_TOP, Inches(4.2), Inches(3.8), chart_data)
    if c.get("note"):
        _add_text_box(slide, Inches(5.3), Inches(5.0), Inches(4.2), Inches(0.3),
                      c["note"], BODY_FONT, 8, MID, italic=True)


def _draw_stepper_bar(slide, labels, active_count, y):
    """Reusable stepper bar: numbered circles + connecting lines + labels."""
    n = len(labels)
    cs = Inches(0.42)  # circle diameter
    total_w = CW
    spacing = int(total_w / max(n - 1, 1)) if n > 1 else 0
    start_x = LM
    for i in range(n):
        cx = Emu(start_x + i * spacing) if n > 1 else Emu(start_x + total_w // 2 - cs // 2)
        active = i < active_count
        color = COLORS[i % len(COLORS)] if active else LIGHT
        _add_oval(slide, cx, y, cs, cs, color)
        tc = WHITE if active else MID
        _add_text_box(slide, cx, y, cs, cs, str(i + 1), TITLE_FONT, 15, tc,
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        lw = Inches(1.4)
        label_x = Emu(cx - (lw - cs) // 2)
        _add_text_box(slide, label_x, Emu(y + cs + Inches(0.04)),
                      lw, Inches(0.25), labels[i],
                      BODY_FONT, 9, color if active else MID,
                      bold=active, align=PP_ALIGN.CENTER)
        if i < n - 1:
            line_x1 = Emu(cx + cs + Inches(0.08))
            next_cx = Emu(start_x + (i + 1) * spacing)
            line_x2 = Emu(next_cx - Inches(0.08))
            _add_rect(slide, line_x1, Emu(y + cs // 2), Emu(line_x2 - line_x1), Inches(0.025), LIGHT)


def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _slide_title(slide, c.get("title", "Process"))
    steps = c.get("steps", [])
    n = min(len(steps), 5)
    if n == 0: return

    labels = [s.get("title", f"Step {i+1}") for i, s in enumerate(steps[:n])]
    _draw_stepper_bar(slide, labels, n, Emu(CONTENT_TOP + Inches(0.05)))

    card_top = Emu(CONTENT_TOP + Inches(0.85))
    gap = Inches(0.2)

    if n <= 3:
        card_w = int((CW - (n - 1) * gap) / n)
        card_h = Inches(3.0)
        for i, step in enumerate(steps[:n]):
            cx = Emu(LM + i * (card_w + gap))
            _add_rounded_rect(slide, cx, card_top, card_w, card_h, OFF_WHITE)
            _add_rect(slide, cx, card_top, Inches(0.06), card_h, COLORS[i % len(COLORS)])
            _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(card_top + Inches(0.15)),
                          Emu(card_w - Inches(0.4)), Inches(0.5),
                          step.get("title", ""), BODY_FONT, 13, DARK, bold=True)
            if step.get("detail"):
                _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(card_top + Inches(0.65)),
                              Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.85)),
                              step["detail"], BODY_FONT, 10, MID, line_spacing=14)
    elif n == 4:
        card_w = int((CW - gap) / 2)
        card_h = Inches(1.4)
        for i, step in enumerate(steps[:4]):
            col = i % 2
            row = i // 2
            cx = Emu(LM + col * (card_w + gap))
            cy = Emu(card_top + row * (card_h + gap))
            _add_rounded_rect(slide, cx, cy, card_w, card_h, OFF_WHITE)
            _add_rect(slide, cx, cy, Inches(0.06), card_h, COLORS[i % len(COLORS)])
            _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.1)),
                          Emu(card_w - Inches(0.4)), Inches(0.4),
                          step.get("title", ""), BODY_FONT, 12, DARK, bold=True)
            if step.get("detail"):
                _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.5)),
                              Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.65)),
                              step["detail"], BODY_FONT, 10, MID, line_spacing=14)
    else:
        card_w = int((CW - 2 * gap) / 3)
        card_h = Inches(1.35)
        for i, step in enumerate(steps[:5]):
            if i < 3:
                cx = Emu(LM + i * (card_w + gap))
                cy = card_top
            else:
                offset = (CW - 2 * card_w - gap) // 2
                cx = Emu(LM + offset + (i - 3) * (card_w + gap))
                cy = Emu(card_top + card_h + gap)
            _add_rounded_rect(slide, cx, cy, card_w, card_h, OFF_WHITE)
            _add_rect(slide, cx, cy, Inches(0.06), card_h, COLORS[i % len(COLORS)])
            _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.1)),
                          Emu(card_w - Inches(0.4)), Inches(0.35),
                          step.get("title", ""), BODY_FONT, 12, DARK, bold=True)
            if step.get("detail"):
                _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.45)),
                              Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.55)),
                              step["detail"], BODY_FONT, 10, MID, line_spacing=14)


def build_process_flow_reveal(prs, c):
    """Generates N slides, one per step, with progressive stepper."""
    steps = c.get("steps", [])
    n = min(len(steps), 5)
    if n == 0: return

    labels = [s.get("title", f"Step {i+1}") for i, s in enumerate(steps[:n])]

    for step_idx in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _accent(slide)
        _slide_title(slide, c.get("title", "Process"))

        _draw_stepper_bar(slide, labels, step_idx + 1, Emu(CONTENT_TOP + Inches(0.05)))

        card_y = Emu(CONTENT_TOP + Inches(0.85))
        card_w = Inches(6.0)
        card_h = Inches(3.0)
        card_x = Emu(LM + (CW - card_w) // 2)

        step = steps[step_idx]
        _add_rounded_rect(slide, card_x, card_y, card_w, card_h, OFF_WHITE)
        _add_rect(slide, card_x, card_y, Inches(0.08), card_h, COLORS[step_idx % len(COLORS)])
        _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.2)),
                      Emu(card_w - Inches(0.5)), Inches(0.6),
                      step.get("title", ""), BODY_FONT, 16, DARK, bold=True)
        if step.get("detail"):
            _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.85)),
                          Emu(card_w - Inches(0.5)), Emu(card_h - Inches(1.1)),
                          step["detail"], BODY_FONT, 12, MID, line_spacing=16)


def build_matrix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Framework"), size=24)
    quads = c.get("quadrants", [{},{},{},{}])
    accents = [GREEN, BLUE, PURPLE, COBALT]
    qW = Inches(4.05); qH = Inches(1.85); gap = Inches(0.2)
    sX = LM; sY = CONTENT_TOP
    for i, q in enumerate(quads[:4]):
        col = i % 2; row = i // 2
        x = Emu(sX + col * (qW + gap)); y = Emu(sY + row * (qH + gap))
        _add_rounded_rect(slide, x, y, qW, qH, OFF_WHITE)
        _add_rect(slide, x, y, Inches(0.06), qH, accents[i])
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.12)),
                      Emu(qW - Inches(0.4)), Inches(0.35),
                      q.get("label", ""), TITLE_FONT, 12, accents[i], bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.5)),
                      Emu(qW - Inches(0.4)), Emu(qH - Inches(0.65)),
                      q.get("detail", ""), BODY_FONT, 11, DARK, line_spacing=15)


def build_methods(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Approach"), size=24)
    fields = c.get("fields", []); n = len(fields)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int(avail / n), Inches(0.82))
    gap = Inches(0.06) if n > 4 else Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.3))
    for i, f in enumerate(fields):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.10), rowH, col)
        _add_text_box(slide, Emu(LM + Inches(0.22)), y, Inches(1.8), rowH,
                      f.get("label", ""), BODY_FONT, 13, col, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Emu(LM + Inches(2.1)), y, Inches(6.3), rowH,
                      f.get("value", ""), BODY_FONT, 13, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_hypotheses(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Hypotheses"), size=24)
    hyps = c.get("hypotheses", []); n = len(hyps)
    if n == 0: return
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int(avail / n), Inches(0.88))
    gap = Inches(0.12)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.5))
    for i, h in enumerate(hyps):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.10), rowH, col)
        cs = Inches(0.44); cx = Emu(LM + Inches(0.18)); cy = Emu(y + (rowH - cs) // 2)
        _add_oval(slide, cx, cy, cs, cs, col)
        _add_text_box(slide, cx, cy, cs, cs, f"H{i+1}", TITLE_FONT, 12, WHITE,
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Emu(LM + Inches(0.72)), y, Inches(5.95), rowH,
                      h.get("text", ""), BODY_FONT, 13, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        if h.get("status"):
            sc = GREEN if h["status"] == "Confirmed" else (COBALT if h["status"] == "Rejected" else MID)
            _add_text_box(slide, Inches(7.8), y, Inches(1.2), rowH,
                          h["status"], BODY_FONT, 10, sc, bold=True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def build_wsn_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Key Finding"), size=24)
    labels = c.get("labels", ["What", "So What", "Now What"])
    cols = [(labels[0], GREEN, c.get("what", {})),
            (labels[1], BLUE, c.get("soWhat", {})),
            (labels[2], PURPLE, c.get("nowWhat", {}))]
    colW = Inches(2.75); gap = Inches(0.2); cardH = Inches(3.8)
    startX = LM; startY = CONTENT_TOP
    for i, (label, color, data) in enumerate(cols):
        x = Emu(startX + i * (colW + gap))
        _add_rounded_rect(slide, x, startY, colW, cardH, OFF_WHITE)
        _add_rect(slide, x, startY, Inches(0.06), cardH, color)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(startY + Inches(0.15)),
                      Emu(colW - Inches(0.4)), Inches(0.35),
                      label, TITLE_FONT, 14, color, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(startY + Inches(0.55)),
                      Emu(colW - Inches(0.4)), Inches(0.7),
                      data.get("headline", ""), BODY_FONT, 12, DARK, bold=True,
                      line_spacing=16)
        if data.get("detail"):
            _add_text_box(slide, Emu(x + Inches(0.2)), Emu(startY + Inches(1.3)),
                          Emu(colW - Inches(0.4)), Inches(1.9),
                          data["detail"], BODY_FONT, 11, MID, line_spacing=14)


def build_wsn_reveal(prs, c):
    """3 progressive slides: What -> So What -> Now What with running summary."""
    labels = c.get("labels", ["What", "So What", "Now What"])
    WSN_KEYS = ["what", "soWhat", "nowWhat"]
    WSN_COLORS = [GREEN, BLUE, PURPLE]

    for step in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _accent(slide)
        _slide_title(slide, c.get("title", "Key Finding"), size=24)

        # Stepper bar at top
        _draw_stepper_bar(slide, labels, step + 1, Emu(CONTENT_TOP + Inches(0.05)))

        # Featured card for current step
        data = c.get(WSN_KEYS[step], {})
        card_y = Emu(CONTENT_TOP + Inches(0.85))
        card_w = Inches(7.0)
        card_h = Inches(2.4)
        card_x = Emu(LM + (CW - card_w) // 2)

        _add_rounded_rect(slide, card_x, card_y, card_w, card_h, OFF_WHITE)
        _add_rect(slide, card_x, card_y, Inches(0.10), card_h, WSN_COLORS[step])
        _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.15)),
                      Emu(card_w - Inches(0.5)), Inches(0.6),
                      data.get("headline", ""), BODY_FONT, 15, DARK, bold=True,
                      line_spacing=18)
        if data.get("detail"):
            _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.8)),
                          Emu(card_w - Inches(0.5)), Emu(card_h - Inches(1.0)),
                          data["detail"], BODY_FONT, 12, MID, line_spacing=16)

        # Running summary (only for step >= 1)
        if step > 0:
            summary_y = Emu(card_y + card_h + Inches(0.15))
            # Thin divider
            _add_rect(slide, LM, summary_y, Inches(3.0), Inches(0.02), LIGHT)
            # Previous step summaries
            for j in range(step):
                iy = Emu(summary_y + Inches(0.1) + j * Inches(0.28))
                prev_data = c.get(WSN_KEYS[j], {})
                summary_text = prev_data.get("summary", prev_data.get("headline", ""))
                # Small bullet square
                _add_rect(slide, LM, Emu(iy + Inches(0.05)), Inches(0.08), Inches(0.08), LIGHT)
                _add_text_box(slide, Emu(LM + Inches(0.18)), iy,
                              Emu(CW - Inches(0.18)), Inches(0.25),
                              summary_text, BODY_FONT, 9, MID,
                              valign=MSO_ANCHOR.MIDDLE)

def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Findings & Recommendations"), size=24)
    items = c.get("items", []); n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.12)) / n), Inches(0.85))
    gap = Inches(0.12)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.3))
    for i, item in enumerate(items[:5]):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, Inches(3.9), rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.10), rowH, col)
        _add_text_box(slide, Emu(LM + Inches(0.2)), y, Inches(3.5), rowH,
                      item.get("finding", ""), BODY_FONT, 12, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        # Arrow in colored circle
        acs = Inches(0.30); acx = Emu(Inches(4.97)); acy = Emu(y + (rowH - acs) // 2)
        _add_oval(slide, acx, acy, acs, acs, col)
        _add_text_box(slide, acx, acy, acs, acs,
                      "\u2192", BODY_FONT, 13, WHITE, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_rounded_rect(slide, Inches(5.4), y, Inches(4.1), rowH, OFF_WHITE)
        _add_rect(slide, Inches(5.4), y, Inches(0.10), rowH, col)
        _add_text_box(slide, Inches(5.6), y, Inches(3.75), rowH,
                      item.get("recommendation", ""), BODY_FONT, 12, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Complete Findings"), size=22)
    items = c.get("items", []); n = min(len(items), 8)
    avail = H - CONTENT_TOP - Inches(0.2)
    rowH = min(int((avail - (n - 1) * Inches(0.08)) / n), Inches(0.56))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.3))
    for i, item in enumerate(items[:8]):
        y = Emu(startY + i * (rowH + gap))
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        _add_rect(slide, LM, y, Inches(4.1), rowH, bg)
        _add_rect(slide, LM, y, Inches(0.04), rowH, GREEN)
        _add_text_box(slide, Emu(LM + Inches(0.15)), y, Inches(3.85), rowH,
                      item.get("finding", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, Inches(5.15), y, Inches(4.35), rowH, bg)
        _add_rect(slide, Inches(5.15), y, Inches(0.04), rowH, BLUE)
        _add_text_box(slide, Inches(5.3), y, Inches(4.1), rowH,
                      item.get("recommendation", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Open Questions"), size=26)
    questions = c.get("questions", [])
    cardW = Inches(4.1); cardH = Inches(1.7); gX = Inches(0.4); gY = Inches(0.2)
    gridY = CONTENT_TOP
    for i, question in enumerate(questions[:4]):
        col = i % 2; row = i // 2
        x = Emu(LM + col * (cardW + gX)); y = Emu(gridY + row * (cardH + gY))
        _add_rounded_rect(slide, x, y, cardW, cardH, OFF_WHITE)
        _add_rect(slide, x, y, Inches(0.06), cardH, COLORS[i % len(COLORS)])
        _add_oval(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.12)),
                  Inches(0.35), Inches(0.35), COLORS[i % len(COLORS)])
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.12)),
                      Inches(0.35), Inches(0.35), str(i + 1),
                      TITLE_FONT, 14, WHITE, bold=True, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.55)),
                      Emu(cardW - Inches(0.4)), Inches(1.0),
                      question, BODY_FONT, 13, DARK, line_spacing=17)


def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Agenda"), size=26)
    items = c.get("items", []); n = len(items)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int(avail / n), Inches(0.72))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.3))
    for i, item in enumerate(items):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.10), rowH, col)
        cs = Inches(0.40); cx = Emu(LM + Inches(0.22)); cy = Emu(y + (rowH - cs) // 2)
        _add_oval(slide, cx, cy, cs, cs, col)
        _add_text_box(slide, cx, cy, cs, cs, str(i + 1), TITLE_FONT, 14, WHITE,
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        title_text = item if isinstance(item, str) else item.get("title", "")
        _add_text_box(slide, Emu(LM + Inches(0.8)), y, Inches(5.5), rowH,
                      title_text, BODY_FONT, 15, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)
        if isinstance(item, dict) and item.get("detail"):
            _add_text_box(slide, Inches(7.5), y, Inches(1.8), rowH,
                          item["detail"], BODY_FONT, 12, MID,
                          align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_progressive_reveal(prs, c):
    takeaways = c.get("takeaways", [])
    for n in range(min(len(takeaways), 5)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _accent(slide); _slide_title(slide, c.get("title", "Building the Picture"))
        cur = takeaways[n]
        _add_rounded_rect(slide, LM, CONTENT_TOP, CW, Inches(2.4), OFF_WHITE)
        _add_rect(slide, LM, CONTENT_TOP, Inches(0.06), Inches(2.4), GREEN)
        _add_text_box(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)),
                      Emu(CW - Inches(0.4)), Inches(0.6),
                      cur.get("headline", ""), BODY_FONT, 16, DARK, bold=True)
        if cur.get("detail"):
            _add_text_box(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.75)),
                          Emu(CW - Inches(0.4)), Inches(1.4),
                          cur["detail"], BODY_FONT, 12, MID, line_spacing=16)
        # Thin subtle divider
        _add_rect(slide, LM, Inches(3.65), Inches(3.0), Inches(0.02), LIGHT)
        # Running summary items (no header, all gray)
        for j in range(n + 1):
            ty = Emu(Inches(3.75) + j * Inches(0.28))
            _add_rect(slide, LM, Emu(ty + Inches(0.04)), Inches(0.08), Inches(0.08), LIGHT)
            _add_text_box(slide, Emu(LM + Inches(0.18)), ty,
                          Emu(CW - Inches(0.18)), Inches(0.25),
                          takeaways[j].get("summary", takeaways[j].get("headline", "")),
                          BODY_FONT, 9, MID,
                          valign=MSO_ANCHOR.MIDDLE)


def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    _add_text_box(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(1.2),
                  c.get("title", "Thank You"), TITLE_FONT, 44, WHITE, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
    _add_rect(slide, Inches(3.75), Inches(2.75), Inches(2.5), Inches(0.04), GREEN_MID)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.5), Inches(2.95), Inches(9), Inches(0.5),
                      c["subtitle"], BODY_FONT, 16, WHITE, align=PP_ALIGN.CENTER)
    if c.get("contact"):
        _add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
                      c["contact"], BODY_FONT, 12, GREEN_LIGHT, align=PP_ALIGN.CENTER)


def build_timeline(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Timeline"))
    milestones = c.get("milestones", []); n = min(len(milestones), 6)
    if n == 0: return
    lineY = Inches(2.8); startX = 0.9; endX = 9.2
    step = (endX - startX) / max(n - 1, 1)
    _add_rect(slide, Inches(startX), lineY, Inches(endX - startX), Inches(0.04), LIGHT)
    STATUS_CLR = {"complete": GREEN, "current": GOLD, "upcoming": MID}
    for i, m in enumerate(milestones[:n]):
        cx = Inches(startX + i * step); col = STATUS_CLR.get(m.get("status", "upcoming"), MID)
        ds = Inches(0.22)
        _add_oval(slide, Emu(cx - ds // 2), Emu(lineY - ds // 2), ds, ds, col)
        _add_text_box(slide, Emu(cx - Inches(0.6)), Emu(lineY - Inches(0.65)),
                      Inches(1.2), Inches(0.3), m.get("date", ""),
                      BODY_FONT, 10, col, bold=True, align=PP_ALIGN.CENTER)
        _add_text_box(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.25)),
                      Inches(1.4), Inches(0.4), m.get("title", ""),
                      BODY_FONT, 11, DARK, bold=True, align=PP_ALIGN.CENTER)
        if m.get("detail"):
            _add_text_box(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.65)),
                          Inches(1.4), Inches(0.7), m["detail"],
                          BODY_FONT, 9, MID, align=PP_ALIGN.CENTER, line_spacing=12)


def build_data_table(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Data"), size=24)
    headers = c.get("headers", []); rows = c.get("rows", [])
    hc = c.get("highlightCol", None)
    nCols = len(headers); nRows = min(len(rows), 10)
    if nCols == 0: return
    tW = CW; tX = LM; colW = tW / nCols; rowH = min(Inches(0.42), Inches(3.8) / max(nRows + 1, 1))
    hdrY = CONTENT_TOP
    _add_rect(slide, tX, hdrY, tW, rowH, GREEN)
    for j, h in enumerate(headers):
        _add_text_box(slide, Emu(tX + j * colW + Inches(0.08)), hdrY,
                      Emu(colW - Inches(0.16)), rowH, h,
                      BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, row in enumerate(rows[:nRows]):
        ry = Emu(hdrY + (i + 1) * rowH)
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        _add_rect(slide, tX, ry, tW, rowH, bg)
        for j, cell in enumerate(row[:nCols]):
            is_hl = hc is not None and j == hc
            _add_text_box(slide, Emu(tX + j * colW + Inches(0.08)), ry,
                          Emu(colW - Inches(0.16)), rowH, str(cell),
                          BODY_FONT, 10, GREEN if is_hl else DARK,
                          bold=is_hl, valign=MSO_ANCHOR.MIDDLE)
    if c.get("note"):
        _add_text_box(slide, tX, Emu(hdrY + (nRows + 1) * rowH + Inches(0.1)),
                      tW, Inches(0.3), c["note"], BODY_FONT, 8, MID, italic=True)


def build_multi_stat(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Key Metrics"), size=24)
    stats = c.get("stats", []); n = min(len(stats), 4)
    if n == 0: return
    gap = Inches(0.2); statW = (CW - (n - 1) * gap) / n
    for i, s in enumerate(stats[:n]):
        x = Emu(LM + i * (statW + gap)); col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, x, CONTENT_TOP, statW, Inches(3.5), OFF_WHITE)
        _add_rect(slide, x, CONTENT_TOP, statW, Inches(0.06), col)
        _add_text_box(slide, x, Emu(CONTENT_TOP + Inches(0.2)), statW, Inches(1.4),
                      s.get("value", "—"), TITLE_FONT, 48, col, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, x, Emu(CONTENT_TOP + Inches(1.7)), statW, Inches(0.45),
                      s.get("label", ""), BODY_FONT, 12, DARK, bold=True, align=PP_ALIGN.CENTER)
        if s.get("detail"):
            _add_text_box(slide, x, Emu(CONTENT_TOP + Inches(2.2)), statW, Inches(0.9),
                          s["detail"], BODY_FONT, 10, MID, align=PP_ALIGN.CENTER, line_spacing=14)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True, align=PP_ALIGN.CENTER)


def build_persona(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Persona"), size=22)
    _add_rounded_rect(slide, LM, CONTENT_TOP, Inches(4.0), Inches(3.9), OFF_WHITE)
    _add_rect(slide, LM, CONTENT_TOP, Inches(0.08), Inches(3.9), GREEN)
    _add_text_box(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)),
                  Inches(3.6), Inches(0.5), c.get("name", ""),
                  TITLE_FONT, 22, GREEN, bold=True)
    _add_text_box(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.6)),
                  Inches(3.6), Inches(0.3), c.get("archetype", ""),
                  BODY_FONT, 12, GOLD, bold=True)
    for i, t in enumerate(c.get("traits", [])[:5]):
        ty = Emu(CONTENT_TOP + Inches(1.1) + i * Inches(0.4))
        _add_oval(slide, Emu(LM + Inches(0.2)), Emu(ty + Inches(0.04)),
                  Inches(0.12), Inches(0.12), COLORS[i % len(COLORS)])
        _add_text_box(slide, Emu(LM + Inches(0.45)), ty, Inches(3.4), Inches(0.35),
                      t, BODY_FONT, 11, DARK)
    _add_rounded_rect(slide, Inches(5.1), CONTENT_TOP, Inches(4.4), Inches(3.9), OFF_WHITE)
    _add_rect(slide, Inches(5.1), CONTENT_TOP, Inches(0.08), Inches(3.9), GOLD)
    _add_text_box(slide, Inches(5.3), Emu(CONTENT_TOP + Inches(0.1)),
                  Inches(4.0), Inches(0.25), "STRATEGY", BODY_FONT, 9, MID, bold=True)
    _add_text_box(slide, Inches(5.3), Emu(CONTENT_TOP + Inches(0.4)),
                  Inches(4.0), Inches(1.5), c.get("strategy", ""),
                  BODY_FONT, 12, DARK, line_spacing=17)
    if c.get("detail"):
        _add_text_box(slide, Inches(5.3), Emu(CONTENT_TOP + Inches(2.0)),
                      Inches(4.0), Inches(1.5), c["detail"],
                      BODY_FONT, 10, MID, line_spacing=14)


def build_risk_tradeoff(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Risk & Reward"), size=24)
    risks = c.get("risks", []); rewards = c.get("rewards", [])
    RED = RGBColor(0xC2, 0x3B, 0x22)
    SEV = {"high": RED, "medium": GOLD, "low": GREEN}
    _add_rect(slide, LM, CONTENT_TOP, Inches(4.0), Inches(0.4), RED)
    _add_text_box(slide, Emu(LM + Inches(0.1)), CONTENT_TOP, Inches(3.8), Inches(0.4),
                  "RISKS", BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(risks[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.5) + i * Inches(0.65))
        sc = SEV.get(r.get("severity", "medium"), GOLD)
        _add_rounded_rect(slide, LM, ry, Inches(4.0), Inches(0.55), OFF_WHITE)
        _add_rect(slide, LM, ry, Inches(0.06), Inches(0.55), sc)
        _add_text_box(slide, Emu(LM + Inches(0.15)), ry, Inches(3.7), Inches(0.55),
                      r.get("label", "") + (" — " + r.get("detail", "") if r.get("detail") else ""),
                      BODY_FONT, 10, DARK, valign=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, Inches(5.15), CONTENT_TOP, Inches(4.35), Inches(0.4), GREEN)
    _add_text_box(slide, Inches(5.25), CONTENT_TOP, Inches(4.15), Inches(0.4),
                  "REWARDS", BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(rewards[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.5) + i * Inches(0.65))
        _add_rounded_rect(slide, Inches(5.15), ry, Inches(4.35), Inches(0.55), OFF_WHITE)
        _add_rect(slide, Inches(5.15), ry, Inches(0.06), Inches(0.55), GREEN)
        _add_text_box(slide, Inches(5.3), ry, Inches(4.05), Inches(0.55),
                      r.get("label", "") + (" — " + r.get("detail", "") if r.get("detail") else ""),
                      BODY_FONT, 10, DARK, valign=MSO_ANCHOR.MIDDLE)


def build_appendix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide)
    _add_rect(slide, LM, Inches(0.25), Inches(2.0), Inches(0.04), MID)
    _add_text_box(slide, LM, Inches(0.35), CW, Inches(0.4),
                  c.get("title", "Appendix"), TITLE_FONT, 16, MID, bold=True)
    y_cursor = 0.85
    for s in c.get("sections", []):
        if y_cursor > 4.8: break
        _add_text_box(slide, LM, Inches(y_cursor), CW, Inches(0.22),
                      s.get("label", ""), BODY_FONT, 9, GREEN, bold=True)
        _add_text_box(slide, LM, Inches(y_cursor + 0.22), CW, Inches(0.75),
                      s.get("content", ""), BODY_FONT, 8, MID, line_spacing=11)
        y_cursor += 1.05


def build_before_after(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Transformation"), size=24)
    bef = c.get("before", {}); aft = c.get("after", {}); interv = c.get("intervention", "")
    _add_rounded_rect(slide, LM, CONTENT_TOP, Inches(3.5), Inches(3.5), OFF_WHITE)
    _add_rect(slide, LM, CONTENT_TOP, Inches(3.5), Inches(0.4), COBALT)
    _add_text_box(slide, Emu(LM + Inches(0.12)), CONTENT_TOP, Inches(3.2), Inches(0.4),
                  bef.get("label", "Before"), TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _add_text_box(slide, Emu(LM + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.5)), Inches(3.2), Inches(2.7),
                  bef.get("detail", ""), BODY_FONT, 11, DARK, line_spacing=16)
    _add_text_box(slide, Inches(4.55), Emu(CONTENT_TOP + Inches(1.0)), Inches(0.7), Inches(0.5),
                  "\u2192", TITLE_FONT, 28, GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if interv:
        _add_text_box(slide, Inches(4.35), Emu(CONTENT_TOP + Inches(1.6)), Inches(1.1), Inches(1.5),
                      interv, BODY_FONT, 9, MID, align=PP_ALIGN.CENTER, line_spacing=13)
    _add_rounded_rect(slide, Inches(5.55), CONTENT_TOP, Inches(3.95), Inches(3.5), OFF_WHITE)
    _add_rect(slide, Inches(5.55), CONTENT_TOP, Inches(3.95), Inches(0.4), GREEN)
    _add_text_box(slide, Inches(5.67), CONTENT_TOP, Inches(3.7), Inches(0.4),
                  aft.get("label", "After"), TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _add_text_box(slide, Inches(5.67), Emu(CONTENT_TOP + Inches(0.5)), Inches(3.7), Inches(2.7),
                  aft.get("detail", ""), BODY_FONT, 11, DARK, line_spacing=16)


def build_summary(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Summary"), size=26)
    sections = c.get("sections", []); n = min(len(sections), 4)
    if n == 0: return
    gap = Inches(0.2); colW = (CW - (n - 1) * gap) / n
    for i, sec in enumerate(sections[:n]):
        x = Emu(LM + i * (colW + gap)); col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, x, CONTENT_TOP, colW, Inches(3.8), OFF_WHITE)
        _add_rect(slide, x, CONTENT_TOP, colW, Inches(0.06), col)
        _add_text_box(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.12)),
                      Emu(colW - Inches(0.24)), Inches(0.3),
                      sec.get("heading", ""), TITLE_FONT, 11, col, bold=True)
        for j, p in enumerate(sec.get("points", [])[:5]):
            py = Emu(CONTENT_TOP + Inches(0.5) + j * Inches(0.6))
            _add_rect(slide, Emu(x + Inches(0.12)), Emu(py + Inches(0.06)),
                      Inches(0.08), Inches(0.08), col)
            _add_text_box(slide, Emu(x + Inches(0.28)), py,
                          Emu(colW - Inches(0.44)), Inches(0.55),
                          p, BODY_FONT, 10, DARK, line_spacing=13)


def build_quote_full(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    _add_text_box(slide, Inches(1.0), Inches(1.0), Inches(8.0), Inches(2.8),
                  c.get("quote", ""), BODY_FONT, 24, WHITE, italic=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=34)
    _add_rect(slide, Inches(3.5), Inches(4.0), Inches(3.0), Inches(0.04), GREEN_MID)
    if c.get("attribution"):
        _add_text_box(slide, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.4),
                      c["attribution"], BODY_FONT, 13, GREEN_MID, bold=True, align=PP_ALIGN.CENTER)
    if c.get("context"):
        _add_text_box(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.35),
                      c["context"], BODY_FONT, 10, RGBColor(0xCC, 0xCC, 0xCC), italic=True,
                      align=PP_ALIGN.CENTER)


def build_stat_hero(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Key Metric"), size=22)
    hero = c.get("hero", {}); col = _resolve_section_color(c)
    _add_rounded_rect(slide, LM, Inches(1.2), CW, Inches(1.8), OFF_WHITE)
    _add_rect(slide, LM, Inches(1.2), Inches(0.08), Inches(1.8), col)
    _add_text_box(slide, LM, Inches(1.2), CW, Inches(1.3),
                  hero.get("value", "—"), TITLE_FONT, 64, col, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if hero.get("label"):
        _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(7.0), Inches(0.45),
                      hero["label"], BODY_FONT, 14, DARK, bold=True, align=PP_ALIGN.CENTER)
    supporting = c.get("supporting", [])
    n = min(len(supporting), 4)
    if n > 0:
        gap = Inches(0.2); sw = (CW - (n - 1) * gap) / n
        for i, s in enumerate(supporting[:n]):
            x = Emu(LM + i * (sw + gap)); sc = COLORS[i % len(COLORS)]
            _add_rounded_rect(slide, x, Inches(3.25), sw, Inches(1.6), OFF_WHITE)
            _add_rect(slide, x, Inches(3.25), sw, Inches(0.06), sc)
            _add_text_box(slide, x, Inches(3.35), sw, Inches(0.7),
                          s.get("value", ""), TITLE_FONT, 26, sc, bold=True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _add_text_box(slide, x, Inches(4.05), sw, Inches(0.35),
                          s.get("label", ""), BODY_FONT, 10, MID, align=PP_ALIGN.CENTER)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True, align=PP_ALIGN.CENTER)


def build_in_brief_featured(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "In Brief"), size=24)
    _add_rounded_rect(slide, LM, CONTENT_TOP, CW, Inches(1.3), OFF_WHITE)
    _add_rect(slide, LM, CONTENT_TOP, Inches(0.10), Inches(1.3), GOLD)
    _add_text_box(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Emu(CW - Inches(0.5)), Inches(1.3),
                  c.get("featured", ""), BODY_FONT, 16, DARK, bold=True,
                  valign=MSO_ANCHOR.MIDDLE, line_spacing=22)
    for i, s in enumerate(c.get("supporting", [])[:4]):
        y = Emu(Inches(2.7) + i * Inches(0.65)); col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, Inches(0.55), OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.06), Inches(0.55), col)
        _add_text_box(slide, Emu(LM + Inches(0.2)), y, Emu(CW - Inches(0.4)), Inches(0.55),
                      s, BODY_FONT, 12, DARK, valign=MSO_ANCHOR.MIDDLE)


def build_in_brief_reveal(prs, c):
    """Spotlight reveal: each slide features one item large while others stay small."""
    items = c.get("items", [])
    n = min(len(items), 6)
    if n == 0:
        return

    for k in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _accent(slide); _slide_title(slide, c.get("title", "In Brief"), size=24)

        y_cursor = Emu(CONTENT_TOP + Inches(0.12))
        for i in range(n):
            col = COLORS[i % len(COLORS)]
            if i == k:
                # Featured card -- large, bold
                h = Inches(1.2)
                _add_rect(slide, LM, y_cursor, CW, h, OFF_WHITE)
                _add_rect(slide, LM, y_cursor, Inches(0.12), h, col)
                _add_text_box(slide, Emu(LM + Inches(0.25)), y_cursor,
                              Emu(CW - Inches(0.5)), h,
                              items[i], BODY_FONT, 15, DARK, bold=True,
                              valign=MSO_ANCHOR.MIDDLE, line_spacing=21)
            else:
                # Small row
                h = Inches(0.5)
                _add_rect(slide, LM, y_cursor, CW, h, OFF_WHITE)
                _add_rect(slide, LM, y_cursor, Inches(0.06), h, col)
                _add_text_box(slide, Emu(LM + Inches(0.18)), y_cursor,
                              Emu(CW - Inches(0.4)), h,
                              items[i], BODY_FONT, 11, DARK,
                              valign=MSO_ANCHOR.MIDDLE)
            # Extra gap around the featured card
            gap = Inches(0.15) if i == k or (i + 1 < n and i + 1 == k) else Inches(0.08)
            y_cursor = Emu(y_cursor + h + gap)


def build_persona_duo(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Archetype Comparison"), size=22)
    personas = c.get("personas", [{}, {}])
    for idx, p in enumerate(personas[:2]):
        x = LM if idx == 0 else Inches(5.15)
        w = Inches(4.0) if idx == 0 else Inches(4.35)
        col = COLORS[idx % len(COLORS)]
        _add_rounded_rect(slide, x, CONTENT_TOP, w, Inches(3.9), OFF_WHITE)
        _add_rect(slide, x, CONTENT_TOP, Inches(0.08), Inches(3.9), col)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)),
                      Emu(w - Inches(0.3)), Inches(0.4), p.get("name", ""),
                      TITLE_FONT, 18, col, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.5)),
                      Emu(w - Inches(0.3)), Inches(0.25), p.get("archetype", ""),
                      BODY_FONT, 10, GOLD, bold=True)
        for j, t in enumerate(p.get("traits", [])[:4]):
            ty = Emu(CONTENT_TOP + Inches(0.9) + j * Inches(0.38))
            _add_oval(slide, Emu(x + Inches(0.2)), Emu(ty + Inches(0.04)),
                      Inches(0.1), Inches(0.1), col)
            _add_text_box(slide, Emu(x + Inches(0.4)), ty,
                          Emu(w - Inches(0.55)), Inches(0.33), t, BODY_FONT, 10, DARK)
        if p.get("strategy"):
            _add_rect(slide, Emu(x + Inches(0.1)), Emu(CONTENT_TOP + Inches(2.5)),
                      Emu(w - Inches(0.2)), Inches(0.04), LIGHT)
            _add_text_box(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(2.65)),
                          Emu(w - Inches(0.3)), Inches(1.1), p["strategy"],
                          BODY_FONT, 10, MID, line_spacing=14)


def build_process_flow_vertical(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Process"), size=24)
    steps = c.get("steps", []); count = min(len(steps), 3)
    if count == 0: return
    avail = 3.8; stepH = (avail - (count - 1) * 0.4) / count
    for i, step in enumerate(steps[:count]):
        y = Inches(1.05 + i * (stepH + 0.4)); col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, Inches(stepH), OFF_WHITE)
        _add_rect(slide, LM, y, Inches(0.08), Inches(stepH), col)
        _add_text_box(slide, Emu(LM + Inches(0.2)), Emu(y + Inches(0.05)),
                      Inches(0.5), Inches(0.5), str(i + 1),
                      TITLE_FONT, 24, col, bold=True)
        _add_text_box(slide, Emu(LM + Inches(0.75)), Emu(y + Inches(0.05)),
                      Inches(3.0), Inches(0.4), step.get("title", ""),
                      BODY_FONT, 13, DARK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if step.get("detail"):
            _add_text_box(slide, Emu(LM + Inches(0.75)), Emu(y + Inches(0.45)),
                          Inches(7.5), Emu(Inches(stepH) - Inches(0.55)),
                          step["detail"], BODY_FONT, 11, MID, line_spacing=15)
        if i < count - 1:
            _add_text_box(slide, Emu(LM + Inches(0.3)), Emu(y + Inches(stepH + 0.05)),
                          Inches(0.5), Inches(0.25), "\u2193",
                          TITLE_FONT, 18, LIGHT, align=PP_ALIGN.CENTER)


def build_text_cards(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Key Points"), size=24)
    items = c.get("items", []); n = min(len(items), 6)
    if n == 0: return
    cols = 2 if n <= 4 else 3
    rows_count = (n + cols - 1) // cols
    cardW = (CW - Emu((cols - 1) * Inches(0.2))) / cols
    cardH = (Inches(3.8) - (rows_count - 1) * Inches(0.2)) / rows_count
    for i, item in enumerate(items[:n]):
        col_idx = i % cols; row = i // cols
        x = Emu(LM + col_idx * (cardW + Inches(0.2)))
        y = Emu(CONTENT_TOP + row * (cardH + Inches(0.2)))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, x, y, cardW, cardH, OFF_WHITE)
        _add_rect(slide, x, y, cardW, Inches(0.06), col)
        title_text = item if isinstance(item, str) else item.get("title", "")
        detail_text = "" if isinstance(item, str) else item.get("detail", "")
        if detail_text:
            _add_text_box(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.1)),
                          Emu(cardW - Inches(0.24)), Inches(0.35),
                          title_text, BODY_FONT, 12, DARK, bold=True)
            _add_text_box(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.5)),
                          Emu(cardW - Inches(0.24)), Emu(cardH - Inches(0.6)),
                          detail_text, BODY_FONT, 10, MID, line_spacing=14)
        else:
            _add_text_box(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.1)),
                          Emu(cardW - Inches(0.24)), Emu(cardH - Inches(0.2)),
                          title_text, BODY_FONT, 12, DARK, valign=MSO_ANCHOR.MIDDLE, line_spacing=17)


def build_text_columns(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Overview"), size=24)
    columns = c.get("columns", []); n = min(len(columns), 3)
    if n == 0: return
    gap = Inches(0.25); colW = (CW - (n - 1) * gap) / n
    for i, col_data in enumerate(columns[:n]):
        x = Emu(LM + i * (colW + gap)); col = COLORS[i % len(COLORS)]
        heading = "" if isinstance(col_data, str) else col_data.get("heading", "")
        body = col_data if isinstance(col_data, str) else col_data.get("body", "")
        if heading:
            _add_text_box(slide, x, CONTENT_TOP, colW, Inches(0.3),
                          heading, TITLE_FONT, 11, col, bold=True)
            _add_rect(slide, x, Emu(CONTENT_TOP + Inches(0.35)), Inches(0.8), Inches(0.04), col)
            _add_text_box(slide, x, Emu(CONTENT_TOP + Inches(0.5)), colW, Inches(3.3),
                          body, BODY_FONT, 11, DARK, line_spacing=16)
        else:
            _add_text_box(slide, x, CONTENT_TOP, colW, Inches(3.8),
                          body, BODY_FONT, 11, DARK, line_spacing=16)
        if i < n - 1:
            _add_rect(slide, Emu(LM + (i + 1) * colW + i * gap + Inches(0.12)),
                      CONTENT_TOP, Inches(0.02), Inches(3.5), LIGHT)


def build_text_narrative(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Overview"), size=24)
    lede = c.get("lede", "")
    body = c.get("body", "")
    if isinstance(body, list): body = "\n\n".join(body)
    _add_rounded_rect(slide, LM, CONTENT_TOP, CW, Inches(1.2), OFF_WHITE)
    _add_rect(slide, LM, CONTENT_TOP, Inches(0.08), Inches(1.2), GOLD)
    _add_text_box(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Emu(CW - Inches(0.5)), Inches(1.2),
                  lede, BODY_FONT, 15, DARK, bold=True, line_spacing=22,
                  valign=MSO_ANCHOR.MIDDLE)
    _add_text_box(slide, LM, Inches(2.6), CW, Inches(2.5),
                  body, BODY_FONT, 12, MID, line_spacing=18)


def build_text_nested(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Detail"), size=24)
    items = c.get("items", []); n = min(len(items), 4)
    if n == 0: return
    avail = H - CONTENT_TOP - Inches(0.3); gap = Inches(0.12)
    cardH = int((avail - (n - 1) * gap) / n)
    for i, item in enumerate(items[:n]):
        y = Emu(CONTENT_TOP + i * (cardH + gap))
        text = item if isinstance(item, str) else item.get("text", "")
        children = [] if isinstance(item, str) else item.get("children", [])
        col = COLORS[i % len(COLORS)]
        _add_rect(slide, LM, y, Inches(2.2), cardH, col)
        _add_text_box(slide, Emu(LM + Inches(0.12)), y, Inches(1.95), cardH,
                      text, BODY_FONT, 11, WHITE, bold=True,
                      valign=MSO_ANCHOR.MIDDLE, line_spacing=15)
        _add_rounded_rect(slide, Inches(3.25), y, Inches(6.25), cardH, OFF_WHITE)
        child_texts = []
        for ch in children:
            ch_text = ch if isinstance(ch, str) else ch.get("text", "")
            grandchildren = [] if isinstance(ch, str) else ch.get("children", [])
            child_texts.append(ch_text)
            for gc in grandchildren:
                gc_text = gc if isinstance(gc, str) else gc.get("text", "")
                child_texts.append("  \u2022 " + gc_text)
        _add_text_box(slide, Inches(3.38), y, Inches(5.95), cardH,
                      "\n".join(child_texts), BODY_FONT, 10, DARK,
                      valign=MSO_ANCHOR.MIDDLE, line_spacing=14)


def build_text_split(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Key Point"), size=24)
    headline = c.get("headline", ""); detail = c.get("detail", "")
    points = c.get("points", [])
    _add_rounded_rect(slide, LM, CONTENT_TOP, Inches(4.0), Inches(3.8), OFF_WHITE)
    _add_rect(slide, LM, CONTENT_TOP, Inches(0.08), Inches(3.8), GREEN)
    _add_text_box(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Inches(3.5), Inches(2.0),
                  headline, TITLE_FONT, 18, DARK, bold=True, line_spacing=26,
                  valign=MSO_ANCHOR.MIDDLE)
    if detail:
        _add_text_box(slide, Emu(LM + Inches(0.25)), Inches(3.0), Inches(3.5), Inches(1.5),
                      detail, BODY_FONT, 11, MID, line_spacing=16)
    for i, p in enumerate(points[:6]):
        py = Emu(CONTENT_TOP + Inches(0.05) + i * Inches(0.6))
        col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, Inches(5.15), py, Inches(4.35), Inches(0.5), OFF_WHITE)
        _add_rect(slide, Inches(5.15), py, Inches(0.06), Inches(0.5), col)
        _add_text_box(slide, Inches(5.3), py, Inches(4.05), Inches(0.5),
                      p, BODY_FONT, 11, DARK, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)


def build_text_annotated(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Analysis"), size=24)
    items = c.get("items", []); n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.85))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))
    for i, item in enumerate(items[:n]):
        y = Emu(startY + i * (rowH + gap)); col = COLORS[i % len(COLORS)]
        _add_rounded_rect(slide, LM, y, CW, rowH, OFF_WHITE)
        _add_rect(slide, LM, y, Inches(1.5), rowH, col)
        _add_text_box(slide, LM, y, Inches(1.5), rowH,
                      item.get("label", ""), BODY_FONT, 10, WHITE, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Emu(LM + Inches(1.7)), y, Emu(CW - Inches(1.9)), rowH,
                      item.get("text", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


def build_icon_cards(prs, c):
    """Row of 2-3 cards, each with an icon/image above and title+detail below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Key Points"), size=24)
    items = c.get("items", []); n = min(len(items), 3)
    if n == 0: return
    gap = Inches(0.25)
    cardW = Emu((CW - (n - 1) * gap) / n)
    iconSize = Inches(0.8)
    iconTop = CONTENT_TOP
    cardTop = Emu(CONTENT_TOP + Inches(1.1))
    cardH = Emu(Inches(3.8) - Inches(1.1))
    for i, item in enumerate(items[:n]):
        x = Emu(LM + i * (cardW + gap))
        col = COLORS[i % len(COLORS)]
        icon_x = Emu(x + (cardW - iconSize) // 2)
        img_path = item.get("imagePath", "") if isinstance(item, dict) else ""
        shape_type = DEFAULT_ICON_SHAPES[i % len(DEFAULT_ICON_SHAPES)]
        _add_icon_or_image(slide, img_path, icon_x, iconTop, iconSize, iconSize, shape_type, col)
        _add_rect(slide, x, cardTop, cardW, cardH, OFF_WHITE)
        title_text = item.get("title", "") if isinstance(item, dict) else str(item)
        detail_text = item.get("detail", "") if isinstance(item, dict) else ""
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(cardTop + Inches(0.12)),
                      Emu(cardW - Inches(0.3)), Inches(0.4),
                      title_text, BODY_FONT, 13, DARK, bold=True)
        if detail_text:
            _add_text_box(slide, Emu(x + Inches(0.15)), Emu(cardTop + Inches(0.5)),
                          Emu(cardW - Inches(0.3)), Emu(cardH - Inches(0.6)),
                          detail_text, BODY_FONT, 10, MID, line_spacing=14)


def build_feature_cards(prs, c):
    """1-2 full-width rows with an icon/image on the left and text on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _accent(slide); _slide_title(slide, c.get("title", "Features"), size=24)
    items = c.get("items", []); n = min(len(items), 2)
    if n == 0: return
    gap = Inches(0.25)
    rowH = Emu((Inches(3.8) - (n - 1) * gap) / n)
    iconSize = Inches(1.1)
    for i, item in enumerate(items[:n]):
        y = Emu(CONTENT_TOP + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        _add_rect(slide, LM, y, CW, rowH, OFF_WHITE)
        icon_y = Emu(y + (rowH - iconSize) // 2)
        icon_x = Emu(LM + Inches(0.3))
        img_path = item.get("imagePath", "") if isinstance(item, dict) else ""
        shape_type = DEFAULT_ICON_SHAPES[i % len(DEFAULT_ICON_SHAPES)]
        _add_icon_or_image(slide, img_path, icon_x, icon_y, iconSize, iconSize, shape_type, col)
        text_x = Emu(LM + Inches(1.7))
        text_w = Emu(CW - Inches(1.9))
        title_text = item.get("title", "") if isinstance(item, dict) else str(item)
        detail_text = item.get("detail", "") if isinstance(item, dict) else ""
        _add_text_box(slide, text_x, Emu(y + Inches(0.2)),
                      text_w, Inches(0.4),
                      title_text, BODY_FONT, 14, DARK, bold=True)
        if detail_text:
            _add_text_box(slide, text_x, Emu(y + Inches(0.65)),
                          text_w, Emu(rowH - Inches(0.8)),
                          detail_text, BODY_FONT, 11, MID, line_spacing=16)


BUILDERS = {
    "title": build_title, "in_brief": build_in_brief,
    "section_divider": build_section_divider, "stat_callout": build_stat_callout,
    "quote": build_quote, "comparison": build_comparison,
    "text_graph": build_text_graph, "process_flow": build_process_flow,
    "matrix": build_matrix, "methods": build_methods,
    "hypotheses": build_hypotheses, "wsn_dense": build_wsn_dense,
    "wsn_reveal": build_wsn_reveal, "findings_recs": build_findings_recs,
    "findings_recs_dense": build_findings_recs_dense,
    "open_questions": build_open_questions, "agenda": build_agenda,
    "progressive_reveal": build_progressive_reveal, "closer": build_closer,
    "timeline": build_timeline, "data_table": build_data_table,
    "multi_stat": build_multi_stat, "persona": build_persona,
    "risk_tradeoff": build_risk_tradeoff, "appendix": build_appendix,
    "before_after": build_before_after, "summary": build_summary,
    "quote_full": build_quote_full, "stat_hero": build_stat_hero,
    "in_brief_featured": build_in_brief_featured, "in_brief_reveal": build_in_brief_reveal,
    "persona_duo": build_persona_duo,
    "process_flow_vertical": build_process_flow_vertical,
    "process_flow_reveal": build_process_flow_reveal,
    "text_cards": build_text_cards, "text_columns": build_text_columns,
    "text_narrative": build_text_narrative, "text_nested": build_text_nested,
    "text_split": build_text_split, "text_annotated": build_text_annotated,
    "icon_cards": build_icon_cards, "feature_cards": build_feature_cards,
}

def build_deck(slide_configs, output_path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    for slide_type, data in slide_configs:
        builder = BUILDERS.get(slide_type)
        if builder: builder(prs, data)
    prs.save(output_path); return output_path
