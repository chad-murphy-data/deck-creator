"""
Editorial v2 — Complete template (21 builders)
Purpose-driven slide layouts where each builder's visual treatment
matches its rhetorical job. Not a reskin — a reimagining.

BUILDERS:
Structural:
  1. build_title — dark green, gold rule, clean footer
  2. build_section_divider — dark green, breadcrumbs + hero + upcoming
  3. build_closer — dark green, bookend with title, italic closing thought
  4. build_agenda — white, newspaper masthead, equal-weight columns

Single-point:
  5. build_stat_callout — white, giant number left, context right
  6. build_quote_full — dark green, large serif italic, beat of silence

Comparison:
  7. build_comparison — white, row-based clean cards, warm/cool tinting
  8. build_comparison_reveal — white, 2-slide asymmetric reveal

Progressive:
  9. build_in_brief — white, progressive reveal, each bullet gets hero treatment
  10. build_wsn_reveal — 3-slide What/So What/Now What with sidebar compression

Analysis:
  11. build_findings_recs — finding→recommendation paired rows
  12. build_process_flow — horizontal cards with colored top bars (<=4 steps)
  13. build_process_flow_vertical — film strip bands (5+ steps)
  14. build_process_flow_accordion — alternating left/right with center spine
  15. build_timeline — route map with alternating stops

Character/Text:
  16. build_persona — magazine profile (single or duo), image placeholder
  17. build_text_narrative — lede + two-column body in warm/cool cards

Wrap-up:
  18. build_open_questions — 2x2 provocation cards with watermark ?
  19. build_data_table — styled table with highlighted column
  20. build_summary — 2-slide reveal (sections → takeaways)
  21. build_appendix — intentionally muted endnotes

Design language:
- White bg for content slides, DK_GREEN for structural slides (title/divider/closer)
- Gold top rule on dark slides, DK_GREEN top rule on white slides
- Running head: title.upper() in BODY_FONT 9pt QUIET (or 18pt for persona)
- Fidelity Slab for titles/headlines, Fidelity Sans for body
- Warm cards (#F0EDE6) for left/primary, Cool cards (#EAEEF0) for right/secondary
- ACCENTS cycle for item markers (DK_GREEN, COBALT, PURPLE, GOLD, TEAL, RED, OCHRE)
- Left-edge accent bars as primary visual markers
- Gold for emphasis, separators, and numbers
- Image placeholders rendered as dark tinted rectangles with "IMG" label
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

W = Inches(10)
H = Inches(5.625)

# ── Palette ──────────────────────────────────────────────
DK_GREEN    = RGBColor(0x04, 0x40, 0x14)
GOLD        = RGBColor(0xD4, 0xA8, 0x43)
CHARCOAL    = RGBColor(0x2D, 0x2D, 0x2D)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
MID         = RGBColor(0x66, 0x66, 0x66)
QUIET       = RGBColor(0x99, 0x99, 0x99)
FAINT       = RGBColor(0xCC, 0xCC, 0xCC)
RULE_CLR    = RGBColor(0xDD, 0xDD, 0xDD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM       = RGBColor(0xE8, 0xE0, 0xCC)
WARM_CARD   = RGBColor(0xF0, 0xED, 0xE6)
COOL_CARD   = RGBColor(0xEA, 0xEE, 0xF0)
LIGHT_BOX   = RGBColor(0xF7, 0xF6, 0xF4)
COBALT      = RGBColor(0x04, 0x54, 0x7C)
LIGHT_GREEN = RGBColor(0xA8, 0xC4, 0xA0)

TITLE_FONT = "Fidelity Slab"
BODY_FONT  = "Fidelity Sans"
LM = Inches(0.65)
CW = Emu(W - Inches(1.3))

ACCENTS = [
    RGBColor(0x04, 0x40, 0x14), RGBColor(0x04, 0x54, 0x7C),
    RGBColor(0x5B, 0x2C, 0x8F), RGBColor(0xD4, 0xA8, 0x43),
    RGBColor(0x1B, 0x7A, 0x6E), RGBColor(0xC2, 0x3B, 0x22),
    RGBColor(0xCC, 0x7A, 0x2E),
]


# ── Helpers ──────────────────────────────────────────────
def _rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def _rect_outline(slide, x, y, w, h, line_color, line_width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape

def _oval(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def _tb(slide, x, y, w, h, text, font_name, font_size, color,
        bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
        line_spacing=None):
    txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = txBox.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font_name; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.italic = italic
    p.alignment = align; txBox.text_frame.auto_size = None
    bodyPr = txBox.text_frame._txBody.bodyPr
    anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
    bodyPr.set("anchor", anchor_map.get(valign, "t"))
    if line_spacing:
        pPr = p._p.get_or_add_pPr()
        spcElem = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPts = etree.SubElement(spcElem, qn("a:spcPts"))
        spcPts.set("val", str(int(line_spacing * 100)))
    return txBox

def _line(slide, x, y, w, color, thickness=1):
    _rect(slide, x, y, w, Pt(thickness), color)

def _set_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def _running_head(slide, title):
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, LM, Inches(0.2), Inches(7), Inches(0.4),
        title.upper(), BODY_FONT, 14, QUIET, bold=True)

def _parse_item(item):
    if isinstance(item, dict):
        return item.get("title", item.get("label", "")), item.get("detail", "")
    return str(item), ""


# ══════════════════════════════════════════════════════════
# 1. TITLE
# ══════════════════════════════════════════════════════════
def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DK_GREEN)
    title = c.get("title", ""); subtitle = c.get("subtitle", "")
    author = c.get("author", ""); date = c.get("date", "")
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), GOLD)
    tsize = 32 if len(title) > 60 else 40
    _tb(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(2.5),
        title, TITLE_FONT, tsize, WHITE, bold=True,
        line_spacing=tsize * 1.25, valign=MSO_ANCHOR.BOTTOM)
    _line(slide, Inches(0.8), Inches(3.3), Inches(2.0), GOLD, 3)
    if subtitle:
        _tb(slide, Inches(0.8), Inches(3.55), Inches(8.0), Inches(1.0),
            subtitle, BODY_FONT, 14, CREAM, line_spacing=20)
    footer_y = Inches(4.85)
    if author:
        _tb(slide, Inches(0.8), footer_y, Inches(4), Inches(0.35),
            author, BODY_FONT, 10, LIGHT_GREEN)
    if date:
        _tb(slide, Inches(0.8), Emu(footer_y + Inches(0.3)), Inches(4), Inches(0.3),
            date, BODY_FONT, 10, LIGHT_GREEN)


# ══════════════════════════════════════════════════════════
# 2. SECTION DIVIDER (with currentSection support)
# ══════════════════════════════════════════════════════════
def build_section_divider(prs, c, current=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DK_GREEN)
    title = c.get("title", ""); items = c.get("items", [])
    n = len(items)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), GOLD)
    # Breadcrumbs
    completed = items[:current]
    if completed:
        parts = []
        for i, item in enumerate(completed):
            t, _ = _parse_item(item)
            parts.append(f"{i + 1}. {t}")
        breadcrumb = "     ".join(parts)
        _tb(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.35),
            breadcrumb, BODY_FONT, 9, LIGHT_GREEN, valign=MSO_ANCHOR.MIDDLE)
        _line(slide, Inches(0.3), Inches(0.65), Inches(9.4),
              RGBColor(0x3A, 0x6A, 0x3A), 1)
    # Hero
    if current < n:
        curr_title, curr_detail = _parse_item(items[current])
        hero_y = Inches(0.9) if completed else Inches(0.5)
        _tb(slide, Inches(0.5), hero_y, Inches(1), Inches(0.7),
            str(current + 1), TITLE_FONT, 48, GOLD, bold=True)
        _tb(slide, Inches(0.5), Emu(hero_y + Inches(0.75)), Inches(8.5), Inches(1.2),
            curr_title, TITLE_FONT, 32, WHITE, bold=True, line_spacing=40)
        _line(slide, Inches(0.5), Emu(hero_y + Inches(2.0)), Inches(2.0), GOLD, 3)
        if curr_detail:
            _tb(slide, Inches(0.5), Emu(hero_y + Inches(2.25)), Inches(6), Inches(0.4),
                curr_detail, BODY_FONT, 13, CREAM, italic=True)
    # Upcoming — monochrome
    upcoming = items[current + 1:]
    if upcoming:
        up_start = Inches(3.9) if completed else Inches(3.5)
        _line(slide, Inches(0.3), Emu(up_start - Inches(0.15)), Inches(9.4),
              RGBColor(0x3A, 0x6A, 0x3A), 1)
        up_n = len(upcoming)
        item_w = Inches(9.0) / min(up_n, 4)
        for i, item in enumerate(upcoming[:4]):
            item_title, _ = _parse_item(item)
            idx = current + 1 + i
            x = Emu(Inches(0.5) + i * item_w)
            _tb(slide, x, up_start, Inches(0.3), Inches(0.3),
                str(idx + 1), TITLE_FONT, 14, GOLD, bold=True)
            _tb(slide, Emu(x + Inches(0.35)), up_start, Emu(item_w - Inches(0.5)), Inches(0.65),
                item_title, BODY_FONT, 11, CREAM, line_spacing=15)
        if up_n > 4:
            _tb(slide, Inches(0.5), Emu(up_start + Inches(0.7)), Inches(2), Inches(0.25),
                f"+{up_n - 4} more", BODY_FONT, 9, LIGHT_GREEN)


# ══════════════════════════════════════════════════════════
# 3. CLOSER
# ══════════════════════════════════════════════════════════
def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DK_GREEN)
    title = c.get("title", ""); subtitle = c.get("subtitle", "")
    contact = c.get("contact", "")
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), GOLD)
    _tb(slide, Inches(0.8), Inches(0.8), Inches(8.4), Inches(1.2),
        title, TITLE_FONT, 36, WHITE, bold=True,
        line_spacing=44, valign=MSO_ANCHOR.BOTTOM)
    _line(slide, Inches(0.8), Inches(2.15), Inches(2.0), GOLD, 3)
    if subtitle:
        _tb(slide, Inches(0.8), Inches(2.5), Inches(8.0), Inches(1.8),
            subtitle, TITLE_FONT, 18, CREAM, italic=True, line_spacing=26)
    if contact:
        _tb(slide, Inches(0.8), Inches(4.9), Inches(8), Inches(0.4),
            contact, BODY_FONT, 10, LIGHT_GREEN)


# ══════════════════════════════════════════════════════════
# 4. AGENDA (newspaper layout)
# ══════════════════════════════════════════════════════════
def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", "Agenda"); items = c.get("items", [])
    n = len(items)
    # Masthead
    _tb(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.4),
        title.upper(), BODY_FONT, 9, QUIET, bold=True, align=PP_ALIGN.CENTER)
    _line(slide, Inches(0.3), Inches(0.55), Inches(9.4), CHARCOAL, 2)
    _line(slide, Inches(0.3), Inches(0.6), Inches(9.4), CHARCOAL, 1)
    _tb(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.7),
        "Agenda", TITLE_FONT, 36, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _line(slide, Inches(0.3), Inches(1.6), Inches(9.4), RULE_CLR, 1)
    if not items: return
    cols = 2 if n <= 4 else min(3, n)
    col_gap = Inches(0.2)
    total_w = Inches(9.0)
    col_w = Emu((total_w - (cols - 1) * col_gap) / cols)
    start_x = Inches(0.5); start_y = Inches(1.85); row_h = Inches(1.25)
    for i, item in enumerate(items):
        item_title, item_detail = _parse_item(item)
        col = i % cols; row = i // cols
        color = ACCENTS[i % len(ACCENTS)]
        x = Emu(start_x + col * (col_w + col_gap))
        y = Emu(start_y + row * row_h)
        _rect(slide, x, y, Inches(0.04), Inches(0.95), color)
        _tb(slide, Emu(x + Inches(0.12)), y, Inches(0.35), Inches(0.35),
            str(i + 1), TITLE_FONT, 16, color, bold=True)
        _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.32)),
            Emu(col_w - Inches(0.2)), Inches(0.35),
            item_title, TITLE_FONT, 13, CHARCOAL, bold=True, line_spacing=16)
        if item_detail:
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.68)),
                Emu(col_w - Inches(0.2)), Inches(0.25),
                item_detail, BODY_FONT, 9, QUIET, line_spacing=12)
        if col < cols - 1 and (i + 1) < n:
            div_x = Emu(x + col_w + col_gap // 2)
            _rect(slide, div_x, y, Inches(0.01), Inches(0.95), RULE_CLR)


# ══════════════════════════════════════════════════════════
# 5. STAT CALLOUT (number as hero)
# ══════════════════════════════════════════════════════════
def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    stat = c.get("stat", ""); headline = c.get("headline", "")
    detail = c.get("detail", ""); source = c.get("source", ""); title = c.get("title", "")
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    if title:
        _tb(slide, LM, Inches(0.2), Inches(4), Inches(0.4),
            title.upper(), BODY_FONT, 14, QUIET, bold=True)
    _tb(slide, LM, Inches(0.6), Inches(5.2), Inches(4.0),
        stat, TITLE_FONT, 120, DK_GREEN, bold=True, valign=MSO_ANCHOR.MIDDLE)
    rx = Inches(5.8); rw = Inches(3.6)
    _tb(slide, rx, Inches(1.2), rw, Inches(1.0),
        headline, TITLE_FONT, 22, CHARCOAL, bold=True, line_spacing=28)
    _line(slide, rx, Inches(2.3), Inches(1.2), GOLD, 2)
    if detail:
        _tb(slide, rx, Inches(2.55), rw, Inches(2.2),
            detail, BODY_FONT, 11, MID, line_spacing=16)
    if source:
        _tb(slide, rx, Inches(5.0), rw, Inches(0.35),
            source, BODY_FONT, 8, QUIET, italic=True)


# ══════════════════════════════════════════════════════════
# 6. QUOTE FULL (beat of silence)
# ══════════════════════════════════════════════════════════
def build_quote_full(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DK_GREEN)
    quote = c.get("quote", ""); attribution = c.get("attribution", "")
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), GOLD)
    _tb(slide, Inches(0.8), Inches(0.4), Inches(1.5), Inches(1.5),
        "\u201C", TITLE_FONT, 144, GOLD, bold=True)
    if len(quote) > 200: qsize, qsp = 20, 28
    elif len(quote) > 120: qsize, qsp = 24, 32
    else: qsize, qsp = 28, 36
    _tb(slide, Inches(1.2), Inches(1.4), Inches(7.6), Inches(2.8),
        quote, TITLE_FONT, qsize, WHITE, italic=True,
        valign=MSO_ANCHOR.MIDDLE, line_spacing=qsp)
    if attribution:
        _line(slide, Inches(1.2), Inches(4.35), Inches(0.8), GOLD, 2)
        _tb(slide, Inches(2.2), Inches(4.15), Inches(5), Inches(0.45),
            attribution, BODY_FONT, 12, GOLD, bold=True)
    context = c.get("context", "")
    if context:
        _tb(slide, Inches(1.2), Inches(4.75), Inches(7.6), Inches(0.6),
            context, BODY_FONT, 10, LIGHT_GREEN, italic=True, line_spacing=14)


# ══════════════════════════════════════════════════════════
# 7. COMPARISON (clean rows, V1)
# ══════════════════════════════════════════════════════════
def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", "")
    left_label = c.get("leftLabel", ""); right_label = c.get("rightLabel", "")
    left_items = c.get("leftItems", []); right_items = c.get("rightItems", [])
    _running_head(slide, title)
    left_x = Inches(0.3); right_x = Inches(5.15)
    left_w = Inches(4.7); right_w = Inches(4.55)
    _tb(slide, Emu(left_x + Inches(0.25)), Inches(0.7), Inches(4), Inches(0.4),
        left_label, TITLE_FONT, 15, DK_GREEN, bold=True)
    _tb(slide, Emu(right_x + Inches(0.25)), Inches(0.7), Inches(4), Inches(0.4),
        right_label, TITLE_FONT, 15, COBALT, bold=True)
    _line(slide, Emu(left_x + Inches(0.25)), Inches(1.12), Inches(1.5), DK_GREEN, 2)
    _line(slide, Emu(right_x + Inches(0.25)), Inches(1.12), Inches(1.5), COBALT, 2)
    n = min(len(left_items), len(right_items))
    row_h = Inches(0.82); gap_y = Inches(0.1); start_y = Inches(1.3)
    for i in range(n):
        y = Emu(start_y + i * (row_h + gap_y))
        _rect(slide, left_x, y, left_w, row_h, WARM_CARD)
        _rect(slide, left_x, y, Inches(0.06), row_h, DK_GREEN)
        _rect(slide, right_x, y, right_w, row_h, COOL_CARD)
        _rect(slide, right_x, y, Inches(0.06), row_h, COBALT)
        _tb(slide, Emu(left_x + Inches(0.25)), y, Emu(left_w - Inches(0.45)), row_h,
            left_items[i], BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=16)
        _tb(slide, Emu(right_x + Inches(0.25)), y, Emu(right_w - Inches(0.45)), row_h,
            right_items[i], BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


# ══════════════════════════════════════════════════════════
# 8. COMPARISON REVEAL (2-slide asymmetric)
# ══════════════════════════════════════════════════════════
def build_comparison_reveal(prs, c):
    title = c.get("title", "")
    left_label = c.get("leftLabel", ""); right_label = c.get("rightLabel", "")
    left_items = c.get("leftItems", []); right_items = c.get("rightItems", [])
    n = min(len(left_items), len(right_items))
    row_h = Inches(0.82); gap_y = Inches(0.1); start_y = Inches(1.5)

    # Slide 1: LEFT dominant
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s1, WHITE); _running_head(s1, title)
    _tb(s1, LM, Inches(0.7), Inches(5.5), Inches(0.5),
        left_label, TITLE_FONT, 20, DK_GREEN, bold=True)
    _line(s1, LM, Inches(1.25), Inches(1.8), DK_GREEN, 2)
    rx_small = Inches(7.0); rw_small = Inches(2.6)
    _tb(s1, rx_small, Inches(0.75), rw_small, Inches(0.35),
        right_label, BODY_FONT, 11, QUIET, bold=True)
    _line(s1, rx_small, Inches(1.12), Inches(1.0), RULE_CLR, 1)
    _rect(s1, Inches(6.65), Inches(1.3), Inches(0.015), Emu(n * (row_h + gap_y)), RULE_CLR)
    for i in range(n):
        y = Emu(start_y + i * (row_h + gap_y))
        _rect(s1, LM, y, Inches(5.7), row_h, WARM_CARD)
        _rect(s1, LM, y, Inches(0.06), row_h, GOLD)
        _tb(s1, Emu(LM + Inches(0.25)), y, Inches(5.2), row_h,
            left_items[i], BODY_FONT, 13, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=17)
        _rect(s1, rx_small, y, rw_small, row_h, LIGHT_BOX)
        _rect(s1, rx_small, y, Inches(0.035), row_h, FAINT)
        _tb(s1, Emu(rx_small + Inches(0.15)), y, Emu(rw_small - Inches(0.25)), row_h,
            right_items[i], BODY_FONT, 9, QUIET, valign=MSO_ANCHOR.MIDDLE, line_spacing=12)

    # Slide 2: RIGHT dominant
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s2, WHITE); _running_head(s2, title)
    lx_small = LM; lw_small = Inches(2.6)
    _tb(s2, lx_small, Inches(0.75), lw_small, Inches(0.35),
        left_label, BODY_FONT, 11, QUIET, bold=True)
    _line(s2, lx_small, Inches(1.12), Inches(1.0), RULE_CLR, 1)
    rx_big = Inches(3.65)
    _tb(s2, rx_big, Inches(0.7), Inches(5.5), Inches(0.5),
        right_label, TITLE_FONT, 20, COBALT, bold=True)
    _line(s2, rx_big, Inches(1.25), Inches(1.8), COBALT, 2)
    _rect(s2, Inches(3.35), Inches(1.3), Inches(0.015), Emu(n * (row_h + gap_y)), RULE_CLR)
    for i in range(n):
        y = Emu(start_y + i * (row_h + gap_y))
        _rect(s2, lx_small, y, lw_small, row_h, LIGHT_BOX)
        _rect(s2, lx_small, y, Inches(0.035), row_h, FAINT)
        _tb(s2, Emu(lx_small + Inches(0.15)), y, Emu(lw_small - Inches(0.25)), row_h,
            left_items[i], BODY_FONT, 9, QUIET, valign=MSO_ANCHOR.MIDDLE, line_spacing=12)
        _rect(s2, rx_big, y, Inches(5.7), row_h, COOL_CARD)
        _rect(s2, rx_big, y, Inches(0.06), row_h, COBALT)
        _tb(s2, Emu(rx_big + Inches(0.25)), y, Inches(5.2), row_h,
            right_items[i], BODY_FONT, 13, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=17)


# ══════════════════════════════════════════════════════════
# 9. IN BRIEF (progressive reveal)
# ══════════════════════════════════════════════════════════
def build_in_brief(prs, c):
    title = c.get("title", ""); bullets = c.get("bullets", [])
    if not bullets: return
    for active_idx in range(len(bullets)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, WHITE)
        _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
        _tb(slide, LM, Inches(0.2), Inches(6), Inches(0.4),
            title.upper(), BODY_FONT, 14, QUIET, bold=True)
        progress = f"{active_idx + 1} of {len(bullets)}"
        _tb(slide, Inches(8.0), Inches(0.2), Inches(1.5), Inches(0.4),
            progress, BODY_FONT, 14, GOLD, bold=True, align=PP_ALIGN.RIGHT)
        bar_w = CW; bar_y = Inches(0.6)
        _rect(slide, LM, bar_y, bar_w, Inches(0.03), RULE_CLR)
        filled_w = Emu(int(bar_w) * (active_idx + 1) / len(bullets))
        _rect(slide, LM, bar_y, filled_w, Inches(0.03), DK_GREEN)
        _tb(slide, LM, Inches(0.9), CW, Inches(1.8),
            bullets[active_idx], TITLE_FONT, 28, BLACK, bold=True,
            line_spacing=36, valign=MSO_ANCHOR.TOP)
        _line(slide, LM, Inches(2.8), Inches(2.5), GOLD, 3)
        if active_idx > 0:
            _tb(slide, LM, Inches(3.1), Inches(3), Inches(0.25),
                "PREVIOUSLY", BODY_FONT, 8, QUIET, bold=True)
            for j in range(active_idx):
                y = Emu(Inches(3.4) + j * Inches(0.42))
                _tb(slide, LM, y, Inches(0.3), Inches(0.35),
                    str(j + 1), TITLE_FONT, 12, FAINT, bold=True)
                prev_text = bullets[j]
                if len(prev_text) > 80: prev_text = prev_text[:77] + "..."
                _tb(slide, Emu(LM + Inches(0.35)), y, Inches(8.0), Inches(0.35),
                    prev_text, BODY_FONT, 10, QUIET, valign=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════
# BUILDERS dict (for build_deck compatibility)
# ══════════════════════════════════════════════════════════
BUILDERS = {
    "title": build_title,
    "section_divider": build_section_divider,
    "closer": build_closer,
    "agenda": build_agenda,
    "stat_callout": build_stat_callout,
    "quote_full": build_quote_full,
    "comparison": build_comparison,
    "comparison_reveal": build_comparison_reveal,
    "in_brief": build_in_brief,
}

# ══════════════════════════════════════════════════════════
# BATCH 2 BUILDERS (added after initial lock-in)
# ══════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════
# 10. WSN REVEAL (3 progressive slides)
# ══════════════════════════════════════════════════════════
def build_wsn_reveal(prs, c):
    title = c.get("title", "")
    zones = [
        ("What We Found", DK_GREEN, c.get("what", {})),
        ("So What", COBALT, c.get("soWhat", {})),
        ("Now What", RGBColor(0x5B, 0x2C, 0x8F), c.get("nowWhat", {})),
    ]
    for active in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, WHITE); _running_head(slide, title)
        step_text = f"Step {active + 1} of 3"
        _tb(slide, Inches(8.0), Inches(0.25), Inches(1.5), Inches(0.3),
            step_text, BODY_FONT, 9, GOLD, bold=True, align=PP_ALIGN.RIGHT)
        for d in range(3):
            dx = Emu(Inches(8.2) + d * Inches(0.25))
            dot_color = zones[d][1] if d <= active else FAINT
            _oval(slide, dx, Inches(0.55), Inches(0.12), Inches(0.12), dot_color)
        sidebar_w = Inches(2.3)
        if active > 0:
            for j in range(active):
                label, color, data = zones[j]
                sy = Emu(Inches(0.9) + j * Inches(1.5))
                _rect(slide, LM, sy, Inches(0.04), Inches(1.2), color)
                _tb(slide, Emu(LM + Inches(0.15)), sy, sidebar_w, Inches(0.25),
                    label.upper(), BODY_FONT, 8, color, bold=True)
                _tb(slide, Emu(LM + Inches(0.15)), Emu(sy + Inches(0.28)),
                    sidebar_w, Inches(0.35),
                    data.get("headline", ""), BODY_FONT, 10, MID, bold=True, line_spacing=13)
                detail = data.get("detail", "")
                if len(detail) > 100: detail = detail[:97] + "..."
                _tb(slide, Emu(LM + Inches(0.15)), Emu(sy + Inches(0.65)),
                    sidebar_w, Inches(0.5), detail, BODY_FONT, 8, QUIET, line_spacing=11)
            sep_x = Emu(LM + sidebar_w + Inches(0.2))
            _rect(slide, sep_x, Inches(0.9), Inches(0.015), Inches(4.0), RULE_CLR)
        label, color, data = zones[active]
        hero_x = Emu(LM + sidebar_w + Inches(0.5)) if active > 0 else LM
        hero_w = Emu(Inches(9.0) - sidebar_w - Inches(0.5)) if active > 0 else CW
        _tb(slide, hero_x, Inches(0.85), hero_w, Inches(0.35),
            label.upper(), BODY_FONT, 10, color, bold=True)
        _line(slide, hero_x, Inches(1.22), Inches(1.5), color, 3)
        _tb(slide, hero_x, Inches(1.45), hero_w, Inches(1.2),
            data.get("headline", ""), TITLE_FONT, 26, BLACK, bold=True, line_spacing=34)
        _line(slide, hero_x, Inches(2.75), Inches(2.0), GOLD, 2)
        detail = data.get("detail", "")
        if detail:
            _tb(slide, hero_x, Inches(3.0), hero_w, Inches(2.0),
                detail, BODY_FONT, 12, MID, line_spacing=17)


# ══════════════════════════════════════════════════════════
# 11. FINDINGS & RECS
# ══════════════════════════════════════════════════════════
def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); items = c.get("items", [])
    n = min(len(items), 5)
    _running_head(slide, title)
    _tb(slide, Inches(0.3), Inches(0.65), Inches(4), Inches(0.35),
        "FINDING", BODY_FONT, 9, DK_GREEN, bold=True)
    _tb(slide, Inches(5.4), Inches(0.65), Inches(4), Inches(0.35),
        "RECOMMENDATION", BODY_FONT, 9, COBALT, bold=True)
    _line(slide, Inches(0.3), Inches(1.0), Inches(9.4), CHARCOAL, 1)
    avail = Inches(4.2); gap = Inches(0.1)
    row_h = Emu(min(int((avail - (n - 1) * gap) / n), Inches(0.85)))
    start_y = Inches(1.15); find_w = Inches(4.3); rec_w = Inches(4.3)
    for i, item in enumerate(items[:n]):
        y = Emu(start_y + i * (row_h + gap))
        finding = item.get("finding", "")
        rec = item.get("recommendation", item.get("rec", ""))
        _rect(slide, Inches(0.3), y, find_w, row_h, WARM_CARD)
        _rect(slide, Inches(0.3), y, Inches(0.05), row_h, DK_GREEN)
        _tb(slide, Inches(0.5), y, Emu(find_w - Inches(0.3)), row_h,
            finding, BODY_FONT, 11, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)
        _tb(slide, Inches(4.75), y, Inches(0.5), row_h,
            "\u2192", TITLE_FONT, 18, GOLD, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _rect(slide, Inches(5.4), y, rec_w, row_h, COOL_CARD)
        _rect(slide, Inches(5.4), y, Inches(0.05), row_h, COBALT)
        _tb(slide, Inches(5.6), y, Emu(rec_w - Inches(0.3)), row_h,
            rec, BODY_FONT, 11, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)
    _line(slide, Inches(0.3), Emu(start_y + n * (row_h + gap)), Inches(9.4), CHARCOAL, 1)


# ══════════════════════════════════════════════════════════
# 12. PROCESS FLOW (horizontal, <=4 steps)
# ══════════════════════════════════════════════════════════
def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); steps = c.get("steps", [])
    n = len(steps)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.55),
        title, TITLE_FONT, 22, CHARCOAL, bold=True)
    if n == 0: return
    gap = Inches(0.35); total_w = Inches(9.0)
    card_w = Emu((total_w - (n - 1) * gap) / n)
    card_h = Inches(3.8); start_x = Inches(0.5); start_y = Inches(1.0)
    for i, step in enumerate(steps):
        color = ACCENTS[i % len(ACCENTS)]
        label = step.get("label", step.get("title", ""))
        detail = step.get("detail", "")
        x = Emu(start_x + i * (card_w + gap))
        _rect(slide, x, start_y, card_w, Inches(0.06), color)
        circ_size = Inches(0.42); circ_x = Emu(x + Inches(0.2))
        circ_y = Emu(start_y + Inches(0.25))
        _oval(slide, circ_x, circ_y, circ_size, circ_size, color)
        _tb(slide, circ_x, circ_y, circ_size, circ_size,
            str(i + 1), BODY_FONT, 14, WHITE, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(x + Inches(0.2)), Emu(start_y + Inches(0.85)),
            Emu(card_w - Inches(0.4)), Inches(0.6),
            label, TITLE_FONT, 16, CHARCOAL, bold=True, line_spacing=21)
        _line(slide, Emu(x + Inches(0.2)), Emu(start_y + Inches(1.5)),
              Emu(int(card_w * 0.4)), color, 2)
        if detail:
            _tb(slide, Emu(x + Inches(0.2)), Emu(start_y + Inches(1.75)),
                Emu(card_w - Inches(0.4)), Inches(2.2),
                detail, BODY_FONT, 12, MID, line_spacing=17)
        if i < n - 1:
            arrow_x = Emu(x + card_w + (gap - Inches(0.3)) // 2)
            _tb(slide, arrow_x, Emu(start_y + Inches(0.32)), Inches(0.3), Inches(0.3),
                "\u2192", TITLE_FONT, 16, GOLD, bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════
# 13. PROCESS FLOW VERTICAL — film strip (default for 5+)
# ══════════════════════════════════════════════════════════
def build_process_flow_vertical(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); steps = c.get("steps", [])
    n = len(steps)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
        title, TITLE_FONT, 22, CHARCOAL, bold=True)
    if n == 0: return
    start_y = Inches(0.95); avail = Inches(4.35); gap = Inches(0.06)
    band_h = Emu(int((avail - (n - 1) * gap) / n)); band_w = Inches(9.4)
    for i, step in enumerate(steps):
        color = ACCENTS[i % len(ACCENTS)]
        label = step.get("label", step.get("title", ""))
        detail = step.get("detail", "")
        y = Emu(start_y + i * (band_h + gap))
        bg = WARM_CARD if i % 2 == 0 else LIGHT_BOX
        _rect(slide, Inches(0.3), y, band_w, band_h, bg)
        _rect(slide, Inches(0.3), y, Inches(0.07), band_h, color)
        _tb(slide, Inches(0.5), y, Inches(1.2), band_h,
            str(i + 1), TITLE_FONT, 48, FAINT, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Inches(1.5), y, Inches(2.8), band_h,
            label, TITLE_FONT, 14, CHARCOAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if detail:
            _tb(slide, Inches(4.5), y, Inches(5.0), band_h,
                detail, BODY_FONT, 11, MID, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)


# ══════════════════════════════════════════════════════════
# 13b. PROCESS FLOW ACCORDION (alt vertical — left bars on all cards)
# ══════════════════════════════════════════════════════════
def build_process_flow_accordion(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); steps = c.get("steps", [])
    n = len(steps)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
        title, TITLE_FONT, 22, CHARCOAL, bold=True)
    if n == 0: return
    spine_x = Inches(5.0); start_y = Inches(0.95); avail = Inches(4.3)
    row_h = Emu(int(avail / n)); card_w = Inches(3.8)
    spine_end = Emu(start_y + n * row_h)
    _rect(slide, Emu(spine_x - Inches(0.01)), start_y, Inches(0.025), Emu(spine_end - start_y), FAINT)
    for i, step in enumerate(steps):
        color = ACCENTS[i % len(ACCENTS)]
        label = step.get("label", step.get("title", ""))
        detail = step.get("detail", "")
        left_side = (i % 2 == 0)
        y = Emu(start_y + i * row_h); card_h = Emu(row_h - Inches(0.08))
        dot_size = Inches(0.3)
        _oval(slide, Emu(spine_x - dot_size // 2), Emu(y + (card_h - dot_size) // 2),
              dot_size, dot_size, color)
        _tb(slide, Emu(spine_x - dot_size // 2), Emu(y + (card_h - dot_size) // 2),
            dot_size, dot_size, str(i + 1), BODY_FONT, 10, WHITE, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if left_side:
            cx = Emu(spine_x - Inches(0.35) - card_w); bg = WARM_CARD
        else:
            cx = Emu(spine_x + Inches(0.35)); bg = COOL_CARD
        _rect(slide, cx, y, card_w, card_h, bg)
        _rect(slide, cx, y, Inches(0.05), card_h, color)  # always left edge
        _tb(slide, Emu(cx + Inches(0.15)), y, Emu(card_w - Inches(0.3)), Inches(0.35),
            label, TITLE_FONT, 12, CHARCOAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if detail:
            _tb(slide, Emu(cx + Inches(0.15)), Emu(y + Inches(0.35)),
                Emu(card_w - Inches(0.3)), Emu(card_h - Inches(0.4)),
                detail, BODY_FONT, 9, MID, line_spacing=12)


# ══════════════════════════════════════════════════════════
# 14. TIMELINE (route map)
# ══════════════════════════════════════════════════════════
def build_timeline(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); milestones = c.get("milestones", [])
    n = len(milestones)
    _running_head(slide, title)
    if n == 0: return
    STATUS_CLR = {"complete": DK_GREEN, "current": GOLD, "upcoming": QUIET}
    line_y = Inches(2.8); line_x = Inches(0.6); line_w = Inches(8.8)
    _rect(slide, line_x, line_y, line_w, Inches(0.035), DK_GREEN)
    step = int(line_w / max(n - 1, 1)) if n > 1 else 0
    for i, m in enumerate(milestones):
        status = m.get("status", "upcoming"); color = STATUS_CLR.get(status, QUIET)
        date = m.get("date", ""); mtitle = m.get("title", m.get("label", ""))
        detail = m.get("detail", "")
        cx = Emu(line_x + i * step) if n > 1 else Emu(line_x + int(line_w / 2))
        dot_size = Inches(0.22)
        _oval(slide, Emu(cx - dot_size // 2), Emu(line_y - dot_size // 2 + Inches(0.017)),
              dot_size, dot_size, color)
        above = (i % 2 == 0)
        card_w = Inches(1.25) if n <= 5 else Inches(1.1)
        card_x = Emu(cx - int(card_w) // 2)
        if above:
            if date:
                _tb(slide, card_x, Inches(0.55), card_w, Inches(0.3),
                    date, BODY_FONT, 9, color, bold=True, align=PP_ALIGN.CENTER)
            _tb(slide, card_x, Inches(0.85), card_w, Inches(0.7),
                mtitle, TITLE_FONT, 11, CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM, line_spacing=14)
            _rect(slide, Emu(cx - Inches(0.01)), Inches(1.6), Inches(0.02), Inches(1.18), color)
            if detail:
                _tb(slide, card_x, Inches(1.6), card_w, Inches(0.5),
                    detail, BODY_FONT, 7, QUIET, align=PP_ALIGN.CENTER, line_spacing=9)
        else:
            _rect(slide, Emu(cx - Inches(0.01)), Inches(2.86), Inches(0.02), Inches(0.28), color)
            _tb(slide, card_x, Inches(3.15), card_w, Inches(0.7),
                mtitle, TITLE_FONT, 11, CHARCOAL, bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP, line_spacing=14)
            if date:
                _tb(slide, card_x, Inches(3.85), card_w, Inches(0.3),
                    date, BODY_FONT, 9, color, bold=True, align=PP_ALIGN.CENTER)
            if detail:
                _tb(slide, card_x, Inches(4.15), card_w, Inches(0.5),
                    detail, BODY_FONT, 7, QUIET, align=PP_ALIGN.CENTER, line_spacing=9)


# Update BUILDERS dict with batch 2
BUILDERS.update({
    "wsn_reveal": build_wsn_reveal,
    "findings_recs": build_findings_recs,
    "process_flow": build_process_flow,
    "process_flow_vertical": build_process_flow_vertical,
    "process_flow_accordion": build_process_flow_accordion,
    "timeline": build_timeline,
})


# ══════════════════════════════════════════════════════════
# BATCH 3 BUILDERS
# ══════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════
# 16. PERSONA — Magazine profile (single or duo)
# ══════════════════════════════════════════════════════════
def build_persona(prs, c):
    personas = c.get("personas", None)
    if personas and len(personas) >= 2:
        _build_persona_duo(prs, c)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, LM, Inches(0.2), Inches(7), Inches(0.4),
        c.get("title", "Persona").upper(), BODY_FONT, 18, QUIET, bold=True)

    name = c.get("name", "")
    archetype = c.get("archetype", "")
    traits = c.get("traits", [])
    detail = c.get("detail", "")
    strategy = c.get("strategy", "")

    panel_w = Inches(3.2); panel_y = Inches(0.7); panel_h = Inches(4.6)
    _rect(slide, Inches(0.3), panel_y, panel_w, panel_h, DK_GREEN)
    _rect(slide, Inches(0.3), panel_y, panel_w, Inches(0.05), GOLD)
    img_size = Inches(1.3)
    _rect(slide, Inches(0.5), Emu(panel_y + Inches(0.2)), img_size, img_size, RGBColor(0x2A, 0x5A, 0x2A))
    _tb(slide, Inches(0.5), Emu(panel_y + Inches(0.2)), img_size, img_size,
        "IMG", BODY_FONT, 11, RGBColor(0x3A, 0x6A, 0x3A), bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(0.5), Emu(panel_y + Inches(1.65)), Emu(panel_w - Inches(0.4)), Inches(0.5),
        name, TITLE_FONT, 20, WHITE, bold=True, line_spacing=25)
    if archetype:
        _tb(slide, Inches(0.5), Emu(panel_y + Inches(2.15)), Emu(panel_w - Inches(0.4)), Inches(0.35),
            archetype, BODY_FONT, 12, GOLD, italic=True)
    if traits:
        _line(slide, Inches(0.5), Emu(panel_y + Inches(2.65)), Inches(1.0), GOLD, 1)
        for i, trait in enumerate(traits):
            ty = Emu(panel_y + Inches(2.85) + i * Inches(0.32))
            _tb(slide, Inches(0.5), ty, Emu(panel_w - Inches(0.4)), Inches(0.28),
                f"\u2022  {trait}", BODY_FONT, 10, CREAM, line_spacing=13)
    rx = Inches(3.8); rw = Inches(5.9)
    if detail:
        detail_h = Inches(2.2)
        _rect(slide, rx, Inches(0.75), rw, detail_h, LIGHT_BOX)
        _rect(slide, rx, Inches(0.75), Inches(0.06), detail_h, DK_GREEN)
        _tb(slide, Emu(rx + Inches(0.25)), Inches(0.9), Emu(rw - Inches(0.5)), Emu(detail_h - Inches(0.3)),
            detail, BODY_FONT, 13, CHARCOAL, line_spacing=19)
    if strategy:
        strat_y = Inches(3.6)
        _rect(slide, rx, strat_y, rw, Inches(1.4), WARM_CARD)
        _rect(slide, rx, strat_y, Inches(0.06), Inches(1.4), GOLD)
        _tb(slide, Emu(rx + Inches(0.2)), strat_y, rw, Inches(0.3),
            "STRATEGY", BODY_FONT, 8, GOLD, bold=True)
        _tb(slide, Emu(rx + Inches(0.2)), Emu(strat_y + Inches(0.35)),
            Emu(rw - Inches(0.4)), Inches(0.9),
            strategy, BODY_FONT, 12, CHARCOAL, line_spacing=17)


def _build_persona_duo(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, LM, Inches(0.2), Inches(7), Inches(0.4),
        c.get("title", "Personas").upper(), BODY_FONT, 18, QUIET, bold=True)
    personas = c.get("personas", [])[:2]
    colors = [DK_GREEN, COBALT]
    accents = [GOLD, RGBColor(0x7A, 0xB8, 0xD4)]
    for idx, p in enumerate(personas):
        base_x = Inches(0.3) + idx * Inches(4.85)
        col_w = Inches(4.55); color = colors[idx]; accent = accents[idx]
        _rect(slide, base_x, Inches(0.7), col_w, Inches(1.5), color)
        _rect(slide, base_x, Inches(0.7), col_w, Inches(0.04), accent)
        name = p.get("name", "")
        img_size = Inches(0.9)
        darker = RGBColor(0x2A, 0x5A, 0x2A) if idx == 0 else RGBColor(0x03, 0x3A, 0x5A)
        _rect(slide, Emu(base_x + Inches(0.15)), Inches(0.8), img_size, img_size, darker)
        _tb(slide, Emu(base_x + Inches(0.15)), Inches(0.8), img_size, img_size,
            "IMG", BODY_FONT, 9, RGBColor(0x3A, 0x6A, 0x3A) if idx == 0 else RGBColor(0x2A, 0x5A, 0x7A),
            bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(base_x + Inches(0.15)), Inches(1.4), Emu(col_w - Inches(0.3)), Inches(0.35),
            name, TITLE_FONT, 16, WHITE, bold=True)
        archetype = p.get("archetype", "")
        if archetype:
            _tb(slide, Emu(base_x + Inches(0.15)), Inches(1.75), Emu(col_w - Inches(0.3)), Inches(0.3),
                archetype, BODY_FONT, 10, accent, italic=True)
        traits = p.get("traits", [])
        trait_y = Inches(2.4)
        for i, trait in enumerate(traits):
            ty = Emu(trait_y + i * Inches(0.28))
            _rect(slide, Emu(base_x + Inches(0.15)), Emu(ty + Inches(0.07)),
                  Inches(0.06), Inches(0.06), color)
            _tb(slide, Emu(base_x + Inches(0.35)), ty, Emu(col_w - Inches(0.5)), Inches(0.25),
                trait, BODY_FONT, 10, CHARCOAL, line_spacing=13)
        detail = p.get("detail", "")
        if detail:
            detail_y = Emu(trait_y + len(traits) * Inches(0.28) + Inches(0.2))
            _line(slide, Emu(base_x + Inches(0.15)), detail_y, Inches(1.0), RULE_CLR, 1)
            _tb(slide, Emu(base_x + Inches(0.15)), Emu(detail_y + Inches(0.15)),
                Emu(col_w - Inches(0.3)), Inches(1.5),
                detail, BODY_FONT, 10, MID, line_spacing=14)
    _rect(slide, Inches(4.97), Inches(0.7), Inches(0.015), Inches(4.5), RULE_CLR)


# ══════════════════════════════════════════════════════════
# 17. TEXT NARRATIVE — Magazine article (lede + two-column body)
# ══════════════════════════════════════════════════════════
def build_text_narrative(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); lede = c.get("lede", ""); body = c.get("body", "")
    _running_head(slide, title)
    _tb(slide, Inches(0.5), Inches(0.7), Inches(9.0), Inches(1.2),
        lede, TITLE_FONT, 22, BLACK, bold=True, line_spacing=30)
    _line(slide, Inches(0.5), Inches(2.0), Inches(2.5), GOLD, 3)
    if isinstance(body, list): body_text = "\n\n".join(body)
    else: body_text = body
    mid = len(body_text) // 2
    split_at = body_text.find('. ', mid)
    if split_at == -1 or split_at > mid + 100: split_at = mid
    else: split_at += 2
    col1 = body_text[:split_at].strip(); col2 = body_text[split_at:].strip()
    col_w = Inches(4.2); col_gap = Inches(0.3); body_y = Inches(2.25); col_h = Inches(3.0)
    _rect(slide, Inches(0.5), body_y, col_w, col_h, WARM_CARD)
    _rect(slide, Inches(0.5), body_y, col_w, Inches(0.05), DK_GREEN)
    _tb(slide, Inches(0.7), Emu(body_y + Inches(0.2)), Emu(col_w - Inches(0.4)), Emu(col_h - Inches(0.3)),
        col1, BODY_FONT, 12, CHARCOAL, line_spacing=17)
    col2_x = Emu(Inches(0.5) + col_w + col_gap)
    _rect(slide, col2_x, body_y, col_w, col_h, COOL_CARD)
    _rect(slide, col2_x, body_y, col_w, Inches(0.05), COBALT)
    _tb(slide, Emu(col2_x + Inches(0.2)), Emu(body_y + Inches(0.2)), Emu(col_w - Inches(0.4)), Emu(col_h - Inches(0.3)),
        col2, BODY_FONT, 12, CHARCOAL, line_spacing=17)


# ══════════════════════════════════════════════════════════
# 18. OPEN QUESTIONS — Provocation cards (2x2)
# ══════════════════════════════════════════════════════════
def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); questions = c.get("questions", [])
    n = min(len(questions), 4)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.55),
        title, TITLE_FONT, 22, CHARCOAL, bold=True)
    if n == 0: return
    cols = 2; rows = (n + 1) // 2
    gap_x = Inches(0.25); gap_y = Inches(0.2)
    total_w = Inches(9.0); total_h = Inches(3.7)
    card_w = Emu((total_w - (cols - 1) * gap_x) / cols)
    card_h = Emu((total_h - (rows - 1) * gap_y) / rows)
    start_x = Inches(0.5); start_y = Inches(1.0)
    for i, q in enumerate(questions[:n]):
        if isinstance(q, dict): q_text = q.get("text", q.get("question", str(q)))
        else: q_text = str(q)
        col = i % cols; row = i // cols
        color = ACCENTS[i % len(ACCENTS)]
        x = Emu(start_x + col * (card_w + gap_x)); y = Emu(start_y + row * (card_h + gap_y))
        _rect(slide, x, y, card_w, card_h, LIGHT_BOX)
        _rect(slide, x, y, card_w, Inches(0.05), color)
        _tb(slide, Emu(x + card_w - Inches(1.0)), y, Inches(1.0), card_h,
            "?", TITLE_FONT, 80, FAINT, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.BOTTOM)
        _tb(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.15)), Inches(0.4), Inches(0.4),
            str(i + 1), TITLE_FONT, 22, color, bold=True)
        _tb(slide, Emu(x + Inches(0.2)), Emu(y + Inches(0.6)),
            Emu(card_w - Inches(0.5)), Emu(card_h - Inches(0.8)),
            q_text, BODY_FONT, 13, CHARCOAL, line_spacing=18)


# ══════════════════════════════════════════════════════════
# 19. DATA TABLE
# ══════════════════════════════════════════════════════════
def build_data_table(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", ""); headers = c.get("headers", [])
    rows = c.get("rows", []); note = c.get("note", "")
    highlight_col = c.get("highlightCol", -1)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.55),
        title, TITLE_FONT, 22, CHARCOAL, bold=True)
    if not headers: return
    n_cols = len(headers); n_rows = len(rows)
    table_x = Inches(0.5); table_w = Inches(9.0)
    col_w = Emu(table_w / n_cols); header_h = Inches(0.5)
    row_h = Inches(0.5); start_y = Inches(1.0)
    for j, h in enumerate(headers):
        hx = Emu(table_x + j * col_w)
        _rect(slide, hx, start_y, col_w, header_h, DK_GREEN)
        _tb(slide, Emu(hx + Inches(0.12)), start_y, Emu(col_w - Inches(0.24)), header_h,
            h, BODY_FONT, 11, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, row in enumerate(rows):
        ry = Emu(start_y + header_h + i * row_h)
        for j, cell in enumerate(row):
            cx = Emu(table_x + j * col_w)
            if j == highlight_col: bg = WARM_CARD
            elif i % 2 == 0: bg = LIGHT_BOX
            else: bg = WHITE
            _rect(slide, cx, ry, col_w, row_h, bg)
            is_hl = (j == highlight_col)
            _tb(slide, Emu(cx + Inches(0.12)), ry, Emu(col_w - Inches(0.24)), row_h,
                str(cell), BODY_FONT, 11, DK_GREEN if is_hl else CHARCOAL,
                bold=is_hl, valign=MSO_ANCHOR.MIDDLE)
    bottom_y = Emu(start_y + header_h + n_rows * row_h)
    _line(slide, table_x, bottom_y, table_w, CHARCOAL, 1)
    if note:
        _tb(slide, table_x, Emu(bottom_y + Inches(0.15)), Inches(8), Inches(0.3),
            note, BODY_FONT, 8, QUIET, italic=True)


# ══════════════════════════════════════════════════════════
# 20. SUMMARY — 2-slide reveal (sections → takeaways)
# ══════════════════════════════════════════════════════════
def build_summary(prs, c):
    title = c.get("title", ""); heading = c.get("heading", "")
    sections = c.get("sections", []); points = c.get("points", [])
    n_sec = len(sections)
    # Slide 1: Sections dominant
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide1, WHITE)
    _rect(slide1, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide1, Inches(0.5), Inches(0.25), Inches(8), Inches(0.55),
        title, TITLE_FONT, 22, BLACK, bold=True)
    if heading:
        _tb(slide1, Inches(0.5), Inches(0.8), Inches(8), Inches(0.3),
            heading, BODY_FONT, 12, MID, italic=True)
    if n_sec > 0:
        sec_y = Inches(1.25); gap = Inches(0.2)
        sec_w = Emu((Inches(9.0) - (n_sec - 1) * gap) / n_sec); sec_h = Inches(2.5)
        for i, sec in enumerate(sections):
            x = Emu(Inches(0.5) + i * (sec_w + gap)); color = ACCENTS[i % len(ACCENTS)]
            _rect(slide1, x, sec_y, sec_w, sec_h, LIGHT_BOX)
            _rect(slide1, x, sec_y, sec_w, Inches(0.05), color)
            _tb(slide1, Emu(x + Inches(0.15)), Emu(sec_y + Inches(0.15)),
                Emu(sec_w - Inches(0.3)), Inches(0.4),
                sec.get("heading", ""), TITLE_FONT, 14, color, bold=True, line_spacing=17)
            for j, pt in enumerate(sec.get("points", [])):
                py = Emu(sec_y + Inches(0.65) + j * Inches(0.5))
                _rect(slide1, Emu(x + Inches(0.15)), Emu(py + Inches(0.05)),
                      Inches(0.06), Inches(0.06), color)
                _tb(slide1, Emu(x + Inches(0.3)), py,
                    Emu(sec_w - Inches(0.45)), Inches(0.45),
                    pt, BODY_FONT, 11, CHARCOAL, line_spacing=14)
    if points:
        tease_y = Inches(4.05)
        _rect(slide1, Inches(0.3), tease_y, Inches(9.4), Inches(1.2), WARM_CARD)
        _rect(slide1, Inches(0.3), tease_y, Inches(0.04), Inches(1.2), FAINT)
        _tb(slide1, Inches(0.5), Emu(tease_y + Inches(0.08)), Inches(3), Inches(0.2),
            "KEY TAKEAWAYS", BODY_FONT, 8, QUIET, bold=True)
        for i, pt in enumerate(points):
            short = pt if len(pt) <= 70 else pt[:67] + "..."
            py = Emu(tease_y + Inches(0.32) + i * Inches(0.2))
            _tb(slide1, Inches(0.5), py, Inches(8.5), Inches(0.2),
                f"{i+1}.  {short}", BODY_FONT, 9, QUIET, line_spacing=11)
    # Slide 2: Takeaways dominant
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide2, WHITE)
    _rect(slide2, Inches(0), Inches(0), W, Inches(0.04), DK_GREEN)
    _tb(slide2, Inches(0.5), Inches(0.25), Inches(8), Inches(0.55),
        title, TITLE_FONT, 22, BLACK, bold=True)
    if n_sec > 0:
        comp_y = Inches(0.95); comp_h = Inches(0.9)
        _rect(slide2, Inches(0.3), comp_y, Inches(9.4), comp_h, LIGHT_BOX)
        comp_gap = Inches(0.15)
        comp_item_w = Emu((Inches(9.0) - (n_sec - 1) * comp_gap) / n_sec)
        for i, sec in enumerate(sections):
            cx = Emu(Inches(0.5) + i * (comp_item_w + comp_gap)); color = ACCENTS[i % len(ACCENTS)]
            _rect(slide2, cx, comp_y, Inches(0.04), comp_h, color)
            _tb(slide2, Emu(cx + Inches(0.15)), comp_y, Emu(comp_item_w - Inches(0.2)), comp_h,
                sec.get("heading", ""), BODY_FONT, 10, QUIET, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if points:
        take_y = Inches(2.1)
        _rect(slide2, Inches(0.3), take_y, Inches(9.4), Inches(0.04), DK_GREEN)
        _tb(slide2, Inches(0.5), Emu(take_y + Inches(0.12)), Inches(3), Inches(0.25),
            "KEY TAKEAWAYS", BODY_FONT, 9, DK_GREEN, bold=True)
        n_pts = len(points); cols = 2
        pt_gap_x = Inches(0.2); pt_gap_y = Inches(0.12)
        pt_w = Emu((Inches(9.0) - (cols - 1) * pt_gap_x) / cols)
        pt_h = Inches(0.7); pt_start_y = Inches(2.5)
        for i, pt in enumerate(points):
            col = i % cols; row = i // cols; color = ACCENTS[i % len(ACCENTS)]
            px = Emu(Inches(0.5) + col * (pt_w + pt_gap_x))
            py = Emu(pt_start_y + row * (pt_h + pt_gap_y))
            _rect(slide2, px, py, pt_w, pt_h, WARM_CARD)
            _rect(slide2, px, py, Inches(0.05), pt_h, color)
            _tb(slide2, Emu(px + Inches(0.2)), py, Emu(pt_w - Inches(0.35)), pt_h,
                pt, BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


# ══════════════════════════════════════════════════════════
# 21. APPENDIX — Intentionally muted endnotes
# ══════════════════════════════════════════════════════════
def build_appendix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    title = c.get("title", "Appendix"); label = c.get("label", "")
    sections = c.get("sections", [])
    _tb(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.4),
        title, TITLE_FONT, 18, QUIET, bold=True)
    if label:
        _tb(slide, Inches(0.5), Inches(0.6), Inches(3), Inches(0.25),
            label, BODY_FONT, 9, FAINT, bold=True)
    _line(slide, Inches(0.5), Inches(0.9), Inches(2.0), RULE_CLR, 1)
    start_y = Inches(1.1); content_x = Inches(2.5); content_w = Inches(7.0)
    for i, sec in enumerate(sections):
        sec_label = sec.get("label", sec.get("heading", ""))
        sec_content = sec.get("content", "")
        y = Emu(start_y + i * Inches(1.2))
        _tb(slide, Inches(0.5), y, Inches(1.8), Inches(0.3),
            sec_label, BODY_FONT, 10, QUIET, bold=True)
        _line(slide, Inches(0.5), Emu(y + Inches(0.3)), Inches(1.5), RULE_CLR, 1)
        _tb(slide, content_x, y, content_w, Inches(1.0),
            sec_content, BODY_FONT, 10, MID, line_spacing=14)


# Update BUILDERS with batch 3
BUILDERS.update({
    "persona": build_persona,
    "persona_duo": build_persona,  # same builder handles both
    "text_narrative": build_text_narrative,
    "open_questions": build_open_questions,
    "data_table": build_data_table,
    "summary": build_summary,
    "appendix": build_appendix,
})


def build_deck(slide_configs, output_path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    for slide_type, data in slide_configs:
        if slide_type == "skip": continue
        builder = BUILDERS.get(slide_type)
        if builder: builder(prs, data)
    prs.save(output_path); return output_path
