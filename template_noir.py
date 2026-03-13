"""
Noir template — full dark mode, vivid accents, commanding presence.
Signature: near-black backgrounds everywhere, electric accent colors,
oversized typography, bold color blocks, no subtlety.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

W = Inches(10)
H = Inches(5.625)

# ── Palette ──────────────────────────────────────────────
BLACK      = RGBColor(0x0D, 0x0D, 0x0D)   # primary bg
NEAR_BLACK = RGBColor(0x14, 0x14, 0x14)   # card/elevated bg
DK_GREEN   = RGBColor(0x04, 0x40, 0x14)   # guardrail: title/section/closer
MINT_GLOW  = RGBColor(0x00, 0xFF, 0xC8)   # bright mint for dark green slides

# Vivid accent palette — designed to pop on black
CORAL      = RGBColor(0xFF, 0x5C, 0x5C)   # electric coral
CYAN       = RGBColor(0x00, 0xD4, 0xFF)   # bright cyan
VIOLET     = RGBColor(0xA855, 0xF7, 0x00)[0:3] if False else RGBColor(0xA8, 0x55, 0xF7)  # vivid purple
AMBER      = RGBColor(0xFF, 0xBB, 0x33)   # hot gold
LIME       = RGBColor(0x4A, 0xDE, 0x80)   # vivid green

ACCENTS = [CORAL, CYAN, VIOLET, AMBER, LIME]

# Text colors on dark
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BRIGHT     = RGBColor(0xF0, 0xF0, 0xF0)   # primary text
DIM        = RGBColor(0x99, 0x99, 0x99)    # secondary text
FAINT      = RGBColor(0x55, 0x55, 0x55)    # captions, rules
DARK_RULE  = RGBColor(0x2A, 0x2A, 0x2A)   # subtle dividers on black

SECTION_COLORS = {
    "green": DK_GREEN, "blue": RGBColor(0x38, 0x80, 0xF3),
    "purple": RGBColor(0x5B, 0x2C, 0x8F), "cobalt": RGBColor(0x04, 0x54, 0x7C),
    "gold": RGBColor(0xD4, 0xA8, 0x43),
    "red": RGBColor(0xC2, 0x3B, 0x22), "teal": RGBColor(0x1B, 0x7A, 0x6E),
    "ochre": RGBColor(0xCC, 0x7A, 0x2E), "slate": RGBColor(0x5A, 0x6A, 0x7A),
    "plum": RGBColor(0x8E, 0x45, 0x85),
}

def _resolve_section_color(c, default=None):
    sc = c.get("sectionColor", "")
    return SECTION_COLORS.get(sc, default or DK_GREEN)

TITLE_FONT = "Fidelity Slab"
TITLE_FONT_FALLBACK = "Trebuchet MS"
BODY_FONT  = "Fidelity Sans"
BODY_FONT_FALLBACK = "Calibri"

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "logo.png")

LM = Inches(0.7)
CW = Inches(8.6)


# ── Helpers ──────────────────────────────────────────────

def _font_fallback(font_name):
    if font_name == TITLE_FONT: return TITLE_FONT, TITLE_FONT_FALLBACK
    if font_name == BODY_FONT: return BODY_FONT, BODY_FONT_FALLBACK
    return font_name, None

def _set_font_with_fallback(font_obj, font_name):
    primary, fallback = _font_fallback(font_name)
    font_obj.name = primary
    if fallback:
        el = font_obj._element
        latin = el.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(el, qn('a:latin'))
            latin.set('typeface', primary)
        cs = el.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(el, qn('a:cs'))
        cs.set('typeface', fallback)

def _rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

def _oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

DEFAULT_ICON_SHAPES = [MSO_SHAPE.OVAL, MSO_SHAPE.RECTANGLE, MSO_SHAPE.ISOSCELES_TRIANGLE]

def _icon_or_image(slide, img_path, x, y, w, h, shape_type, fill_color):
    """Add a user-supplied image or a coloured placeholder shape."""
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y, w, h)
    else:
        s = slide.shapes.add_shape(shape_type, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill_color
        s.line.fill.background()

def _tb(slide, x, y, w, h, text, font=None, size=12, color=None,
        bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
        line_spacing=None):
    font = font or BODY_FONT; color = color or BRIGHT
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    _set_font_with_fallback(p.font, font); p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.italic = italic
    p.alignment = align
    if line_spacing: p.line_spacing = Pt(line_spacing)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', {MSO_ANCHOR.TOP:'t', MSO_ANCHOR.MIDDLE:'ctr', MSO_ANCHOR.BOTTOM:'b'}[valign])
    return box

def _dark_bg(slide):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = BLACK

def _green_bg(slide):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = DK_GREEN

def _content_header(slide, title, accent_color=CORAL, size=28):
    """Bold title with vivid color block accent."""
    _rect(slide, LM, Inches(0.25), Inches(0.12), Inches(0.7), accent_color)
    _tb(slide, Emu(LM + Inches(0.25)), Inches(0.25), Inches(8.0), Inches(0.7), title,
        TITLE_FONT, size, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

def _rounded_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    pg = s._element.find('.//' + qn('a:prstGeom'))
    if pg is not None:
        av = pg.find(qn('a:avLst'))
        if av is None: av = etree.SubElement(pg, qn('a:avLst'))
        for g in av.findall(qn('a:gd')): av.remove(g)
        g = etree.SubElement(av, qn('a:gd')); g.set('name','adj'); g.set('fmla','val 5000')
    return s

CONTENT_TOP = Inches(1.15)


def _draw_stepper_bar(slide, labels, active_count, y):
    """Reusable stepper bar: numbered circles + connecting lines + labels."""
    n = len(labels)
    cs = Inches(0.42)
    total_w = CW
    spacing = int(total_w / max(n - 1, 1)) if n > 1 else 0
    start_x = LM
    for i in range(n):
        cx = Emu(start_x + i * spacing) if n > 1 else Emu(start_x + total_w // 2 - cs // 2)
        active = i < active_count
        color = ACCENTS[i % len(ACCENTS)] if active else FAINT
        _oval(slide, cx, y, cs, cs, color)
        tc = WHITE if active else DIM
        _tb(slide, cx, y, cs, cs, str(i + 1), TITLE_FONT, 15, tc,
            bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        lw = Inches(1.4)
        label_x = Emu(cx - (lw - cs) // 2)
        _tb(slide, label_x, Emu(y + cs + Inches(0.04)),
            lw, Inches(0.25), labels[i],
            BODY_FONT, 9, color if active else DIM,
            bold=active, align=PP_ALIGN.CENTER)
        if i < n - 1:
            line_x1 = Emu(cx + cs + Inches(0.08))
            next_cx = Emu(start_x + (i + 1) * spacing)
            line_x2 = Emu(next_cx - Inches(0.08))
            _rect(slide, line_x1, Emu(y + cs // 2), Emu(line_x2 - line_x1), Inches(0.025), FAINT)


# ── Builders ─────────────────────────────────────────────

def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _green_bg(slide)
    # Bold bottom accent bar
    _rect(slide, Inches(0), Emu(H - Inches(0.12)), W, Inches(0.12), MINT_GLOW)
    # Title
    _tb(slide, Inches(0.7), Inches(0.8), Inches(8.6), Inches(2.2),
        c.get("title", "Title"), TITLE_FONT, 42, WHITE, bold=True,
        valign=MSO_ANCHOR.BOTTOM)
    # Vivid accent bar
    _rect(slide, Inches(0.7), Inches(3.15), Inches(3.5), Inches(0.1), MINT_GLOW)
    if c.get("subtitle"):
        _tb(slide, Inches(0.7), Inches(3.45), Inches(8.6), Inches(0.5),
            c["subtitle"], BODY_FONT, 17, RGBColor(0xA0, 0xD8, 0xB0))
    meta = []
    if c.get("author"): meta.append(c["author"])
    if c.get("date"): meta.append(c["date"])
    if meta:
        _tb(slide, Inches(0.7), Inches(4.4), Inches(5), Inches(0.7),
            "\n".join(meta), BODY_FONT, 13, RGBColor(0x80, 0xB0, 0x80))
    # Logo in upper-right
    logo = c.get("logo_path", LOGO_PATH)
    if logo and os.path.isfile(logo):
        slide.shapes.add_picture(logo, Inches(7.6), Inches(0.3), Inches(2.0))
    img = c.get("imagePath", "")
    if img and os.path.exists(img):
        slide.shapes.add_picture(img, Inches(5.4), Inches(0.3), Inches(4.3), Inches(5.0))


def build_in_brief(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "In Brief"), CORAL)
    bullets = c.get("bullets", [])
    n = len(bullets)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.95))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, b in enumerate(bullets):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Elevated dark card
        _rect(slide, LM, y, CW, rowH, NEAR_BLACK)
        # Bold color left strip
        _rect(slide, LM, y, Inches(0.14), rowH, col)
        # Number in color
        _tb(slide, Emu(LM + Inches(0.3)), y, Inches(0.5), rowH,
            str(i + 1), TITLE_FONT, 24, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        # Text
        _tb(slide, Emu(LM + Inches(0.8)), y, Inches(7.5), rowH,
            b, BODY_FONT, 13, BRIGHT, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=18)


def build_section_divider(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    # Bottom glow bar
    _rect(slide, Inches(0), Emu(H - Inches(0.1)), W, Inches(0.1), MINT_GLOW)
    if c.get("sectionNumber"):
        _tb(slide, Inches(0.7), Inches(0.8), Inches(2), Inches(1.0),
            f"0{c['sectionNumber']}", TITLE_FONT, 60, MINT_GLOW, bold=True,
            valign=MSO_ANCHOR.BOTTOM)
    _tb(slide, Inches(0.7), Inches(2.0), Inches(8.6), Inches(1.2),
        c.get("title", "Section"), TITLE_FONT, 40, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(0.7), Inches(3.3), Inches(3.0), Inches(0.08), MINT_GLOW)
    if c.get("subtitle"):
        _tb(slide, Inches(0.7), Inches(3.6), Inches(8.6), Inches(0.5),
            c["subtitle"], BODY_FONT, 15, RGBColor(0xA0, 0xD8, 0xB0))


def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    sc = _resolve_section_color(c)
    _content_header(slide, c.get("title", "Key Metric"), AMBER, size=22)
    # Visual anchor circle behind stat (dark tint for noir)
    cs = Inches(2.2)
    cx = Emu((W - cs) // 2)
    cy = Emu(Inches(1.5))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, cs, cs)
    oval_fill = oval.fill
    oval_fill.solid()
    r, g, b = sc[0], sc[1], sc[2]
    dark_r = int(r * 0.2)
    dark_g = int(g * 0.2)
    dark_b = int(b * 0.2)
    oval_fill.fore_color.rgb = RGBColor(dark_r, dark_g, dark_b)
    oval.line.fill.background()
    # Massive color block with the stat
    _rect(slide, Inches(0.7), Inches(1.25), Inches(8.6), Inches(2.2), NEAR_BLACK)
    # Centered colored rule
    _rect(slide, Inches(3.75), Inches(1.25), Inches(2.5), Inches(0.1), AMBER)
    _tb(slide, Inches(0.7), Inches(1.35), Inches(8.6), Inches(2.1),
        c.get("stat", "—"), TITLE_FONT, 96, AMBER, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if c.get("headline"):
        _tb(slide, Inches(1.0), Inches(3.6), Inches(8.0), Inches(0.6),
            c["headline"], BODY_FONT, 16, WHITE, bold=True,
            align=PP_ALIGN.CENTER)
    if c.get("detail"):
        _tb(slide, Inches(1.5), Inches(4.2), Inches(7.0), Inches(0.7),
            c["detail"], BODY_FONT, 12, DIM, align=PP_ALIGN.CENTER,
            line_spacing=17)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.3),
            c["source"], BODY_FONT, 9, FAINT, italic=True, align=PP_ALIGN.CENTER)


def build_quote(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "In Their Words"), VIOLET, size=22)
    # Large open quote in accent color
    _tb(slide, Inches(0.5), Inches(1.0), Inches(1.0), Inches(1.2),
        "\u201C", TITLE_FONT, 100, RGBColor(0x30, 0x30, 0x30), bold=True)
    # Quote
    _tb(slide, Inches(1.2), Inches(1.5), Inches(7.5), Inches(2.2),
        c.get("quote", ""), BODY_FONT, 20, WHITE, italic=True,
        valign=MSO_ANCHOR.MIDDLE, line_spacing=28)
    # Color attribution bar
    _rect(slide, Inches(1.2), Inches(3.9), Inches(2.0), Inches(0.06), VIOLET)
    if c.get("attribution"):
        _tb(slide, Inches(1.2), Inches(4.1), Inches(7), Inches(0.35),
            c["attribution"], BODY_FONT, 13, VIOLET, bold=True)
    if c.get("context"):
        _tb(slide, Inches(1.2), Inches(4.45), Inches(7), Inches(0.35),
            c["context"], BODY_FONT, 10, DIM, italic=True)


def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Comparison"), CYAN, size=24)

    cardH = Inches(3.7)
    # Left card
    _rect(slide, LM, CONTENT_TOP, Inches(4.0), cardH, NEAR_BLACK)
    _rect(slide, LM, CONTENT_TOP, Inches(4.0), Inches(0.55), CORAL)
    _tb(slide, Emu(LM + Inches(0.15)), CONTENT_TOP, Inches(3.7), Inches(0.55),
        c.get("leftLabel", "Before"), TITLE_FONT, 14, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(c.get("leftItems", [])):
        _tb(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.7) + i * Inches(0.6)),
            Inches(3.5), Inches(0.55), item, BODY_FONT, 12, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)

    # Right card
    _rect(slide, Inches(4.95), CONTENT_TOP, Inches(4.35), cardH, NEAR_BLACK)
    _rect(slide, Inches(4.95), CONTENT_TOP, Inches(4.35), Inches(0.55), CYAN)
    _tb(slide, Inches(5.1), CONTENT_TOP, Inches(4.0), Inches(0.55),
        c.get("rightLabel", "After"), TITLE_FONT, 14, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(c.get("rightItems", [])):
        _tb(slide, Inches(5.15), Emu(CONTENT_TOP + Inches(0.7) + i * Inches(0.6)),
            Inches(3.95), Inches(0.55), item, BODY_FONT, 12, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)


def build_text_graph(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Title"), LIME, size=24)
    texts = c.get("text", [])
    if not isinstance(texts, list): texts = [texts]
    for i, t in enumerate(texts):
        _tb(slide, LM, Emu(CONTENT_TOP + i * Inches(1.2)),
            Inches(3.8), Inches(1.1), t, BODY_FONT, 12, BRIGHT,
            line_spacing=17)
    _rect(slide, Inches(4.8), CONTENT_TOP, Inches(0.06), Inches(3.5), FAINT)

    from pptx.chart.data import CategoryChartData
    raw = c.get("chartData", [{"name": "S1", "labels": ["A","B","C"], "values": [25,45,30]}])
    cd_obj = CategoryChartData()
    if raw:
        cd = raw[0]
        cd_obj.categories = cd.get("labels", ["A","B","C"])
        cd_obj.add_series(cd.get("name","Series 1"), cd.get("values",[25,45,30]))
    ct = {"line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}.get(
        c.get("chartType","bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(ct, Inches(5.05), CONTENT_TOP, Inches(4.0), Inches(3.6), cd_obj)


def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Process"), CYAN, size=24)
    steps = c.get("steps", [])
    n = min(len(steps), 5)
    if n == 0: return

    labels = [s.get("title", f"Step {i+1}") for i, s in enumerate(steps[:n])]
    _draw_stepper_bar(slide, labels, n, Emu(CONTENT_TOP + Inches(0.05)))

    card_top = Emu(CONTENT_TOP + Inches(0.85))
    gap = Inches(0.2)

    if n <= 3:
        card_w = int((CW - (n - 1) * gap) / n)
        card_h = Inches(3.0)
        for i, step in enumerate(steps[:n]):
            cx = Emu(LM + i * (card_w + gap))
            _rounded_rect(slide, cx, card_top, card_w, card_h, NEAR_BLACK)
            _rect(slide, cx, card_top, Inches(0.06), card_h, ACCENTS[i % len(ACCENTS)])
            _tb(slide, Emu(cx + Inches(0.2)), Emu(card_top + Inches(0.15)),
                Emu(card_w - Inches(0.4)), Inches(0.5),
                step.get("title", ""), BODY_FONT, 13, BRIGHT, bold=True)
            if step.get("detail"):
                _tb(slide, Emu(cx + Inches(0.2)), Emu(card_top + Inches(0.65)),
                    Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.85)),
                    step["detail"], BODY_FONT, 10, DIM, line_spacing=14)
    elif n == 4:
        card_w = int((CW - gap) / 2)
        card_h = Inches(1.4)
        for i, step in enumerate(steps[:4]):
            col = i % 2
            row = i // 2
            cx = Emu(LM + col * (card_w + gap))
            cy = Emu(card_top + row * (card_h + gap))
            _rounded_rect(slide, cx, cy, card_w, card_h, NEAR_BLACK)
            _rect(slide, cx, cy, Inches(0.06), card_h, ACCENTS[i % len(ACCENTS)])
            _tb(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.1)),
                Emu(card_w - Inches(0.4)), Inches(0.4),
                step.get("title", ""), BODY_FONT, 12, BRIGHT, bold=True)
            if step.get("detail"):
                _tb(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.5)),
                    Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.65)),
                    step["detail"], BODY_FONT, 10, DIM, line_spacing=14)
    else:
        card_w = int((CW - 2 * gap) / 3)
        card_h = Inches(1.35)
        for i, step in enumerate(steps[:5]):
            if i < 3:
                cx = Emu(LM + i * (card_w + gap))
                cy = card_top
            else:
                offset = (CW - 2 * card_w - gap) // 2
                cx = Emu(LM + offset + (i - 3) * (card_w + gap))
                cy = Emu(card_top + card_h + gap)
            _rounded_rect(slide, cx, cy, card_w, card_h, NEAR_BLACK)
            _rect(slide, cx, cy, Inches(0.06), card_h, ACCENTS[i % len(ACCENTS)])
            _tb(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.1)),
                Emu(card_w - Inches(0.4)), Inches(0.35),
                step.get("title", ""), BODY_FONT, 12, BRIGHT, bold=True)
            if step.get("detail"):
                _tb(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.45)),
                    Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.55)),
                    step["detail"], BODY_FONT, 10, DIM, line_spacing=14)


def build_process_flow_reveal(prs, c):
    """Generates N slides, one per step, with progressive stepper."""
    steps = c.get("steps", [])
    n = min(len(steps), 5)
    if n == 0: return

    labels = [s.get("title", f"Step {i+1}") for i, s in enumerate(steps[:n])]

    for step_idx in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _content_header(slide, c.get("title", "Process"), CYAN, size=24)

        _draw_stepper_bar(slide, labels, step_idx + 1, Emu(CONTENT_TOP + Inches(0.05)))

        card_y = Emu(CONTENT_TOP + Inches(0.85))
        card_w = Inches(6.0)
        card_h = Inches(3.0)
        card_x = Emu(LM + (CW - card_w) // 2)

        step = steps[step_idx]
        _rounded_rect(slide, card_x, card_y, card_w, card_h, NEAR_BLACK)
        _rect(slide, card_x, card_y, Inches(0.08), card_h, ACCENTS[step_idx % len(ACCENTS)])
        _tb(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.2)),
            Emu(card_w - Inches(0.5)), Inches(0.6),
            step.get("title", ""), BODY_FONT, 16, BRIGHT, bold=True)
        if step.get("detail"):
            _tb(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.85)),
                Emu(card_w - Inches(0.5)), Emu(card_h - Inches(1.1)),
                step["detail"], BODY_FONT, 12, DIM, line_spacing=16)


def build_matrix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Framework"), VIOLET, size=24)
    quads = c.get("quadrants", [{}, {}, {}, {}])
    qW = Inches(4.15)
    qH = Inches(1.75)
    gap = Inches(0.2)
    sX = LM

    for i, q in enumerate(quads[:4]):
        col_idx = i % 2; row = i // 2
        x = Emu(sX + col_idx * (qW + gap))
        y = Emu(CONTENT_TOP + row * (qH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, y, qW, qH, NEAR_BLACK)
        # Bold color top bar
        _rect(slide, x, y, qW, Inches(0.08), col)
        _tb(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.15)),
            Emu(qW - Inches(0.3)), Inches(0.35),
            q.get("label", ""), TITLE_FONT, 13, col, bold=True)
        _tb(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.55)),
            Emu(qW - Inches(0.3)), Emu(qH - Inches(0.65)),
            q.get("detail", ""), BODY_FONT, 11, DIM, line_spacing=15)


def build_methods(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Approach"), LIME, size=24)
    fields = c.get("fields", [])
    n = len(fields)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.06)) / n), Inches(0.8))
    gap = Inches(0.06)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, f in enumerate(fields):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        bg = NEAR_BLACK if i % 2 == 0 else BLACK
        _rect(slide, LM, y, CW, rowH, bg)
        _rect(slide, LM, y, Inches(0.1), rowH, col)
        _tb(slide, Emu(LM + Inches(0.25)), y, Inches(1.8), rowH,
            f.get("label", ""), TITLE_FONT, 11, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(LM + Inches(2.15)), y, Inches(6.1), rowH,
            f.get("value", ""), BODY_FONT, 12, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


def build_hypotheses(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Hypotheses"), CYAN, size=24)
    hyps = c.get("hypotheses", [])
    n = len(hyps)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.12)) / n), Inches(0.88))
    gap = Inches(0.12)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.5))

    for i, h in enumerate(hyps):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, CW, rowH, NEAR_BLACK)
        _rect(slide, LM, y, Inches(0.12), rowH, col)
        # H-number
        cs = Inches(0.44)
        cx = Emu(LM + Inches(0.22)); cy = Emu(y + (rowH - cs) // 2)
        _oval(slide, cx, cy, cs, cs, col)
        _tb(slide, cx, cy, cs, cs, f"H{i+1}", TITLE_FONT, 12, BLACK,
            bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(LM + Inches(0.75)), y, Inches(5.8), rowH,
            h.get("text", ""), BODY_FONT, 13, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)
        if h.get("status"):
            is_good = h["status"] == "Confirmed"
            is_bad = h["status"] == "Rejected"
            sc = LIME if is_good else (CORAL if is_bad else DIM)
            _rect(slide, Inches(7.5), Emu(y + (rowH - Inches(0.3)) // 2),
                  Inches(1.3), Inches(0.3), sc)
            _tb(slide, Inches(7.5), Emu(y + (rowH - Inches(0.3)) // 2),
                Inches(1.3), Inches(0.3), h["status"],
                BODY_FONT, 9, BLACK, bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE)


def build_wsn_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Key Finding"), CORAL, size=24)
    labels = c.get("labels", ["What", "So What", "Now What"])
    cols_data = [
        (labels[0], CORAL, c.get("what", {})),
        (labels[1], CYAN, c.get("soWhat", {})),
        (labels[2], VIOLET, c.get("nowWhat", {})),
    ]
    colW = Inches(2.7)
    gap_w = Inches(0.15)
    cardH = Inches(3.85)

    for i, (label, color, data) in enumerate(cols_data):
        x = Emu(LM + i * (colW + gap_w))
        _rect(slide, x, CONTENT_TOP, colW, cardH, NEAR_BLACK)
        # Bold color header
        _rect(slide, x, CONTENT_TOP, colW, Inches(0.5), color)
        _tb(slide, Emu(x + Inches(0.12)), CONTENT_TOP, Emu(colW - Inches(0.24)), Inches(0.5),
            label, TITLE_FONT, 13, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.6)),
            Emu(colW - Inches(0.24)), Inches(0.7),
            data.get("headline", ""), BODY_FONT, 12, WHITE, bold=True,
            line_spacing=17)
        if data.get("detail"):
            _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(1.35)),
                Emu(colW - Inches(0.24)), Inches(1.9),
                data["detail"], BODY_FONT, 10, DIM, line_spacing=14)


def build_wsn_reveal(prs, c):
    """3 progressive slides: What -> So What -> Now What with running summary."""
    labels = c.get("labels", ["What", "So What", "Now What"])
    WSN_KEYS = ["what", "soWhat", "nowWhat"]
    WSN_COLORS = [CORAL, CYAN, VIOLET]

    for step in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _rect(slide, LM, Inches(0.25), Inches(0.12), Inches(0.7), CORAL)
        _tb(slide, Emu(LM + Inches(0.25)), Inches(0.25), Inches(8.0), Inches(0.7),
            c.get("title", "Key Finding"), TITLE_FONT, 26, WHITE, bold=True,
            valign=MSO_ANCHOR.MIDDLE)

        # Stepper bar at top
        _draw_stepper_bar(slide, labels, step + 1, Emu(CONTENT_TOP + Inches(0.05)))

        # Featured card for current step
        data = c.get(WSN_KEYS[step], {})
        card_y = Emu(CONTENT_TOP + Inches(0.85))
        card_w = Inches(7.0)
        card_h = Inches(2.4)
        card_x = Emu(LM + (CW - card_w) // 2)

        _rect(slide, card_x, card_y, card_w, card_h, NEAR_BLACK)
        _rect(slide, card_x, card_y, Inches(0.10), card_h, WSN_COLORS[step])
        _tb(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.15)),
            Emu(card_w - Inches(0.5)), Inches(0.6),
            data.get("headline", ""), BODY_FONT, 15, BRIGHT, bold=True,
            line_spacing=18)
        if data.get("detail"):
            _tb(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.8)),
                Emu(card_w - Inches(0.5)), Emu(card_h - Inches(1.0)),
                data["detail"], BODY_FONT, 12, DIM, line_spacing=16)

        # Running summary (only for step >= 1)
        if step > 0:
            summary_y = Emu(card_y + card_h + Inches(0.15))
            # Thin divider
            _rect(slide, LM, summary_y, Inches(3.0), Inches(0.02), FAINT)
            # Previous step summaries
            for j in range(step):
                iy = Emu(summary_y + Inches(0.1) + j * Inches(0.28))
                prev_data = c.get(WSN_KEYS[j], {})
                summary_text = prev_data.get("summary", prev_data.get("headline", ""))
                # Small bullet square
                _rect(slide, LM, Emu(iy + Inches(0.05)), Inches(0.08), Inches(0.08), FAINT)
                _tb(slide, Emu(LM + Inches(0.18)), iy,
                    Emu(CW - Inches(0.18)), Inches(0.25),
                    summary_text, BODY_FONT, 9, DIM,
                    valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Findings & Recommendations"), AMBER, size=22)
    items = c.get("items", [])
    n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.25)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.82))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, item in enumerate(items[:5]):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Finding
        _rect(slide, LM, y, Inches(3.7), rowH, NEAR_BLACK)
        _rect(slide, LM, y, Inches(0.1), rowH, col)
        _tb(slide, Emu(LM + Inches(0.2)), y, Inches(3.3), rowH,
            item.get("finding", ""), BODY_FONT, 11, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)
        # Arrow
        _tb(slide, Inches(4.55), y, Inches(0.4), rowH,
            "\u2192", BODY_FONT, 18, col, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE)
        # Rec
        _rect(slide, Inches(5.1), y, Inches(4.2), rowH, NEAR_BLACK)
        _rect(slide, Inches(5.1), y, Inches(0.1), rowH, col)
        _tb(slide, Inches(5.3), y, Inches(3.8), rowH,
            item.get("recommendation", ""), BODY_FONT, 11, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Complete Findings"), AMBER, size=22)
    items = c.get("items", [])
    n = min(len(items), 8)
    avail = H - CONTENT_TOP - Inches(0.2)
    rowH = min(int((avail - (n - 1) * Inches(0.08)) / n), Inches(0.56))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, item in enumerate(items[:8]):
        y = Emu(startY + i * (rowH + gap))
        bg = NEAR_BLACK if i % 2 == 0 else BLACK
        _rect(slide, LM, y, Inches(4.0), rowH, bg)
        _rect(slide, LM, y, Inches(0.06), rowH, CORAL)
        _tb(slide, Emu(LM + Inches(0.15)), y, Inches(3.7), rowH,
            item.get("finding", ""), BODY_FONT, 11, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)
        _rect(slide, Inches(4.85), y, Inches(4.45), rowH, bg)
        _rect(slide, Inches(4.85), y, Inches(0.06), rowH, CYAN)
        _tb(slide, Inches(5.0), y, Inches(4.15), rowH,
            item.get("recommendation", ""), BODY_FONT, 11, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE)


def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Open Questions"), AMBER, size=24)
    questions = c.get("questions", [])
    cardW = Inches(4.15)
    cardH = Inches(1.75)
    gX = Inches(0.2)
    gY = Inches(0.2)

    for i, question in enumerate(questions[:4]):
        col_idx = i % 2; row = i // 2
        x = Emu(LM + col_idx * (cardW + gX))
        y = Emu(CONTENT_TOP + row * (cardH + gY))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, y, cardW, cardH, NEAR_BLACK)
        _rect(slide, x, y, cardW, Inches(0.08), col)
        _tb(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.12)),
            Inches(0.6), Inches(0.55), str(i + 1),
            TITLE_FONT, 28, col, bold=True)
        _tb(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.7)),
            Emu(cardW - Inches(0.3)), Inches(0.9),
            question, BODY_FONT, 12, BRIGHT, line_spacing=17)


def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Agenda"), CYAN, size=26)
    items = c.get("items", [])
    n = len(items)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int(avail / n), Inches(0.72))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, item in enumerate(items):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, CW, rowH, NEAR_BLACK)
        _rect(slide, LM, y, Inches(0.12), rowH, col)
        # Bold number
        _tb(slide, Emu(LM + Inches(0.25)), y, Inches(0.5), rowH,
            str(i + 1), TITLE_FONT, 24, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        title_text = item if isinstance(item, str) else item.get("title", "")
        _tb(slide, Emu(LM + Inches(0.8)), y, Inches(5.5), rowH,
            title_text, BODY_FONT, 15, WHITE, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        if isinstance(item, dict) and item.get("detail"):
            _tb(slide, Inches(7.2), y, Inches(1.8), rowH,
                item["detail"], BODY_FONT, 11, DIM,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_progressive_reveal(prs, c):
    takeaways = c.get("takeaways", [])
    for n in range(min(len(takeaways), 5)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _rect(slide, LM, Inches(0.25), Inches(0.12), Inches(0.7), LIME)
        _tb(slide, Emu(LM + Inches(0.25)), Inches(0.25), Inches(8.0), Inches(0.7),
            c.get("title", "Building the Picture"), TITLE_FONT, 24, WHITE,
            bold=True, valign=MSO_ANCHOR.MIDDLE)

        cur = takeaways[n]
        col = ACCENTS[n % len(ACCENTS)]
        _rect(slide, LM, Inches(1.1), CW, Inches(2.2), NEAR_BLACK)
        _rect(slide, LM, Inches(1.1), Inches(0.1), Inches(2.2), col)
        _tb(slide, Emu(LM + Inches(0.25)), Inches(1.2), Inches(8.1), Inches(0.55),
            cur.get("headline", ""), BODY_FONT, 15, WHITE, bold=True)
        if cur.get("detail"):
            _tb(slide, Emu(LM + Inches(0.25)), Inches(1.8), Inches(8.1), Inches(1.3),
                cur["detail"], BODY_FONT, 11, DIM, line_spacing=15)

        # Thin subtle divider
        _rect(slide, LM, Inches(3.5), Inches(3.0), Inches(0.02), FAINT)
        # Running summary items (no header, all gray)
        for j in range(n + 1):
            ty = Emu(Inches(3.6) + j * Inches(0.28))
            _rect(slide, LM, Emu(ty + Inches(0.04)), Inches(0.08), Inches(0.08), FAINT)
            _tb(slide, Emu(LM + Inches(0.18)), ty, Inches(8.0), Inches(0.25),
                takeaways[j].get("summary", takeaways[j].get("headline", "")),
                BODY_FONT, 9, DIM,
                valign=MSO_ANCHOR.MIDDLE)


def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    _rect(slide, Inches(0), Emu(H - Inches(0.12)), W, Inches(0.12), MINT_GLOW)
    _tb(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.4),
        c.get("title", "Thank You"), TITLE_FONT, 48, WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
    _rect(slide, Inches(3.25), Inches(2.8), Inches(3.5), Inches(0.08), MINT_GLOW)
    if c.get("subtitle"):
        _tb(slide, Inches(0.5), Inches(3.1), Inches(9), Inches(0.5),
            c["subtitle"], BODY_FONT, 17, RGBColor(0xA0, 0xD8, 0xB0),
            align=PP_ALIGN.CENTER)
    if c.get("contact"):
        _tb(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(0.4),
            c["contact"], BODY_FONT, 13, RGBColor(0x80, 0xB0, 0x80),
            align=PP_ALIGN.CENTER)




def build_timeline(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Timeline"), CYAN)
    milestones = c.get("milestones", [])
    n = min(len(milestones), 6)
    if n == 0: return
    lineY = Inches(2.8)
    startX = 0.7; endX = 9.3
    step = (endX - startX) / max(n - 1, 1)
    _rect(slide, Inches(startX), lineY, Inches(endX - startX), Inches(0.04), FAINT)
    STATUS_CLR = {"complete": LIME, "current": AMBER, "upcoming": FAINT}
    for i, m in enumerate(milestones[:n]):
        cx = Inches(startX + i * step)
        col = STATUS_CLR.get(m.get("status", "upcoming"), FAINT)
        ds = Inches(0.22)
        _oval(slide, Emu(cx - ds // 2), Emu(lineY - ds // 2), ds, ds, col)
        _tb(slide, Emu(cx - Inches(0.6)), Emu(lineY - Inches(0.65)),
            Inches(1.2), Inches(0.3), m.get("date", ""),
            BODY_FONT, 10, col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.25)),
            Inches(1.4), Inches(0.4), m.get("title", ""),
            BODY_FONT, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
        if m.get("detail"):
            _tb(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.65)),
                Inches(1.4), Inches(0.7), m["detail"],
                BODY_FONT, 9, DIM, align=PP_ALIGN.CENTER, line_spacing=12)


def build_data_table(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Data"), LIME, size=24)
    headers = c.get("headers", []); rows = c.get("rows", [])
    hc = c.get("highlightCol", None)
    nCols = len(headers); nRows = min(len(rows), 10)
    if nCols == 0: return
    tW = 8.6; tX = 0.7; colW = tW / nCols; rowH = min(0.42, 3.8 / max(nRows + 1, 1)); hdrY = 1.25
    _rect(slide, Inches(tX), Inches(hdrY), Inches(tW), Inches(rowH), CORAL)
    for j, h in enumerate(headers):
        _tb(slide, Inches(tX + j * colW + 0.08), Inches(hdrY),
            Inches(colW - 0.16), Inches(rowH), h,
            BODY_FONT, 10, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, row in enumerate(rows[:nRows]):
        ry = hdrY + (i + 1) * rowH
        bg = NEAR_BLACK if i % 2 == 0 else BLACK
        _rect(slide, Inches(tX), Inches(ry), Inches(tW), Inches(rowH), bg)
        for j, cell in enumerate(row[:nCols]):
            is_hl = hc is not None and j == hc
            _tb(slide, Inches(tX + j * colW + 0.08), Inches(ry),
                Inches(colW - 0.16), Inches(rowH), str(cell),
                BODY_FONT, 10, AMBER if is_hl else BRIGHT,
                bold=is_hl, valign=MSO_ANCHOR.MIDDLE)
    if c.get("note"):
        _tb(slide, Inches(tX), Inches(hdrY + (nRows + 1) * rowH + 0.1),
            Inches(tW), Inches(0.3), c["note"], BODY_FONT, 8, FAINT, italic=True)


def build_multi_stat(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Key Metrics"), AMBER, size=24)
    stats = c.get("stats", []); n = min(len(stats), 4)
    if n == 0: return
    totalW = 8.6; gap = 0.2; statW = (totalW - (n - 1) * gap) / n
    for i, s in enumerate(stats[:n]):
        x = Inches(0.7 + i * (statW + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, CONTENT_TOP, Inches(statW), Inches(3.5), NEAR_BLACK)
        _rect(slide, x, CONTENT_TOP, Inches(statW), Inches(0.08), col)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(0.2)), Inches(statW), Inches(1.4),
            s.get("value", "—"), TITLE_FONT, 52, col, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(1.7)), Inches(statW), Inches(0.45),
            s.get("label", ""), BODY_FONT, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
        if s.get("detail"):
            _tb(slide, x, Emu(CONTENT_TOP + Inches(2.2)), Inches(statW), Inches(0.9),
                s["detail"], BODY_FONT, 10, DIM, align=PP_ALIGN.CENTER, line_spacing=14)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.3),
            c["source"], BODY_FONT, 9, FAINT, italic=True, align=PP_ALIGN.CENTER)


def build_persona(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Persona"), VIOLET, size=22)
    _rect(slide, LM, CONTENT_TOP, Inches(3.8), Inches(3.9), NEAR_BLACK)
    _rect(slide, LM, CONTENT_TOP, Inches(0.1), Inches(3.9), VIOLET)
    _tb(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)), Inches(3.4), Inches(0.5),
        c.get("name", ""), TITLE_FONT, 24, VIOLET, bold=True)
    _tb(slide, Emu(LM + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.6)), Inches(3.4), Inches(0.3),
        c.get("archetype", ""), BODY_FONT, 12, AMBER, bold=True)
    traits = c.get("traits", [])
    for i, t in enumerate(traits[:5]):
        ty = Emu(CONTENT_TOP + Inches(1.1) + i * Inches(0.4))
        _rect(slide, Emu(LM + Inches(0.2)), Emu(ty + Inches(0.06)), Inches(0.1), Inches(0.1), ACCENTS[i % len(ACCENTS)])
        _tb(slide, Emu(LM + Inches(0.42)), Emu(ty), Inches(3.3), Inches(0.35), t, BODY_FONT, 11, BRIGHT)
    _rect(slide, Inches(4.75), CONTENT_TOP, Inches(4.55), Inches(3.9), NEAR_BLACK)
    _rect(slide, Inches(4.75), CONTENT_TOP, Inches(0.1), Inches(3.9), AMBER)
    _tb(slide, Inches(4.95), Emu(CONTENT_TOP + Inches(0.1)), Inches(4.15), Inches(0.25),
        "STRATEGY", BODY_FONT, 9, FAINT, bold=True)
    _tb(slide, Inches(4.95), Emu(CONTENT_TOP + Inches(0.4)), Inches(4.15), Inches(1.5),
        c.get("strategy", ""), BODY_FONT, 12, BRIGHT, line_spacing=17)
    if c.get("detail"):
        _tb(slide, Inches(4.95), Emu(CONTENT_TOP + Inches(2.0)), Inches(4.15), Inches(1.5),
            c["detail"], BODY_FONT, 10, DIM, line_spacing=14)


def build_risk_tradeoff(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Risk & Reward"), AMBER, size=24)
    risks = c.get("risks", []); rewards = c.get("rewards", [])
    SEV = {"high": CORAL, "medium": AMBER, "low": LIME}
    _rect(slide, LM, CONTENT_TOP, Inches(4.0), Inches(0.45), CORAL)
    _tb(slide, Emu(LM + Inches(0.12)), CONTENT_TOP, Inches(3.7), Inches(0.45),
        "RISKS", BODY_FONT, 11, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(risks[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.55) + i * Inches(0.6))
        sc = SEV.get(r.get("severity", "medium"), AMBER)
        _rect(slide, LM, ry, Inches(4.0), Inches(0.5), NEAR_BLACK)
        _rect(slide, LM, ry, Inches(0.08), Inches(0.5), sc)
        _tb(slide, Emu(LM + Inches(0.18)), ry, Inches(3.6), Inches(0.5),
            r.get("label", ""), BODY_FONT, 10, BRIGHT, valign=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(4.95), CONTENT_TOP, Inches(4.35), Inches(0.45), LIME)
    _tb(slide, Inches(5.07), CONTENT_TOP, Inches(4.1), Inches(0.45),
        "REWARDS", BODY_FONT, 11, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(rewards[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.55) + i * Inches(0.6))
        _rect(slide, Inches(4.95), ry, Inches(4.35), Inches(0.5), NEAR_BLACK)
        _rect(slide, Inches(4.95), ry, Inches(0.08), Inches(0.5), LIME)
        _tb(slide, Inches(5.13), ry, Inches(4.0), Inches(0.5),
            r.get("label", ""), BODY_FONT, 10, BRIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_appendix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _tb(slide, LM, Inches(0.25), CW, Inches(0.4),
        c.get("title", "Appendix"), TITLE_FONT, 16, FAINT, bold=True)
    _rect(slide, LM, Inches(0.68), Inches(1.2), Inches(0.04), FAINT)
    sections = c.get("sections", []); y_cursor = 0.85
    for s in sections:
        if y_cursor > 4.8: break
        _tb(slide, LM, Inches(y_cursor), CW, Inches(0.22),
            s.get("label", ""), BODY_FONT, 9, AMBER, bold=True)
        _tb(slide, LM, Inches(y_cursor + 0.22), CW, Inches(0.75),
            s.get("content", ""), BODY_FONT, 8, DIM, line_spacing=11)
        y_cursor += 1.05


def build_before_after(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Transformation"), CYAN, size=24)
    bef = c.get("before", {}); aft = c.get("after", {}); interv = c.get("intervention", "")
    _rect(slide, LM, CONTENT_TOP, Inches(3.5), Inches(3.6), NEAR_BLACK)
    _rect(slide, LM, CONTENT_TOP, Inches(3.5), Inches(0.5), CORAL)
    _tb(slide, Emu(LM + Inches(0.12)), CONTENT_TOP, Inches(3.2), Inches(0.5),
        bef.get("label", "Before"), TITLE_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Emu(LM + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.6)), Inches(3.2), Inches(2.7),
        bef.get("detail", ""), BODY_FONT, 11, BRIGHT, line_spacing=16)
    _tb(slide, Inches(4.35), Emu(CONTENT_TOP + Inches(1.0)), Inches(1.0), Inches(0.5),
        "\u2192", TITLE_FONT, 30, AMBER, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(4.2), Emu(CONTENT_TOP + Inches(1.6)), Inches(1.3), Inches(1.5),
        interv, BODY_FONT, 9, DIM, align=PP_ALIGN.CENTER, line_spacing=13)
    _rect(slide, Inches(5.65), CONTENT_TOP, Inches(3.65), Inches(3.6), NEAR_BLACK)
    _rect(slide, Inches(5.65), CONTENT_TOP, Inches(3.65), Inches(0.5), LIME)
    _tb(slide, Inches(5.77), CONTENT_TOP, Inches(3.4), Inches(0.5),
        aft.get("label", "After"), TITLE_FONT, 12, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(5.77), Emu(CONTENT_TOP + Inches(0.6)), Inches(3.4), Inches(2.7),
        aft.get("detail", ""), BODY_FONT, 11, BRIGHT, line_spacing=16)


def build_summary(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Summary"), LIME, size=26)
    sections = c.get("sections", []); n = len(sections)
    if n == 0: return
    colW = (8.6 - (n - 1) * 0.15) / n
    for i, sec in enumerate(sections[:4]):
        x = Inches(0.7 + i * (colW + 0.15))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, CONTENT_TOP, Inches(colW), Inches(3.9), NEAR_BLACK)
        _rect(slide, x, CONTENT_TOP, Inches(colW), Inches(0.08), col)
        _tb(slide, Emu(x + Inches(0.12)), Emu(CONTENT_TOP + Inches(0.12)), Inches(colW - 0.24), Inches(0.3),
            sec.get("heading", ""), TITLE_FONT, 11, col, bold=True)
        points = sec.get("points", [])
        for j, p in enumerate(points[:5]):
            py = Emu(CONTENT_TOP + Inches(0.5) + j * Inches(0.6))
            _rect(slide, Emu(x + Inches(0.12)), Emu(py + Inches(0.06)), Inches(0.08), Inches(0.08), col)
            _tb(slide, Emu(x + Inches(0.28)), Emu(py), Inches(colW - 0.44), Inches(0.55),
                p, BODY_FONT, 10, BRIGHT, line_spacing=13)




def build_quote_full(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _rect(slide, Inches(3.0), Inches(0.8), Inches(4.0), Inches(0.06), VIOLET)
    _tb(slide, Inches(0.8), Inches(1.1), Inches(8.4), Inches(2.8),
        c.get("quote", ""), BODY_FONT, 26, WHITE, italic=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=36)
    _rect(slide, Inches(3.0), Inches(4.1), Inches(4.0), Inches(0.06), VIOLET)
    if c.get("attribution"):
        _tb(slide, Inches(1.0), Inches(4.35), Inches(8.0), Inches(0.4),
            c["attribution"], BODY_FONT, 14, VIOLET, bold=True, align=PP_ALIGN.CENTER)
    if c.get("context"):
        _tb(slide, Inches(1.0), Inches(4.75), Inches(8.0), Inches(0.35),
            c["context"], BODY_FONT, 10, FAINT, italic=True, align=PP_ALIGN.CENTER)


def build_stat_hero(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Key Metric"), AMBER, size=22)
    hero = c.get("hero", {})
    _rect(slide, LM, Inches(1.15), CW, Inches(1.8), NEAR_BLACK)
    _rect(slide, LM, Inches(1.15), CW, Inches(0.08), AMBER)
    _tb(slide, LM, Inches(1.2), CW, Inches(1.3),
        hero.get("value", "—"), TITLE_FONT, 72, AMBER, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if hero.get("label"):
        _tb(slide, Inches(1.5), Inches(2.5), Inches(7.0), Inches(0.4),
            hero["label"], BODY_FONT, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
    supporting = c.get("supporting", [])
    n = min(len(supporting), 4)
    if n > 0:
        sw = (8.6 - (n - 1) * 0.15) / n
        for i, s in enumerate(supporting[:n]):
            x = Inches(0.7 + i * (sw + 0.15))
            sc = ACCENTS[i % len(ACCENTS)]
            _rect(slide, x, Inches(3.15), Inches(sw), Inches(1.5), NEAR_BLACK)
            _rect(slide, x, Inches(3.15), Inches(sw), Inches(0.06), sc)
            _tb(slide, x, Inches(3.25), Inches(sw), Inches(0.7),
                s.get("value", ""), TITLE_FONT, 28, sc, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _tb(slide, x, Inches(3.95), Inches(sw), Inches(0.35),
                s.get("label", ""), BODY_FONT, 10, DIM, align=PP_ALIGN.CENTER)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.3),
            c["source"], BODY_FONT, 9, FAINT, italic=True, align=PP_ALIGN.CENTER)


def build_in_brief_featured(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "In Brief"), CORAL)
    featured = c.get("featured", "")
    supporting = c.get("supporting", [])
    _rect(slide, LM, CONTENT_TOP, CW, Inches(1.3), NEAR_BLACK)
    _rect(slide, LM, CONTENT_TOP, Inches(0.12), Inches(1.3), AMBER)
    _tb(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Emu(CW - Inches(0.3)), Inches(1.3),
        featured, BODY_FONT, 16, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE, line_spacing=22)
    for i, s in enumerate(supporting[:4]):
        y = Emu(Inches(2.65) + i * Inches(0.6))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, CW, Inches(0.5), NEAR_BLACK)
        _rect(slide, LM, y, Inches(0.1), Inches(0.5), col)
        _tb(slide, Emu(LM + Inches(0.22)), Emu(y), Emu(CW - Inches(0.3)), Inches(0.5),
            s, BODY_FONT, 12, BRIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_in_brief_reveal(prs, c):
    """Spotlight reveal: each slide features one item large while others stay small."""
    items = c.get("items", [])
    n = min(len(items), 6)
    if n == 0:
        return

    for k in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)
        _content_header(slide, c.get("title", "In Brief"), CORAL)

        y_cursor = Emu(CONTENT_TOP + Inches(0.12))
        for i in range(n):
            col = ACCENTS[i % len(ACCENTS)]
            if i == k:
                # Featured card -- large, bold
                h = Inches(1.2)
                _rect(slide, LM, y_cursor, CW, h, NEAR_BLACK)
                _rect(slide, LM, y_cursor, Inches(0.12), h, col)
                _tb(slide, Emu(LM + Inches(0.25)), y_cursor,
                    Emu(CW - Inches(0.3)), h,
                    items[i], BODY_FONT, 15, WHITE, bold=True,
                    valign=MSO_ANCHOR.MIDDLE, line_spacing=21)
            else:
                # Small row
                h = Inches(0.5)
                _rect(slide, LM, y_cursor, CW, h, NEAR_BLACK)
                _rect(slide, LM, y_cursor, Inches(0.06), h, col)
                _tb(slide, Emu(LM + Inches(0.18)), y_cursor,
                    Emu(CW - Inches(0.3)), h,
                    items[i], BODY_FONT, 11, BRIGHT,
                    valign=MSO_ANCHOR.MIDDLE)
            # Extra gap around the featured card
            gap = Inches(0.15) if i == k or (i + 1 < n and i + 1 == k) else Inches(0.08)
            y_cursor = Emu(y_cursor + h + gap)


def build_persona_duo(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Archetype Comparison"), VIOLET, size=22)
    personas = c.get("personas", [{}, {}])
    for idx, p in enumerate(personas[:2]):
        x = LM if idx == 0 else Inches(4.85)
        w = Inches(3.95) if idx == 0 else Inches(4.45)
        col = ACCENTS[idx % len(ACCENTS)]
        _rect(slide, x, CONTENT_TOP, w, Inches(3.9), NEAR_BLACK)
        _rect(slide, x, CONTENT_TOP, Inches(0.1), Inches(3.9), col)
        _tb(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.1)), Emu(w - Inches(0.3)), Inches(0.4),
            p.get("name", ""), TITLE_FONT, 18, col, bold=True)
        _tb(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(0.5)), Emu(w - Inches(0.3)), Inches(0.25),
            p.get("archetype", ""), BODY_FONT, 10, AMBER, bold=True)
        traits = p.get("traits", [])
        for j, t in enumerate(traits[:4]):
            ty = Emu(CONTENT_TOP + Inches(0.9) + j * Inches(0.38))
            _rect(slide, Emu(x + Inches(0.2)), Emu(ty + Inches(0.06)), Inches(0.08), Inches(0.08), col)
            _tb(slide, Emu(x + Inches(0.4)), Emu(ty), Emu(w - Inches(0.55)), Inches(0.33), t, BODY_FONT, 10, BRIGHT)
        if p.get("strategy"):
            _rect(slide, Emu(x + Inches(0.1)), Emu(CONTENT_TOP + Inches(2.5)), Emu(w - Inches(0.2)), Inches(0.04), FAINT)
            _tb(slide, Emu(x + Inches(0.2)), Emu(CONTENT_TOP + Inches(2.65)), Emu(w - Inches(0.3)), Inches(1.1),
                p["strategy"], BODY_FONT, 10, DIM, line_spacing=14)


def build_process_flow_vertical(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Process"), CYAN, size=24)
    steps = c.get("steps", []); count = min(len(steps), 3)
    if count == 0: return
    avail = 3.8; stepH = (avail - (count - 1) * 0.4) / count
    for i, step in enumerate(steps[:count]):
        y = Inches(1.25 + i * (stepH + 0.4))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, Emu(y), CW, Inches(stepH), NEAR_BLACK)
        _rect(slide, LM, Emu(y), Inches(0.1), Inches(stepH), col)
        _tb(slide, Emu(LM + Inches(0.2)), Emu(y + Inches(0.05)), Inches(0.5), Inches(0.5),
            str(i + 1), TITLE_FONT, 26, col, bold=True)
        _tb(slide, Emu(LM + Inches(0.75)), Emu(y + Inches(0.05)), Inches(3.5), Inches(0.4),
            step.get("title", ""), BODY_FONT, 12, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if step.get("detail"):
            _tb(slide, Emu(LM + Inches(0.75)), Emu(y + Inches(0.5)), Inches(7.5), Emu(Inches(stepH) - Inches(0.55)),
                step["detail"], BODY_FONT, 10, DIM, line_spacing=15)
        if i < count - 1:
            _tb(slide, Emu(LM + Inches(0.3)), Emu(y + Inches(stepH + 0.05)), Inches(0.5), Inches(0.25),
                "\u2193", TITLE_FONT, 20, FAINT, align=PP_ALIGN.CENTER)


def build_text_cards(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Key Points"), LIME, size=24)
    items = c.get("items", []); n = min(len(items), 6)
    if n == 0: return
    cols = 2 if n <= 4 else 3
    rows_count = (n + cols - 1) // cols
    cardW = (8.6 - (cols - 1) * 0.15) / cols
    cardH = (3.8 - (rows_count - 1) * 0.15) / rows_count
    for i, item in enumerate(items[:n]):
        col_idx = i % cols; row = i // cols
        x = Inches(0.7 + col_idx * (cardW + 0.15))
        y = Emu(CONTENT_TOP + row * (Inches(cardH) + Inches(0.15)))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, x, y, Inches(cardW), Inches(cardH), NEAR_BLACK)
        _rect(slide, x, y, Inches(cardW), Inches(0.06), col)
        title_text = item if isinstance(item, str) else item.get("title", "")
        detail_text = "" if isinstance(item, str) else item.get("detail", "")
        if detail_text:
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.1)), Inches(cardW - 0.24), Inches(0.35),
                title_text, BODY_FONT, 12, WHITE, bold=True)
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.5)), Inches(cardW - 0.24), Emu(Inches(cardH) - Inches(0.6)),
                detail_text, BODY_FONT, 10, DIM, line_spacing=14)
        else:
            _tb(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.1)), Inches(cardW - 0.24), Emu(Inches(cardH) - Inches(0.2)),
                title_text, BODY_FONT, 12, BRIGHT, valign=MSO_ANCHOR.MIDDLE, line_spacing=17)


def build_text_columns(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Overview"), CYAN, size=24)
    columns = c.get("columns", []); n = min(len(columns), 3)
    if n == 0: return
    colW = (8.6 - (n - 1) * 0.2) / n
    for i, col_data in enumerate(columns[:n]):
        x = Inches(0.7 + i * (colW + 0.2))
        col = ACCENTS[i % len(ACCENTS)]
        heading = "" if isinstance(col_data, str) else col_data.get("heading", "")
        body = col_data if isinstance(col_data, str) else col_data.get("body", "")
        if heading:
            _tb(slide, x, CONTENT_TOP, Inches(colW), Inches(0.3), heading, TITLE_FONT, 11, col, bold=True)
            _rect(slide, x, Emu(CONTENT_TOP + Inches(0.35)), Inches(0.8), Inches(0.04), col)
            _tb(slide, x, Emu(CONTENT_TOP + Inches(0.5)), Inches(colW), Inches(3.3),
                body, BODY_FONT, 11, BRIGHT, line_spacing=16)
        else:
            _tb(slide, x, CONTENT_TOP, Inches(colW), Inches(3.8), body, BODY_FONT, 11, BRIGHT, line_spacing=16)
        if i < n - 1:
            _rect(slide, Inches(0.7 + (i + 1) * colW + i * 0.2 + 0.1), CONTENT_TOP, Inches(0.02), Inches(3.5), FAINT)




def build_text_narrative(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Overview"), AMBER)
    lede = c.get("lede", "")
    body = c.get("body", "")
    if isinstance(body, list): body = "\n\n".join(body)
    _rect(slide, LM, CONTENT_TOP, CW, Inches(1.3), NEAR_BLACK)
    _rect(slide, LM, CONTENT_TOP, Inches(0.12), Inches(1.3), AMBER)
    _tb(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Emu(CW - Inches(0.3)), Inches(1.3),
        lede, BODY_FONT, 16, WHITE, bold=True, line_spacing=22,
        valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, LM, Inches(2.65), CW, Inches(2.5),
        body, BODY_FONT, 12, DIM, line_spacing=18)


def build_text_nested(prs, c):
    """Section cards — top-level items as colored label blocks, children as flowing text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Detail"), CYAN)
    items = c.get("items", [])
    n = min(len(items), 4)
    if n == 0: return
    avail = H - CONTENT_TOP - Inches(0.3)
    gap = Inches(0.12)
    cardH = (avail - (n - 1) * gap) / n

    for i, item in enumerate(items[:n]):
        y = Emu(CONTENT_TOP + i * (cardH + gap))
        text = item if isinstance(item, str) else item.get("text", "")
        children = [] if isinstance(item, str) else item.get("children", [])
        col = ACCENTS[i % len(ACCENTS)]
        # Colored label block
        _rect(slide, LM, y, Inches(2.3), Emu(cardH), col)
        _tb(slide, Emu(LM + Inches(0.12)), y, Inches(2.05), Emu(cardH),
            text, BODY_FONT, 11, BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=15)
        # Children on dark card
        _rect(slide, Inches(3.15), y, Inches(6.15), Emu(cardH), NEAR_BLACK)
        child_texts = []
        for ch in children:
            ch_text = ch if isinstance(ch, str) else ch.get("text", "")
            grandchildren = [] if isinstance(ch, str) else ch.get("children", [])
            child_texts.append(ch_text)
            for gc in grandchildren:
                gc_text = gc if isinstance(gc, str) else gc.get("text", "")
                child_texts.append("  \u2022 " + gc_text)
        body = "\n".join(child_texts)
        _tb(slide, Inches(3.3), y, Inches(5.85), Emu(cardH),
            body, BODY_FONT, 10, BRIGHT, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=14)
def build_text_split(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Key Point"), LIME)
    headline = c.get("headline", "")
    detail = c.get("detail", "")
    points = c.get("points", [])
    _rect(slide, LM, CONTENT_TOP, Inches(4.1), Inches(3.9), NEAR_BLACK)
    _rect(slide, LM, CONTENT_TOP, Inches(0.1), Inches(3.9), LIME)
    _tb(slide, Emu(LM + Inches(0.25)), CONTENT_TOP, Inches(3.6), Inches(2.0),
        headline, TITLE_FONT, 20, WHITE, bold=True, line_spacing=28,
        valign=MSO_ANCHOR.MIDDLE)
    if detail:
        _tb(slide, Emu(LM + Inches(0.25)), Inches(3.1), Inches(3.6), Inches(1.4),
            detail, BODY_FONT, 11, DIM, line_spacing=16)
    for i, p in enumerate(points[:6]):
        py = Emu(CONTENT_TOP + Inches(0.05) + i * Inches(0.6))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, Inches(5.05), py, Inches(4.25), Inches(0.5), NEAR_BLACK)
        _rect(slide, Inches(5.05), py, Inches(0.1), Inches(0.5), col)
        _tb(slide, Inches(5.25), Emu(py), Inches(3.9), Inches(0.5),
            p, BODY_FONT, 11, BRIGHT, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)


def build_text_annotated(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Analysis"), VIOLET)
    items = c.get("items", [])
    n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.85))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))
    for i, item in enumerate(items[:n]):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, CW, Emu(rowH), NEAR_BLACK)
        _rect(slide, LM, y, Inches(1.5), Emu(rowH), col)
        _tb(slide, LM, y, Inches(1.5), Emu(rowH),
            item.get("label", ""), BODY_FONT, 10, BLACK, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, Emu(LM + Inches(1.7)), y, Inches(6.7), Emu(rowH),
            item.get("text", ""), BODY_FONT, 11, BRIGHT,
            valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


def build_icon_cards(prs, c):
    """Row of 2-3 cards, each with an icon/image above and title+detail below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Key Points"), LIME, size=24)
    items = c.get("items", []); n = min(len(items), 3)
    if n == 0: return
    gap = Inches(0.25)
    cardW = Emu((CW - (n - 1) * gap) / n)
    iconSize = Inches(0.8)
    iconTop = CONTENT_TOP
    cardTop = Emu(CONTENT_TOP + Inches(1.1))
    cardH = Emu(Inches(3.8) - Inches(1.1))
    for i, item in enumerate(items[:n]):
        x = Emu(LM + i * (cardW + gap))
        col = ACCENTS[i % len(ACCENTS)]
        icon_x = Emu(x + (cardW - iconSize) // 2)
        img_path = item.get("imagePath", "") if isinstance(item, dict) else ""
        shape_type = DEFAULT_ICON_SHAPES[i % len(DEFAULT_ICON_SHAPES)]
        _icon_or_image(slide, img_path, icon_x, iconTop, iconSize, iconSize, shape_type, col)
        _rect(slide, x, cardTop, cardW, cardH, NEAR_BLACK)
        title_text = item.get("title", "") if isinstance(item, dict) else str(item)
        detail_text = item.get("detail", "") if isinstance(item, dict) else ""
        _tb(slide, Emu(x + Inches(0.15)), Emu(cardTop + Inches(0.12)),
            Emu(cardW - Inches(0.3)), Inches(0.4),
            title_text, BODY_FONT, 13, WHITE, bold=True)
        if detail_text:
            _tb(slide, Emu(x + Inches(0.15)), Emu(cardTop + Inches(0.5)),
                Emu(cardW - Inches(0.3)), Emu(cardH - Inches(0.6)),
                detail_text, BODY_FONT, 10, DIM, line_spacing=14)


def build_feature_cards(prs, c):
    """1-2 full-width rows with an icon/image on the left and text on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _content_header(slide, c.get("title", "Features"), LIME, size=24)
    items = c.get("items", []); n = min(len(items), 2)
    if n == 0: return
    gap = Inches(0.25)
    rowH = Emu((Inches(3.8) - (n - 1) * gap) / n)
    iconSize = Inches(1.1)
    for i, item in enumerate(items[:n]):
        y = Emu(CONTENT_TOP + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, y, CW, rowH, NEAR_BLACK)
        icon_y = Emu(y + (rowH - iconSize) // 2)
        icon_x = Emu(LM + Inches(0.3))
        img_path = item.get("imagePath", "") if isinstance(item, dict) else ""
        shape_type = DEFAULT_ICON_SHAPES[i % len(DEFAULT_ICON_SHAPES)]
        _icon_or_image(slide, img_path, icon_x, icon_y, iconSize, iconSize, shape_type, col)
        text_x = Emu(LM + Inches(1.7))
        text_w = Emu(CW - Inches(1.9))
        title_text = item.get("title", "") if isinstance(item, dict) else str(item)
        detail_text = item.get("detail", "") if isinstance(item, dict) else ""
        _tb(slide, text_x, Emu(y + Inches(0.2)),
            text_w, Inches(0.4),
            title_text, BODY_FONT, 14, WHITE, bold=True)
        if detail_text:
            _tb(slide, text_x, Emu(y + Inches(0.65)),
                text_w, Emu(rowH - Inches(0.8)),
                detail_text, BODY_FONT, 11, DIM, line_spacing=16)


# ── Dispatch ─────────────────────────────────────────────
BUILDERS = {
    "title": build_title, "in_brief": build_in_brief,
    "section_divider": build_section_divider, "stat_callout": build_stat_callout,
    "quote": build_quote, "comparison": build_comparison,
    "text_graph": build_text_graph, "process_flow": build_process_flow,
    "matrix": build_matrix, "methods": build_methods,
    "hypotheses": build_hypotheses, "wsn_dense": build_wsn_dense,
    "wsn_reveal": build_wsn_reveal, "findings_recs": build_findings_recs,
    "findings_recs_dense": build_findings_recs_dense,
    "open_questions": build_open_questions, "agenda": build_agenda,
    "progressive_reveal": build_progressive_reveal, "closer": build_closer,
    "timeline": build_timeline, "data_table": build_data_table,
    "multi_stat": build_multi_stat, "persona": build_persona,
    "risk_tradeoff": build_risk_tradeoff, "appendix": build_appendix,
    "before_after": build_before_after, "summary": build_summary,
    "quote_full": build_quote_full, "stat_hero": build_stat_hero,
    "in_brief_featured": build_in_brief_featured, "in_brief_reveal": build_in_brief_reveal,
    "persona_duo": build_persona_duo,
    "process_flow_vertical": build_process_flow_vertical,
    "process_flow_reveal": build_process_flow_reveal,
    "text_cards": build_text_cards, "text_columns": build_text_columns,
    "text_narrative": build_text_narrative, "text_nested": build_text_nested,
    "text_split": build_text_split, "text_annotated": build_text_annotated,
    "icon_cards": build_icon_cards, "feature_cards": build_feature_cards,
}

def build_deck(slide_configs, output_path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    for slide_type, data in slide_configs:
        if slide_type == "skip": continue
        builder = BUILDERS.get(slide_type)
        if builder: builder(prs, data)
    prs.save(output_path); return output_path
