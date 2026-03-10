"""
Editorial template — FT / Monocle / polished infographic aesthetic.
Signature: serif headlines, warm paper background, hairline rules, restrained color.
Typography and whitespace do the heavy lifting.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

W = Inches(10)
H = Inches(5.625)

# ── Palette ──────────────────────────────────────────────
DK_GREEN   = RGBColor(0x04, 0x40, 0x14)   # title/section/closer bg
GOLD       = RGBColor(0xD4, 0xA8, 0x43)   # accent — stat numbers, item markers
GOLD_MUTED = RGBColor(0xC4, 0x9E, 0x3C)   # slightly deeper gold for small text
CREAM      = RGBColor(0xE8, 0xE0, 0xCC)   # muted cream for text on dark bg
PAPER      = RGBColor(0xF5, 0xF3, 0xEE)   # warm off-white content bg
CHARCOAL   = RGBColor(0x2D, 0x2D, 0x2D)   # primary text
MID        = RGBColor(0x55, 0x55, 0x55)    # secondary text
QUIET      = RGBColor(0x88, 0x88, 0x88)    # captions, sources
RULE_CLR   = RGBColor(0xCC, 0xC7, 0xBA)    # hairline rules on content slides
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0xA8, 0xC4, 0xA0)  # soft green for accents on dark bg

# Section colors (for sectionColor wayfinding)
GREEN  = RGBColor(0x36, 0x87, 0x27)
BLUE   = RGBColor(0x38, 0x80, 0xF3)
PURPLE = RGBColor(0x5B, 0x2C, 0x8F)
COBALT = RGBColor(0x04, 0x54, 0x7C)

SECTION_COLORS = {
    "green": DK_GREEN, "blue": BLUE, "purple": PURPLE,
    "cobalt": COBALT, "gold": GOLD,
    "red": RGBColor(0xC2, 0x3B, 0x22), "teal": RGBColor(0x1B, 0x7A, 0x6E),
    "ochre": RGBColor(0xCC, 0x7A, 0x2E), "slate": RGBColor(0x5A, 0x6A, 0x7A),
    "plum": RGBColor(0x8E, 0x45, 0x85),
}

# Muted content accent colors for items (restrained, not loud)
ACCENTS = [
    RGBColor(0x04, 0x40, 0x14),  # dark green
    RGBColor(0x04, 0x54, 0x7C),  # cobalt
    RGBColor(0x5B, 0x2C, 0x8F),  # purple
    RGBColor(0xD4, 0xA8, 0x43),  # gold
    RGBColor(0x1B, 0x7A, 0x6E),  # teal
]

def _resolve_section_color(c, default=None):
    sc = c.get("sectionColor", "")
    return SECTION_COLORS.get(sc, default or DK_GREEN)

TITLE_FONT = "Fidelity Slab"
TITLE_FONT_FALLBACK = "Georgia"
BODY_FONT  = "Fidelity Sans"
BODY_FONT_FALLBACK = "Calibri"

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "logo.png")

LM = Inches(0.8)        # left margin
RM = Inches(0.8)        # right margin
CW = Inches(8.4)        # content width


# ── Helpers ──────────────────────────────────────────────

def _font_fallback(font_name):
    if font_name == TITLE_FONT: return TITLE_FONT, TITLE_FONT_FALLBACK
    if font_name == BODY_FONT: return BODY_FONT, BODY_FONT_FALLBACK
    return font_name, None

def _set_font_with_fallback(font_obj, font_name):
    primary, fallback = _font_fallback(font_name)
    font_obj.name = primary
    if fallback:
        el = font_obj._element
        latin = el.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(el, qn('a:latin'))
            latin.set('typeface', primary)
        cs = el.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(el, qn('a:cs'))
        cs.set('typeface', fallback)

def _rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

def _oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

def _line(slide, x, y, w, color=RULE_CLR, thickness=0.75):
    """Thin horizontal hairline rule."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(thickness))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def _vline(slide, x, y, h, color=RULE_CLR, thickness=0.75):
    """Thin vertical hairline rule."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(thickness), h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def _tb(slide, x, y, w, h, text, font=None, size=12, color=None,
        bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
        line_spacing=None):
    font = font or BODY_FONT; color = color or CHARCOAL
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    _set_font_with_fallback(p.font, font); p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.italic = italic
    p.alignment = align
    if line_spacing: p.line_spacing = Pt(line_spacing)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', {MSO_ANCHOR.TOP:'t', MSO_ANCHOR.MIDDLE:'ctr', MSO_ANCHOR.BOTTOM:'b'}[valign])
    return box

def _paper_bg(slide):
    """Set warm paper background."""
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = PAPER

def _dark_bg(slide):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = DK_GREEN

def _content_title(slide, title, size=26):
    """Editorial-style page title: serif, left-aligned, hairline below."""
    _tb(slide, LM, Inches(0.35), CW, Inches(0.65), title,
        TITLE_FONT, size, CHARCOAL, bold=True, valign=MSO_ANCHOR.BOTTOM)


CONTENT_TOP = Inches(1.25)


# ── Builders ─────────────────────────────────────────────

def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    # Thin gold line as top accent
    _rect(slide, Inches(0), Inches(0), W, Inches(0.06), GOLD)
    _tb(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.8),
        c.get("title", "Title"), TITLE_FONT, 42, WHITE, bold=True,
        valign=MSO_ANCHOR.BOTTOM)
    _line(slide, Inches(0.8), Inches(2.95), Inches(2.2), GOLD, 2)
    if c.get("subtitle"):
        _tb(slide, Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.5),
            c["subtitle"], BODY_FONT, 16, CREAM)
    meta = []
    if c.get("author"): meta.append(c["author"])
    if c.get("date"): meta.append(c["date"])
    if meta:
        _tb(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.7),
            "\n".join(meta), BODY_FONT, 12, LIGHT_GREEN)
    # Logo in upper-right
    logo = c.get("logo_path", LOGO_PATH)
    if logo and os.path.isfile(logo):
        slide.shapes.add_picture(logo, Inches(7.6), Inches(0.3), Inches(2.0))


def build_in_brief(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "In Brief"))
    bullets = c.get("bullets", [])
    n = len(bullets)
    avail = H - CONTENT_TOP - Inches(0.4)
    rowH = min(int((avail - (n - 1) * Inches(0.08)) / n), Inches(0.95))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, b in enumerate(bullets):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Small number
        _tb(slide, LM, y, Inches(0.35), rowH,
            str(i + 1), TITLE_FONT, 20, GOLD, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        # Text
        _tb(slide, Emu(LM + Inches(0.5)), y, Inches(7.8), rowH,
            b, BODY_FONT, 14, CHARCOAL, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=19)
        # Thin separator below (not on last)
        if i < n - 1:
            _line(slide, Emu(LM + Inches(0.5)), Emu(y + rowH + gap // 2),
                  Inches(7.8))


def build_section_divider(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    # Gold top accent
    _rect(slide, Inches(0), Inches(0), W, Inches(0.06), GOLD)
    if c.get("sectionNumber"):
        _tb(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.8),
            f"0{c['sectionNumber']}", TITLE_FONT, 48, GOLD, bold=True,
            valign=MSO_ANCHOR.BOTTOM)
    _tb(slide, Inches(0.8), Inches(2.1), Inches(8), Inches(1.0),
        c.get("title", "Section"), TITLE_FONT, 36, WHITE, bold=True,
        valign=MSO_ANCHOR.MIDDLE)
    _line(slide, Inches(0.8), Inches(3.2), Inches(2.0), GOLD, 2)
    if c.get("subtitle"):
        _tb(slide, Inches(0.8), Inches(3.45), Inches(8), Inches(0.5),
            c["subtitle"], BODY_FONT, 14, CREAM)


def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    sc = _resolve_section_color(c)
    _content_title(slide, c.get("title", "Key Metric"), size=22)
    # Thin rule above stat
    _line(slide, Inches(2.5), Inches(1.45), Inches(5.0), RULE_CLR, 1)
    # Big number
    _tb(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.8),
        c.get("stat", "—"), TITLE_FONT, 80, GOLD, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # Thin rule below stat
    _line(slide, Inches(2.5), Inches(3.3), Inches(5.0), RULE_CLR, 1)
    if c.get("headline"):
        _tb(slide, Inches(1.5), Inches(3.5), Inches(7), Inches(0.6),
            c["headline"], BODY_FONT, 17, CHARCOAL, bold=True,
            align=PP_ALIGN.CENTER)
    if c.get("detail"):
        _tb(slide, Inches(1.5), Inches(4.1), Inches(7), Inches(0.7),
            c["detail"], BODY_FONT, 12, MID, align=PP_ALIGN.CENTER,
            line_spacing=17)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(4.9), Inches(9), Inches(0.3),
            c["source"], BODY_FONT, 9, QUIET, italic=True, align=PP_ALIGN.CENTER)


def build_quote(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "In Their Words"), size=22)
    # Large open-quote in muted gold
    _tb(slide, Inches(0.6), Inches(1.1), Inches(1), Inches(1.2),
        "\u201C", TITLE_FONT, 80, CREAM, bold=True)
    # Quote text
    _tb(slide, Inches(1.3), Inches(1.6), Inches(7.2), Inches(2.0),
        c.get("quote", ""), TITLE_FONT, 18, CHARCOAL, italic=True,
        valign=MSO_ANCHOR.MIDDLE, line_spacing=26)
    # Attribution line
    _line(slide, Inches(1.3), Inches(3.8), Inches(1.2), GOLD, 1.5)
    if c.get("attribution"):
        _tb(slide, Inches(1.3), Inches(3.95), Inches(7), Inches(0.4),
            c["attribution"], BODY_FONT, 12, MID, bold=True)
    if c.get("context"):
        _tb(slide, Inches(1.3), Inches(4.3), Inches(7), Inches(0.35),
            c["context"], BODY_FONT, 10, QUIET, italic=True)


def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Comparison"), size=24)

    # Left column
    _tb(slide, LM, Emu(CONTENT_TOP + Inches(0.05)), Inches(3.8), Inches(0.35),
        c.get("leftLabel", "Before"), TITLE_FONT, 14, COBALT, bold=True)
    _line(slide, LM, Emu(CONTENT_TOP + Inches(0.45)), Inches(3.8), COBALT, 1.5)
    for i, item in enumerate(c.get("leftItems", [])):
        iy = Emu(CONTENT_TOP + Inches(0.6) + i * Inches(0.6))
        _tb(slide, LM, iy, Inches(3.8), Inches(0.55),
            item, BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE)

    # Center divider
    _vline(slide, Inches(4.95), CONTENT_TOP, Inches(3.5), RULE_CLR, 1)

    # Right column
    _tb(slide, Inches(5.3), Emu(CONTENT_TOP + Inches(0.05)), Inches(3.9), Inches(0.35),
        c.get("rightLabel", "After"), TITLE_FONT, 14, DK_GREEN, bold=True)
    _line(slide, Inches(5.3), Emu(CONTENT_TOP + Inches(0.45)), Inches(3.9), DK_GREEN, 1.5)
    for i, item in enumerate(c.get("rightItems", [])):
        iy = Emu(CONTENT_TOP + Inches(0.6) + i * Inches(0.6))
        _tb(slide, Inches(5.3), iy, Inches(3.9), Inches(0.55),
            item, BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE)


def build_text_graph(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Title"), size=24)
    texts = c.get("text", [])
    if not isinstance(texts, list): texts = [texts]
    for i, t in enumerate(texts):
        _tb(slide, LM, Emu(CONTENT_TOP + i * Inches(1.2)),
            Inches(3.9), Inches(1.1), t, BODY_FONT, 12, CHARCOAL,
            line_spacing=17)
    _vline(slide, Inches(5.0), CONTENT_TOP, Inches(3.5), RULE_CLR, 1)

    from pptx.chart.data import CategoryChartData
    raw = c.get("chartData", [{"name": "S1", "labels": ["A","B","C"], "values": [25,45,30]}])
    cd_obj = CategoryChartData()
    if raw:
        cd = raw[0]
        cd_obj.categories = cd.get("labels", ["A","B","C"])
        cd_obj.add_series(cd.get("name","Series 1"), cd.get("values",[25,45,30]))
    ct = {"line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}.get(
        c.get("chartType","bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    slide.shapes.add_chart(ct, Inches(5.3), CONTENT_TOP, Inches(4.0), Inches(3.6), cd_obj)


def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Process"), size=24)
    steps = c.get("steps", [])
    count = min(len(steps), 5)
    if count == 0: return

    total_w = 8.4
    gap_w = 0.3
    step_w = (total_w - (count - 1) * gap_w) / count
    startX = 0.8
    stepY = 1.3
    stepH = 3.8

    for i, step in enumerate(steps[:count]):
        x = Inches(startX + i * (step_w + gap_w))
        y = Inches(stepY)
        sw = Inches(step_w)
        sh = Inches(stepH)
        col = ACCENTS[i % len(ACCENTS)]

        # Thin color line at top of column
        _line(slide, x, y, sw, col, 2.5)
        # Step number — large, editorial
        _tb(slide, x, Emu(y + Inches(0.15)), sw, Inches(0.55),
            str(i + 1), TITLE_FONT, 28, GOLD, bold=True)
        # Step title
        _tb(slide, x, Emu(y + Inches(0.7)), sw, Inches(0.65),
            step.get("title", ""), BODY_FONT, 12, CHARCOAL, bold=True)
        if step.get("detail"):
            _tb(slide, x, Emu(y + Inches(1.4)), sw, Inches(2.2),
                step["detail"], BODY_FONT, 10, MID, line_spacing=14)


def build_matrix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Framework"), size=24)
    quads = c.get("quadrants", [{}, {}, {}, {}])
    qW = Inches(4.0)
    qH = Inches(1.7)
    gap = Inches(0.3)
    sX = LM
    sY = CONTENT_TOP

    for i, q in enumerate(quads[:4]):
        col_idx = i % 2; row = i // 2
        x = Emu(sX + col_idx * (qW + gap))
        y = Emu(sY + row * (qH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Thin top border
        _line(slide, x, y, qW, col, 2)
        _tb(slide, Emu(x + Inches(0.1)), Emu(y + Inches(0.12)),
            Emu(qW - Inches(0.2)), Inches(0.3),
            q.get("label", ""), TITLE_FONT, 12, col, bold=True)
        _tb(slide, Emu(x + Inches(0.1)), Emu(y + Inches(0.45)),
            Emu(qW - Inches(0.2)), Emu(qH - Inches(0.55)),
            q.get("detail", ""), BODY_FONT, 11, MID, line_spacing=15)


def build_methods(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Approach"), size=24)
    fields = c.get("fields", [])
    for i, f in enumerate(fields):
        y = Emu(CONTENT_TOP + i * Inches(0.85))
        col = ACCENTS[i % len(ACCENTS)]
        # Label
        _tb(slide, LM, y, Inches(2.2), Inches(0.75),
            f.get("label", ""), TITLE_FONT, 12, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        # Thin vertical separator
        _vline(slide, Emu(LM + Inches(2.3)), Emu(y + Inches(0.1)),
               Inches(0.55), RULE_CLR, 1)
        # Value
        _tb(slide, Emu(LM + Inches(2.5)), y, Inches(5.7), Inches(0.75),
            f.get("value", ""), BODY_FONT, 12, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE, line_spacing=17)
        # Separator line below
        if i < len(fields) - 1:
            _line(slide, LM, Emu(y + Inches(0.8)), CW)


def build_hypotheses(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Hypotheses"), size=24)
    hyps = c.get("hypotheses", [])
    n = len(hyps)
    avail = H - CONTENT_TOP - Inches(0.35)
    rowH = min(int((avail - (n - 1) * Inches(0.08)) / n), Inches(0.75))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, h in enumerate(hyps):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Thin left accent
        _vline(slide, LM, y, rowH, col, 3)
        # H-number
        _tb(slide, Emu(LM + Inches(0.15)), y, Inches(0.5), rowH,
            f"H{i+1}", TITLE_FONT, 13, col, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        # Text
        _tb(slide, Emu(LM + Inches(0.6)), y, Inches(6.0), rowH,
            h.get("text", ""), BODY_FONT, 12, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE)
        # Status
        if h.get("status"):
            sc = DK_GREEN if h["status"] == "Confirmed" else (RGBColor(0xC2,0x3B,0x22) if h["status"] == "Rejected" else MID)
            _tb(slide, Inches(7.6), y, Inches(1.4), rowH,
                h["status"], BODY_FONT, 10, sc, bold=True,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_wsn_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Key Finding"), size=24)
    cols_data = [
        ("What", DK_GREEN, c.get("what", {})),
        ("So What", COBALT, c.get("soWhat", {})),
        ("Now What", PURPLE, c.get("nowWhat", {})),
    ]
    colW = Inches(2.65)
    gap = Inches(0.2)
    cardH = Inches(3.8)
    startX = LM

    for i, (label, color, data) in enumerate(cols_data):
        x = Emu(startX + i * (colW + gap))
        # Color top line
        _line(slide, x, CONTENT_TOP, colW, color, 2.5)
        _tb(slide, Emu(x + Inches(0.05)), Emu(CONTENT_TOP + Inches(0.12)),
            Emu(colW - Inches(0.1)), Inches(0.3),
            label, TITLE_FONT, 13, color, bold=True)
        _tb(slide, Emu(x + Inches(0.05)), Emu(CONTENT_TOP + Inches(0.5)),
            Emu(colW - Inches(0.1)), Inches(1.1),
            data.get("headline", ""), BODY_FONT, 12, CHARCOAL, bold=True,
            line_spacing=17)
        if data.get("detail"):
            _tb(slide, Emu(x + Inches(0.05)), Emu(CONTENT_TOP + Inches(1.65)),
                Emu(colW - Inches(0.1)), Inches(2.0),
                data["detail"], BODY_FONT, 10, MID, line_spacing=14)


def build_wsn_reveal(prs, c):
    """Progressive 3-slide reveal: What → So What → Now What."""
    WSN_COLORS = [DK_GREEN, COBALT, PURPLE]
    WSN_LABELS = ["What We Found", "So What", "Now What"]
    WSN_KEYS   = ["what", "soWhat", "nowWhat"]

    def _hdr(slide):
        _paper_bg(slide)
        _tb(slide, LM, Inches(0.35), CW, Inches(0.55),
            c.get("title", "Key Finding"), TITLE_FONT, 26, CHARCOAL, bold=True,
            valign=MSO_ANCHOR.BOTTOM)
        _line(slide, LM, Inches(0.95), Inches(1.8), GOLD_MUTED, 1.5)

    def _zone(slide, x, w, label, color, data, condensed=False):
        y = Inches(1.2)
        h = Inches(1.9) if condensed else Inches(3.5)
        _line(slide, x, y, w, color, 2.5)
        _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.1)),
            Inches(2), Inches(0.3), label, TITLE_FONT,
            11 if condensed else 13, color, bold=True)
        _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.45)),
            Emu(w - Inches(0.1)), Inches(0.7) if condensed else Inches(1.0),
            data.get("headline", ""), BODY_FONT,
            11 if condensed else 13, CHARCOAL, bold=True,
            line_spacing=16 if condensed else 18)
        if data.get("detail"):
            detail_y = Emu(y + Inches(1.15)) if condensed else Emu(y + Inches(1.5))
            _tb(slide, Emu(x + Inches(0.05)), detail_y,
                Emu(w - Inches(0.1)), Inches(0.55) if condensed else Inches(1.8),
                data["detail"], BODY_FONT,
                9 if condensed else 11, MID, line_spacing=13 if condensed else 15)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6]); _hdr(s1)
    _zone(s1, LM, Inches(5.5), "What We Found", DK_GREEN, c.get("what", {}))

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); _hdr(s2)
    _zone(s2, LM, Inches(4.0), "What We Found", DK_GREEN, c.get("what", {}))
    _zone(s2, Inches(5.1), Inches(4.1), "So What", COBALT, c.get("soWhat", {}))

    # Slide 3: condensed top + full-width Now What
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); _hdr(s3)
    _zone(s3, LM, Inches(4.0), "What We Found", DK_GREEN, c.get("what", {}), True)
    _zone(s3, Inches(5.1), Inches(4.1), "So What", COBALT, c.get("soWhat", {}), True)
    # Now What — full width
    nwY = Inches(3.35)
    _line(s3, LM, nwY, CW, PURPLE, 2.5)
    nw = c.get("nowWhat", {})
    _tb(s3, Emu(LM + Inches(0.05)), Emu(nwY + Inches(0.1)),
        Inches(2), Inches(0.3), "Now What", TITLE_FONT, 13, PURPLE, bold=True)
    _tb(s3, Emu(LM + Inches(0.05)), Emu(nwY + Inches(0.45)),
        Emu(CW - Inches(0.1)), Inches(0.6),
        nw.get("headline", ""), BODY_FONT, 16, CHARCOAL, bold=True)
    if nw.get("detail"):
        _tb(s3, Emu(LM + Inches(0.05)), Emu(nwY + Inches(1.1)),
            Emu(CW - Inches(0.1)), Inches(0.5),
            nw["detail"], BODY_FONT, 11, MID)


def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Findings & Recommendations"), size=22)
    items = c.get("items", [])
    n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int((avail - (n - 1) * Inches(0.1)) / n), Inches(0.82))
    gap = Inches(0.1)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    # Column headers
    _tb(slide, LM, Emu(startY - Inches(0.3)), Inches(3.8), Inches(0.25),
        "FINDING", BODY_FONT, 9, QUIET, bold=True)
    _tb(slide, Inches(5.4), Emu(startY - Inches(0.3)), Inches(3.8), Inches(0.25),
        "RECOMMENDATION", BODY_FONT, 9, QUIET, bold=True)

    for i, item in enumerate(items[:5]):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Finding
        _vline(slide, LM, y, rowH, col, 3)
        _tb(slide, Emu(LM + Inches(0.15)), y, Inches(3.7), rowH,
            item.get("finding", ""), BODY_FONT, 11, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE)
        # Arrow
        _tb(slide, Inches(4.7), y, Inches(0.5), rowH,
            "\u2192", BODY_FONT, 16, GOLD, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE)
        # Recommendation
        _vline(slide, Inches(5.3), y, rowH, GOLD, 3)
        _tb(slide, Inches(5.5), y, Inches(3.7), rowH,
            item.get("recommendation", ""), BODY_FONT, 11, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Complete Findings"), size=22)
    items = c.get("items", [])
    n = min(len(items), 8)
    avail = H - CONTENT_TOP - Inches(0.2)
    rowH = min(int((avail - (n - 1) * Inches(0.04)) / n), Inches(0.48))
    gap = Inches(0.04)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    # Column headers
    _tb(slide, LM, Emu(startY - Inches(0.25)), Inches(3.8), Inches(0.22),
        "FINDING", BODY_FONT, 8, QUIET, bold=True)
    _tb(slide, Inches(5.3), Emu(startY - Inches(0.25)), Inches(3.9), Inches(0.22),
        "RECOMMENDATION", BODY_FONT, 8, QUIET, bold=True)

    for i, item in enumerate(items[:8]):
        y = Emu(startY + i * (rowH + gap))
        # Finding
        _vline(slide, LM, y, rowH, DK_GREEN, 2)
        _tb(slide, Emu(LM + Inches(0.12)), y, Inches(3.8), rowH,
            item.get("finding", ""), BODY_FONT, 9, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE)
        # Rec
        _vline(slide, Inches(5.2), y, rowH, GOLD, 2)
        _tb(slide, Inches(5.35), y, Inches(3.85), rowH,
            item.get("recommendation", ""), BODY_FONT, 9, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE)
        # Separator
        if i < n - 1:
            _line(slide, LM, Emu(y + rowH), CW, RULE_CLR, 0.5)


def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Open Questions"), size=24)
    questions = c.get("questions", [])
    cardW = Inches(4.0)
    cardH = Inches(1.7)
    gX = Inches(0.4)
    gY = Inches(0.25)

    for i, question in enumerate(questions[:4]):
        col_idx = i % 2; row = i // 2
        x = Emu(LM + col_idx * (cardW + gX))
        y = Emu(CONTENT_TOP + row * (cardH + gY))
        col = ACCENTS[i % len(ACCENTS)]
        # Top color line
        _line(slide, x, y, cardW, col, 2.5)
        # Number
        _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.1)),
            Inches(0.5), Inches(0.5), str(i + 1),
            TITLE_FONT, 24, GOLD, bold=True)
        # Question text
        _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.65)),
            Emu(cardW - Inches(0.1)), Inches(0.9),
            question, BODY_FONT, 12, CHARCOAL, line_spacing=17)


def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Agenda"), size=26)
    items = c.get("items", [])
    n = len(items)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min(int(avail / n), Inches(0.72))
    gap = Inches(0.08)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.25))

    for i, item in enumerate(items):
        y = Emu(startY + i * (rowH + gap))
        # Number
        _tb(slide, LM, y, Inches(0.5), rowH,
            str(i + 1), TITLE_FONT, 22, GOLD, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        title_text = item if isinstance(item, str) else item.get("title", "")
        _tb(slide, Emu(LM + Inches(0.55)), y, Inches(5.5), rowH,
            title_text, BODY_FONT, 15, CHARCOAL, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        if isinstance(item, dict) and item.get("detail"):
            _tb(slide, Inches(7.2), y, Inches(2.0), rowH,
                item["detail"], BODY_FONT, 11, QUIET,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            _line(slide, Emu(LM + Inches(0.55)),
                  Emu(y + rowH + gap // 2), Inches(7.8))


def build_progressive_reveal(prs, c):
    takeaways = c.get("takeaways", [])
    for n in range(min(len(takeaways), 5)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _paper_bg(slide)
        _tb(slide, LM, Inches(0.25), CW, Inches(0.55),
            c.get("title", "Building the Picture"), TITLE_FONT, 24, CHARCOAL, bold=True)
        _line(slide, LM, Inches(0.85), Inches(1.8), GOLD_MUTED, 1.5)

        cur = takeaways[n]
        col = ACCENTS[n % len(ACCENTS)]
        # Main content zone
        _line(slide, LM, Inches(1.1), CW, col, 2.5)
        _tb(slide, Emu(LM + Inches(0.05)), Inches(1.2), Emu(CW - Inches(0.1)), Inches(0.55),
            cur.get("headline", ""), BODY_FONT, 15, CHARCOAL, bold=True)
        if cur.get("detail"):
            _tb(slide, Emu(LM + Inches(0.05)), Inches(1.8), Emu(CW - Inches(0.1)), Inches(1.3),
                cur["detail"], BODY_FONT, 11, MID, line_spacing=15)

        # Separator
        _line(slide, Inches(0), Inches(3.35), W, RULE_CLR, 1)
        _tb(slide, LM, Inches(3.45), Inches(3), Inches(0.22),
            "Running Takeaways", BODY_FONT, 9, GOLD_MUTED, bold=True)

        for j in range(n + 1):
            ty = Emu(Inches(3.75) + j * Inches(0.38))
            active = j == n
            jcol = ACCENTS[j % len(ACCENTS)]
            # Small square marker
            _rect(slide, LM, Emu(ty + Inches(0.07)), Inches(0.1), Inches(0.1), jcol)
            _tb(slide, Emu(LM + Inches(0.2)), ty, Emu(CW - Inches(0.2)), Inches(0.32),
                takeaways[j].get("summary", takeaways[j].get("headline", "")),
                BODY_FONT, 10, CHARCOAL if active else QUIET,
                bold=active, valign=MSO_ANCHOR.MIDDLE)


def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sc = _resolve_section_color(c)
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = sc
    _rect(slide, Inches(0), Inches(0), W, Inches(0.06), GOLD)
    _tb(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(1.2),
        c.get("title", "Thank You"), TITLE_FONT, 44, WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
    _line(slide, Inches(3.75), Inches(2.75), Inches(2.5), GOLD, 2)
    if c.get("subtitle"):
        _tb(slide, Inches(0.5), Inches(2.95), Inches(9), Inches(0.5),
            c["subtitle"], BODY_FONT, 16, CREAM, align=PP_ALIGN.CENTER)
    if c.get("contact"):
        _tb(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
            c["contact"], BODY_FONT, 12, LIGHT_GREEN, align=PP_ALIGN.CENTER)




def build_timeline(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Timeline"), size=24)
    milestones = c.get("milestones", [])
    n = min(len(milestones), 6)
    if n == 0: return
    lineY = Inches(2.8)
    startX = 0.8
    endX = 9.2
    step = (endX - startX) / max(n - 1, 1)
    # Horizontal line
    _line(slide, Inches(startX), lineY, Inches(endX - startX), RULE_CLR, 1.5)
    STATUS_CLR = {"complete": DK_GREEN, "current": GOLD, "upcoming": QUIET}
    for i, m in enumerate(milestones[:n]):
        cx = Inches(startX + i * step)
        col = STATUS_CLR.get(m.get("status", "upcoming"), QUIET)
        # Dot
        ds = Inches(0.18)
        _oval(slide, Emu(cx - ds // 2), Emu(lineY - ds // 2), ds, ds, col)
        # Date above
        _tb(slide, Emu(cx - Inches(0.6)), Emu(lineY - Inches(0.7)),
            Inches(1.2), Inches(0.35), m.get("date", ""),
            BODY_FONT, 10, col, bold=True, align=PP_ALIGN.CENTER)
        # Title below
        _tb(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.2)),
            Inches(1.4), Inches(0.4), m.get("title", ""),
            BODY_FONT, 11, CHARCOAL, bold=True, align=PP_ALIGN.CENTER)
        if m.get("detail"):
            _tb(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.6)),
                Inches(1.4), Inches(0.7), m["detail"],
                BODY_FONT, 9, MID, align=PP_ALIGN.CENTER, line_spacing=12)


def build_data_table(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Data"), size=24)
    headers = c.get("headers", [])
    rows = c.get("rows", [])
    hc = c.get("highlightCol", None)
    nCols = len(headers)
    nRows = min(len(rows), 10)
    if nCols == 0: return
    tW = 8.4
    tX = 0.8
    colW = tW / nCols
    rowH = min(0.42, 3.8 / max(nRows + 1, 1))
    hdrY = 1.25
    # Header row
    _rect(slide, Inches(tX), Inches(hdrY), Inches(tW), Inches(rowH), DK_GREEN)
    for j, h in enumerate(headers):
        _tb(slide, Inches(tX + j * colW + 0.08), Inches(hdrY),
            Inches(colW - 0.16), Inches(rowH), h,
            BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    # Data rows
    for i, row in enumerate(rows[:nRows]):
        ry = hdrY + (i + 1) * rowH
        bg = PAPER if i % 2 == 0 else RGBColor(0xED, 0xEB, 0xE5)
        _rect(slide, Inches(tX), Inches(ry), Inches(tW), Inches(rowH), bg)
        for j, cell in enumerate(row[:nCols]):
            is_hl = hc is not None and j == hc
            _tb(slide, Inches(tX + j * colW + 0.08), Inches(ry),
                Inches(colW - 0.16), Inches(rowH), str(cell),
                BODY_FONT, 10, DK_GREEN if is_hl else CHARCOAL,
                bold=is_hl, valign=MSO_ANCHOR.MIDDLE)
    if c.get("note"):
        _tb(slide, Inches(tX), Inches(hdrY + (nRows + 1) * rowH + 0.1),
            Inches(tW), Inches(0.3), c["note"],
            BODY_FONT, 8, QUIET, italic=True)


def build_multi_stat(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Key Metrics"), size=24)
    stats = c.get("stats", [])
    n = min(len(stats), 4)
    if n == 0: return
    totalW = 8.4
    gap = 0.3
    statW = (totalW - (n - 1) * gap) / n
    for i, s in enumerate(stats[:n]):
        x = Inches(0.8 + i * (statW + gap))
        col = ACCENTS[i % len(ACCENTS)]
        _line(slide, x, Inches(1.4), Inches(statW), col, 2.5)
        _tb(slide, x, Inches(1.55), Inches(statW), Inches(1.4),
            s.get("value", "—"), TITLE_FONT, 52, GOLD, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _tb(slide, x, Inches(3.0), Inches(statW), Inches(0.45),
            s.get("label", ""), BODY_FONT, 13, CHARCOAL, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
        if s.get("detail"):
            _tb(slide, x, Inches(3.5), Inches(statW), Inches(0.8),
                s["detail"], BODY_FONT, 10, MID,
                align=PP_ALIGN.CENTER, line_spacing=14)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(4.9), Inches(9), Inches(0.3),
            c["source"], BODY_FONT, 9, QUIET, italic=True, align=PP_ALIGN.CENTER)


def build_persona(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Persona"), size=22)
    # Name and archetype
    _tb(slide, LM, Inches(1.2), Inches(4.0), Inches(0.55),
        c.get("name", ""), TITLE_FONT, 28, DK_GREEN, bold=True)
    _tb(slide, LM, Inches(1.75), Inches(4.0), Inches(0.35),
        c.get("archetype", ""), BODY_FONT, 14, GOLD, bold=True)
    _line(slide, LM, Inches(2.2), Inches(3.5), GOLD_MUTED, 1.5)
    # Traits
    traits = c.get("traits", [])
    for i, t in enumerate(traits[:5]):
        ty = Inches(2.4) + i * Inches(0.4)
        _rect(slide, LM, Emu(ty + Inches(0.08)), Inches(0.08), Inches(0.08), ACCENTS[i % len(ACCENTS)])
        _tb(slide, Emu(LM + Inches(0.2)), Emu(ty), Inches(3.8), Inches(0.35),
            t, BODY_FONT, 11, CHARCOAL)
    # Right column: strategy
    _vline(slide, Inches(5.1), Inches(1.2), Inches(3.5), RULE_CLR, 1)
    _tb(slide, Inches(5.4), Inches(1.2), Inches(3.8), Inches(0.3),
        "STRATEGY", BODY_FONT, 9, QUIET, bold=True)
    _tb(slide, Inches(5.4), Inches(1.55), Inches(3.8), Inches(1.5),
        c.get("strategy", ""), BODY_FONT, 12, CHARCOAL, line_spacing=17)
    if c.get("detail"):
        _tb(slide, Inches(5.4), Inches(3.2), Inches(3.8), Inches(1.5),
            c["detail"], BODY_FONT, 10, MID, line_spacing=14)


def build_risk_tradeoff(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Risk & Reward"), size=24)
    risks = c.get("risks", [])
    rewards = c.get("rewards", [])
    SEV = {"high": RGBColor(0xC2, 0x3B, 0x22), "medium": GOLD, "low": DK_GREEN}
    # Left: Risks
    _tb(slide, LM, CONTENT_TOP, Inches(3.8), Inches(0.3),
        "RISKS", BODY_FONT, 10, RGBColor(0xC2, 0x3B, 0x22), bold=True)
    _line(slide, LM, Emu(CONTENT_TOP + Inches(0.35)), Inches(3.8), RGBColor(0xC2, 0x3B, 0x22), 1.5)
    for i, r in enumerate(risks[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.5) + i * Inches(0.7))
        sc = SEV.get(r.get("severity", "medium"), GOLD)
        _rect(slide, LM, ry, Inches(0.06), Inches(0.6), sc)
        _tb(slide, Emu(LM + Inches(0.2)), ry, Inches(3.5), Inches(0.3),
            r.get("label", ""), BODY_FONT, 11, CHARCOAL, bold=True)
        if r.get("detail"):
            _tb(slide, Emu(LM + Inches(0.2)), Emu(ry + Inches(0.3)),
                Inches(3.5), Inches(0.3), r["detail"], BODY_FONT, 9, MID)
    # Divider
    _vline(slide, Inches(5.0), CONTENT_TOP, Inches(3.5), RULE_CLR, 1)
    # Right: Rewards
    _tb(slide, Inches(5.3), CONTENT_TOP, Inches(3.9), Inches(0.3),
        "REWARDS", BODY_FONT, 10, DK_GREEN, bold=True)
    _line(slide, Inches(5.3), Emu(CONTENT_TOP + Inches(0.35)), Inches(3.9), DK_GREEN, 1.5)
    for i, r in enumerate(rewards[:5]):
        ry = Emu(CONTENT_TOP + Inches(0.5) + i * Inches(0.7))
        _rect(slide, Inches(5.3), ry, Inches(0.06), Inches(0.6), DK_GREEN)
        _tb(slide, Inches(5.5), ry, Inches(3.6), Inches(0.3),
            r.get("label", ""), BODY_FONT, 11, CHARCOAL, bold=True)
        if r.get("detail"):
            _tb(slide, Inches(5.5), Emu(ry + Inches(0.3)),
                Inches(3.6), Inches(0.3), r["detail"], BODY_FONT, 9, MID)


def build_appendix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _tb(slide, LM, Inches(0.25), CW, Inches(0.45),
        c.get("title", "Appendix"), TITLE_FONT, 18, QUIET, bold=True)
    _line(slide, LM, Inches(0.72), Inches(1.2), RULE_CLR, 1)
    sections = c.get("sections", [])
    y_cursor = 0.85
    for s in sections:
        if y_cursor > 4.8: break
        _tb(slide, LM, Inches(y_cursor), CW, Inches(0.25),
            s.get("label", ""), BODY_FONT, 9, GOLD_MUTED, bold=True)
        _tb(slide, LM, Inches(y_cursor + 0.25), CW, Inches(0.8),
            s.get("content", ""), BODY_FONT, 8, MID, line_spacing=11)
        y_cursor += 1.1


def build_before_after(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Transformation"), size=24)
    bef = c.get("before", {})
    aft = c.get("after", {})
    interv = c.get("intervention", "")
    panelH = Inches(3.5)
    # Before
    _line(slide, LM, CONTENT_TOP, Inches(3.3), COBALT, 2.5)
    _tb(slide, LM, Emu(CONTENT_TOP + Inches(0.1)), Inches(3.3), Inches(0.3),
        bef.get("label", "Before"), TITLE_FONT, 13, COBALT, bold=True)
    _tb(slide, LM, Emu(CONTENT_TOP + Inches(0.5)), Inches(3.3), Inches(2.8),
        bef.get("detail", ""), BODY_FONT, 11, CHARCOAL, line_spacing=16)
    # Intervention arrow zone
    _tb(slide, Inches(4.25), Emu(CONTENT_TOP + Inches(0.8)), Inches(1.4), Inches(0.5),
        "\u2192", TITLE_FONT, 28, GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(4.05), Emu(CONTENT_TOP + Inches(1.4)), Inches(1.8), Inches(1.5),
        interv, BODY_FONT, 10, MID, align=PP_ALIGN.CENTER, line_spacing=14)
    # After
    _line(slide, Inches(6.0), CONTENT_TOP, Inches(3.2), DK_GREEN, 2.5)
    _tb(slide, Inches(6.0), Emu(CONTENT_TOP + Inches(0.1)), Inches(3.2), Inches(0.3),
        aft.get("label", "After"), TITLE_FONT, 13, DK_GREEN, bold=True)
    _tb(slide, Inches(6.0), Emu(CONTENT_TOP + Inches(0.5)), Inches(3.2), Inches(2.8),
        aft.get("detail", ""), BODY_FONT, 11, CHARCOAL, line_spacing=16)


def build_summary(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Summary"), size=26)
    sections = c.get("sections", [])
    n = len(sections)
    if n == 0: return
    colW = (8.4 - (n - 1) * 0.25) / n
    for i, sec in enumerate(sections[:4]):
        x = Inches(0.8 + i * (colW + 0.25))
        col = ACCENTS[i % len(ACCENTS)]
        _line(slide, x, CONTENT_TOP, Inches(colW), col, 2.5)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(0.1)), Inches(colW), Inches(0.35),
            sec.get("heading", ""), TITLE_FONT, 12, col, bold=True)
        points = sec.get("points", [])
        for j, p in enumerate(points[:5]):
            py = Emu(CONTENT_TOP + Inches(0.55) + j * Inches(0.55))
            _rect(slide, x, Emu(py + Inches(0.06)), Inches(0.06), Inches(0.06), col)
            _tb(slide, Emu(x + Inches(0.15)), Emu(py), Inches(colW - 0.2), Inches(0.5),
                p, BODY_FONT, 10, CHARCOAL, line_spacing=13)




def build_quote_full(prs, c):
    """Full-bleed dramatic quote — no title, dark bg, massive text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    _line(slide, Inches(3.5), Inches(1.0), Inches(3.0), GOLD, 2)
    _tb(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(2.5),
        c.get("quote", ""), TITLE_FONT, 26, WHITE, italic=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=36)
    _line(slide, Inches(3.5), Inches(4.0), Inches(3.0), GOLD, 2)
    if c.get("attribution"):
        _tb(slide, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.4),
            c["attribution"], BODY_FONT, 13, GOLD, bold=True, align=PP_ALIGN.CENTER)
    if c.get("context"):
        _tb(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.35),
            c["context"], BODY_FONT, 10, RGBColor(0x88, 0x88, 0x88), italic=True,
            align=PP_ALIGN.CENTER)


def build_stat_hero(prs, c):
    """One hero stat + 2-3 supporting stats below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Key Metric"), size=22)
    hero = c.get("hero", {})
    _tb(slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(1.3),
        hero.get("value", "—"), TITLE_FONT, 72, GOLD, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if hero.get("label"):
        _tb(slide, Inches(1.5), Inches(2.6), Inches(7.0), Inches(0.45),
            hero["label"], BODY_FONT, 16, CHARCOAL, bold=True, align=PP_ALIGN.CENTER)
    _line(slide, Inches(2.0), Inches(3.2), Inches(6.0), RULE_CLR, 1)
    supporting = c.get("supporting", [])
    n = min(len(supporting), 4)
    if n > 0:
        sw = 7.0 / n
        for i, s in enumerate(supporting[:n]):
            x = Inches(1.5 + i * sw)
            col = ACCENTS[i % len(ACCENTS)]
            _tb(slide, x, Inches(3.4), Inches(sw), Inches(0.8),
                s.get("value", ""), TITLE_FONT, 28, col, bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _tb(slide, x, Inches(4.2), Inches(sw), Inches(0.35),
                s.get("label", ""), BODY_FONT, 10, MID,
                align=PP_ALIGN.CENTER)
    if c.get("source"):
        _tb(slide, Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.3),
            c["source"], BODY_FONT, 9, QUIET, italic=True, align=PP_ALIGN.CENTER)


def build_in_brief_featured(prs, c):
    """One featured insight at top, supporting points below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "In Brief"), size=24)
    featured = c.get("featured", "")
    supporting = c.get("supporting", [])
    # Featured block with gold accent
    _line(slide, LM, CONTENT_TOP, CW, GOLD, 2.5)
    _tb(slide, LM, Emu(CONTENT_TOP + Inches(0.15)), CW, Inches(1.1),
        featured, BODY_FONT, 17, CHARCOAL, bold=True, line_spacing=24)
    # Supporting points below separator
    _line(slide, LM, Inches(2.55), CW, RULE_CLR, 1)
    n = len(supporting)
    for i, s in enumerate(supporting[:4]):
        y = Emu(Inches(2.75) + i * Inches(0.65))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, LM, Emu(y + Inches(0.08)), Inches(0.06), Inches(0.06), col)
        _tb(slide, Emu(LM + Inches(0.2)), Emu(y), CW, Inches(0.55),
            s, BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE)


def build_persona_duo(prs, c):
    """Two personas side by side for comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Archetype Comparison"), size=22)
    personas = c.get("personas", [{}, {}])

    for idx, p in enumerate(personas[:2]):
        x = Inches(0.8) if idx == 0 else Inches(5.2)
        w = Inches(4.0)
        col = ACCENTS[idx % len(ACCENTS)]
        _line(slide, x, CONTENT_TOP, w, col, 2.5)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(0.1)), w, Inches(0.45),
            p.get("name", ""), TITLE_FONT, 20, col, bold=True)
        _tb(slide, x, Emu(CONTENT_TOP + Inches(0.55)), w, Inches(0.3),
            p.get("archetype", ""), BODY_FONT, 11, GOLD, bold=True)
        traits = p.get("traits", [])
        for j, t in enumerate(traits[:4]):
            ty = Emu(CONTENT_TOP + Inches(1.0) + j * Inches(0.38))
            _rect(slide, x, Emu(ty + Inches(0.06)), Inches(0.06), Inches(0.06), col)
            _tb(slide, Emu(x + Inches(0.18)), Emu(ty), Emu(w - Inches(0.2)), Inches(0.33),
                t, BODY_FONT, 10, CHARCOAL)
        if p.get("strategy"):
            _line(slide, x, Emu(CONTENT_TOP + Inches(2.6)), Inches(1.5), RULE_CLR, 1)
            _tb(slide, x, Emu(CONTENT_TOP + Inches(2.7)), w, Inches(1.0),
                p["strategy"], BODY_FONT, 10, MID, line_spacing=14)
    # Center divider
    _vline(slide, Inches(4.95), CONTENT_TOP, Inches(3.5), RULE_CLR, 1)


def build_process_flow_vertical(prs, c):
    """Vertical 2-3 step flow with big connector arrows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Process"), size=24)
    steps = c.get("steps", [])
    count = min(len(steps), 3)
    if count == 0: return
    avail = 3.8
    stepH = (avail - (count - 1) * 0.4) / count
    for i, step in enumerate(steps[:count]):
        y = Inches(1.25 + i * (stepH + 0.4))
        col = ACCENTS[i % len(ACCENTS)]
        _line(slide, LM, Emu(y), CW, col, 2.5)
        _tb(slide, LM, Emu(y + Inches(0.08)), Inches(0.5), Inches(0.5),
            str(i + 1), TITLE_FONT, 24, GOLD, bold=True)
        _tb(slide, Emu(LM + Inches(0.6)), Emu(y + Inches(0.08)), Inches(3.0), Inches(0.4),
            step.get("title", ""), BODY_FONT, 14, CHARCOAL, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
        if step.get("detail"):
            _tb(slide, Emu(LM + Inches(0.6)), Emu(y + Inches(0.5)), Inches(7.5), Inches(stepH - 0.5),
                step["detail"], BODY_FONT, 11, MID, line_spacing=15)
        if i < count - 1:
            _tb(slide, Inches(1.0), Emu(y + Inches(stepH + 0.05)), Inches(0.5), Inches(0.3),
                "\u2193", TITLE_FONT, 20, QUIET, align=PP_ALIGN.CENTER)


def build_text_cards(prs, c):
    """Points in a 2x2 or 2x3 grid of cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Key Points"), size=24)
    items = c.get("items", [])
    n = min(len(items), 6)
    if n == 0: return
    cols = 2 if n <= 4 else 3
    rows_count = (n + cols - 1) // cols
    cardW = (8.4 - (cols - 1) * 0.25) / cols
    cardH = (3.8 - (rows_count - 1) * 0.2) / rows_count
    for i, item in enumerate(items[:n]):
        col_idx = i % cols; row = i // cols
        x = Inches(0.8 + col_idx * (cardW + 0.25))
        y = Emu(CONTENT_TOP + row * (Inches(cardH) + Inches(0.2)))
        col = ACCENTS[i % len(ACCENTS)]
        _line(slide, x, y, Inches(cardW), col, 2.5)
        title_text = item if isinstance(item, str) else item.get("title", "")
        detail_text = "" if isinstance(item, str) else item.get("detail", "")
        if detail_text:
            _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.12)), Inches(cardW - 0.1), Inches(0.35),
                title_text, BODY_FONT, 12, CHARCOAL, bold=True)
            _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.5)), Inches(cardW - 0.1), Emu(Inches(cardH) - Inches(0.6)),
                detail_text, BODY_FONT, 10, MID, line_spacing=14)
        else:
            _tb(slide, Emu(x + Inches(0.05)), Emu(y + Inches(0.12)), Inches(cardW - 0.1), Emu(Inches(cardH) - Inches(0.2)),
                title_text, BODY_FONT, 12, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=17)


def build_text_columns(prs, c):
    """2-3 columns of flowing text, magazine style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Overview"), size=24)
    columns = c.get("columns", [])
    n = min(len(columns), 3)
    if n == 0: return
    colW = (8.4 - (n - 1) * 0.3) / n
    for i, col_data in enumerate(columns[:n]):
        x = Inches(0.8 + i * (colW + 0.3))
        col = ACCENTS[i % len(ACCENTS)]
        heading = col_data if isinstance(col_data, str) else col_data.get("heading", "")
        body = "" if isinstance(col_data, str) else col_data.get("body", col_data if isinstance(col_data, str) else "")
        if isinstance(col_data, str):
            body = col_data; heading = ""
        if heading:
            _tb(slide, x, CONTENT_TOP, Inches(colW), Inches(0.35),
                heading, TITLE_FONT, 12, col, bold=True)
            _line(slide, x, Emu(CONTENT_TOP + Inches(0.38)), Inches(1.0), col, 1.5)
            _tb(slide, x, Emu(CONTENT_TOP + Inches(0.5)), Inches(colW), Inches(3.3),
                body, BODY_FONT, 11, CHARCOAL, line_spacing=16)
        else:
            _tb(slide, x, CONTENT_TOP, Inches(colW), Inches(3.8),
                body, BODY_FONT, 11, CHARCOAL, line_spacing=16)
        if i < n - 1:
            _vline(slide, Inches(0.8 + (i + 1) * colW + i * 0.3 + 0.15),
                   CONTENT_TOP, Inches(3.5), RULE_CLR, 0.75)




def build_text_narrative(prs, c):
    """Lede sentence pulled big, then body paragraphs below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Overview"), size=24)
    lede = c.get("lede", "")
    body = c.get("body", "")
    if isinstance(body, list): body = "\n\n".join(body)
    # Lede — big, bold, with gold left accent
    _rect(slide, LM, CONTENT_TOP, Inches(0.06), Inches(1.0), GOLD)
    _tb(slide, Emu(LM + Inches(0.2)), CONTENT_TOP, Inches(8.1), Inches(1.0),
        lede, TITLE_FONT, 16, CHARCOAL, bold=True, line_spacing=23,
        valign=MSO_ANCHOR.MIDDLE)
    # Separator
    _line(slide, LM, Inches(2.4), CW, RULE_CLR, 1)
    # Body text
    _tb(slide, LM, Inches(2.6), CW, Inches(2.5),
        body, BODY_FONT, 12, MID, line_spacing=18)


def build_text_nested(prs, c):
    """Section cards — top-level items as colored label blocks, children as flowing text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Detail"), size=24)
    items = c.get("items", [])
    n = min(len(items), 4)
    if n == 0: return
    avail = H - CONTENT_TOP - Inches(0.3)
    gap = Inches(0.15)
    cardH = (avail - (n - 1) * gap) / n

    for i, item in enumerate(items[:n]):
        y = Emu(CONTENT_TOP + i * (cardH + gap))
        text = item if isinstance(item, str) else item.get("text", "")
        children = [] if isinstance(item, str) else item.get("children", [])
        col = ACCENTS[i % len(ACCENTS)]
        # Colored label block
        _rect(slide, LM, y, Inches(2.2), Emu(cardH), col)
        _tb(slide, Emu(LM + Inches(0.12)), y, Inches(1.95), Emu(cardH),
            text, BODY_FONT, 11, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=15)
        # Children as flowing text
        child_texts = []
        for ch in children:
            ch_text = ch if isinstance(ch, str) else ch.get("text", "")
            grandchildren = [] if isinstance(ch, str) else ch.get("children", [])
            child_texts.append(ch_text)
            for gc in grandchildren:
                gc_text = gc if isinstance(gc, str) else gc.get("text", "")
                child_texts.append("  \u2022 " + gc_text)
        body = "\n".join(child_texts)
        _tb(slide, Inches(3.15), y, Inches(5.9), Emu(cardH),
            body, BODY_FONT, 10, CHARCOAL, valign=MSO_ANCHOR.MIDDLE,
            line_spacing=14)
def build_text_split(prs, c):
    """Left: big message. Right: supporting evidence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Key Point"), size=24)
    headline = c.get("headline", "")
    detail = c.get("detail", "")
    points = c.get("points", [])
    # Left: big message
    _rect(slide, LM, CONTENT_TOP, Inches(0.06), Inches(3.5), DK_GREEN)
    _tb(slide, Emu(LM + Inches(0.2)), CONTENT_TOP, Inches(4.0), Inches(2.0),
        headline, TITLE_FONT, 20, CHARCOAL, bold=True, line_spacing=28,
        valign=MSO_ANCHOR.MIDDLE)
    if detail:
        _tb(slide, Emu(LM + Inches(0.2)), Inches(3.2), Inches(4.0), Inches(1.2),
            detail, BODY_FONT, 11, MID, line_spacing=16)
    # Divider
    _vline(slide, Inches(5.1), CONTENT_TOP, Inches(3.5), RULE_CLR, 1)
    # Right: evidence
    for i, p in enumerate(points[:6]):
        py = Emu(CONTENT_TOP + Inches(0.1) + i * Inches(0.55))
        col = ACCENTS[i % len(ACCENTS)]
        _rect(slide, Inches(5.35), Emu(py + Inches(0.07)), Inches(0.06), Inches(0.06), col)
        _tb(slide, Inches(5.55), Emu(py), Inches(3.6), Inches(0.5),
            p, BODY_FONT, 11, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, line_spacing=15)


def build_text_annotated(prs, c):
    """Labeled paragraphs with colored margin tags."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paper_bg(slide)
    _content_title(slide, c.get("title", "Analysis"), size=24)
    items = c.get("items", [])
    n = min(len(items), 5)
    avail = H - CONTENT_TOP - Inches(0.3)
    rowH = min((avail - (n - 1) * Inches(0.12)) / n, Inches(0.85))
    gap = Inches(0.12)
    total = n * rowH + (n - 1) * gap
    startY = Emu(CONTENT_TOP + int((avail - total) * 0.2))

    for i, item in enumerate(items[:n]):
        y = Emu(startY + i * (rowH + gap))
        col = ACCENTS[i % len(ACCENTS)]
        # Label on the left
        _rect(slide, LM, y, Inches(1.5), Emu(rowH), col)
        _tb(slide, LM, y, Inches(1.5), Emu(rowH),
            item.get("label", ""), BODY_FONT, 10, WHITE, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Text on the right
        _tb(slide, Emu(LM + Inches(1.7)), y, Inches(6.7), Emu(rowH),
            item.get("text", ""), BODY_FONT, 11, CHARCOAL,
            valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


# ── Dispatch ─────────────────────────────────────────────
BUILDERS = {
    "title": build_title, "in_brief": build_in_brief,
    "section_divider": build_section_divider, "stat_callout": build_stat_callout,
    "quote": build_quote, "comparison": build_comparison,
    "text_graph": build_text_graph, "process_flow": build_process_flow,
    "matrix": build_matrix, "methods": build_methods,
    "hypotheses": build_hypotheses, "wsn_dense": build_wsn_dense,
    "wsn_reveal": build_wsn_reveal, "findings_recs": build_findings_recs,
    "findings_recs_dense": build_findings_recs_dense,
    "open_questions": build_open_questions, "agenda": build_agenda,
    "progressive_reveal": build_progressive_reveal, "closer": build_closer,
    "timeline": build_timeline, "data_table": build_data_table,
    "multi_stat": build_multi_stat, "persona": build_persona,
    "risk_tradeoff": build_risk_tradeoff, "appendix": build_appendix,
    "before_after": build_before_after, "summary": build_summary,
    "quote_full": build_quote_full, "stat_hero": build_stat_hero,
    "in_brief_featured": build_in_brief_featured, "persona_duo": build_persona_duo,
    "process_flow_vertical": build_process_flow_vertical,
    "text_cards": build_text_cards, "text_columns": build_text_columns,
    "text_narrative": build_text_narrative, "text_nested": build_text_nested,
    "text_split": build_text_split, "text_annotated": build_text_annotated,
}

def build_deck(slide_configs, output_path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    for slide_type, data in slide_configs:
        if slide_type == "skip": continue
        builder = BUILDERS.get(slide_type)
        if builder: builder(prs, data)
    prs.save(output_path); return output_path
