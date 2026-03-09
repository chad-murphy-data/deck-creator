# -*- coding: utf-8 -*-
"""
Colorful template — python-pptx port.
Signature: colored header bars, multi-color card system with section theming.
With Fidelity Slab / Fidelity Sans fonts and logo support.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

W = Inches(10)
H = Inches(5.625)

# Colors
GREEN = RGBColor(0x36, 0x87, 0x27)
GREEN_LIGHT = RGBColor(0xE2, 0xF0, 0xD9)
GREEN_MID = RGBColor(0x1D, 0xE4, 0xCA)
BLUE = RGBColor(0x38, 0x80, 0xF3)
BLUE_LIGHT = RGBColor(0xE0, 0xEA, 0xFC)
PURPLE = RGBColor(0x5B, 0x2C, 0x8F)
PURPLE_LIGHT = RGBColor(0xED, 0xE4, 0xF5)
ORANGE = RGBColor(0x04, 0x54, 0x7C)
ORANGE_LIGHT = RGBColor(0xE0, 0xEE, 0xF4)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
GOLD_LIGHT = RGBColor(0xF5, 0xEC, 0xD4)
DARK = RGBColor(0x40, 0x3F, 0x3E)
MID = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xE8, 0xE8, 0xE8)
OFF_WHITE = RGBColor(0xF9, 0xF7, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Fidelity Slab"
TITLE_FONT_FALLBACK = "Georgia"
BODY_FONT = "Fidelity Sans"
BODY_FONT_FALLBACK = "Calibri"

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "logo.png")

COLORS = [GREEN, BLUE, PURPLE, ORANGE, GOLD]
LIGHTS = [GREEN_LIGHT, BLUE_LIGHT, PURPLE_LIGHT, ORANGE_LIGHT, GOLD_LIGHT]


def _add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_oval(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _font_fallback(font_name):
    """Return (primary, fallback) for a font name."""
    if font_name == TITLE_FONT:
        return TITLE_FONT, TITLE_FONT_FALLBACK
    if font_name == BODY_FONT:
        return BODY_FONT, BODY_FONT_FALLBACK
    return font_name, None


def _set_font_with_fallback(font_obj, font_name):
    """Set font name and add XML-level fallback via cs element."""
    primary, fallback = _font_fallback(font_name)
    font_obj.name = primary
    if fallback:
        el = font_obj._element  # rPr or defRPr element
        latin = el.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(el, qn('a:latin'))
            latin.set('typeface', primary)
        cs = el.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(el, qn('a:cs'))
        cs.set('typeface', fallback)


def _add_text_box(slide, x, y, w, h, text, font_name=BODY_FONT, font_size=12,
                  color=DARK, bold=False, italic=False, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    _set_font_with_fallback(p.font, font_name)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align

    if line_spacing:
        p.line_spacing = Pt(line_spacing)

    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}
        bodyPr.set('anchor', anchor_map.get(valign, 't'))

    return txBox


def _add_rounded_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    pg = shape._element.find('.//' + qn('a:prstGeom'))
    if pg is not None:
        av = pg.find(qn('a:avLst'))
        if av is None: av = etree.SubElement(pg, qn('a:avLst'))
        for g in av.findall(qn('a:gd')): av.remove(g)
        g = etree.SubElement(av, qn('a:gd')); g.set('name','adj'); g.set('fmla','val 5000')
    return shape


def _header(slide, title, color=GREEN):
    _add_rect(slide, Inches(0), Inches(0), W, Inches(1.0), color)
    _add_text_box(slide, Inches(0.5), Inches(0.15), Inches(8), Inches(0.7),
                  title, TITLE_FONT, 24, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

CONTENT_TOP = Inches(1.15)
_LM = Inches(0.75)
_CW = Inches(8.5)


def _draw_stepper_bar(slide, labels, active_count, y):
    """Reusable stepper bar: numbered circles + connecting lines + labels."""
    n = len(labels)
    cs = Inches(0.42)
    total_w = _CW
    spacing = int(total_w / max(n - 1, 1)) if n > 1 else 0
    start_x = _LM
    for i in range(n):
        cx = Emu(start_x + i * spacing) if n > 1 else Emu(start_x + total_w // 2 - cs // 2)
        active = i < active_count
        color = COLORS[i % len(COLORS)] if active else LIGHT
        _add_oval(slide, cx, y, cs, cs, color)
        tc = WHITE if active else MID
        _add_text_box(slide, cx, y, cs, cs, str(i + 1), TITLE_FONT, 15, tc,
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        lw = Inches(1.4)
        label_x = Emu(cx - (lw - cs) // 2)
        _add_text_box(slide, label_x, Emu(y + cs + Inches(0.04)),
                      lw, Inches(0.25), labels[i],
                      BODY_FONT, 9, color if active else MID,
                      bold=active, align=PP_ALIGN.CENTER)
        if i < n - 1:
            line_x1 = Emu(cx + cs + Inches(0.08))
            next_cx = Emu(start_x + (i + 1) * spacing)
            line_x2 = Emu(next_cx - Inches(0.08))
            _add_rect(slide, line_x1, Emu(y + cs // 2), Emu(line_x2 - line_x1), Inches(0.025), LIGHT)


# ============================================================
# BUILDERS
# ============================================================

def build_title(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    _add_text_box(slide, Inches(0.5), Inches(0.8), Inches(6.5), Inches(2.0),
                  c.get("title", "Title"), TITLE_FONT, 36, GREEN, bold=True,
                  valign=MSO_ANCHOR.BOTTOM)
    # Multi-color bar
    uw = Inches(1.5)
    bar_colors = [GREEN, BLUE, PURPLE, ORANGE]
    for i, col in enumerate(bar_colors):
        _add_rect(slide, Emu(Inches(0.5) + i * uw), Inches(2.92),
                  Emu(uw - Inches(0.08)), Inches(0.05), col)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.5), Inches(3.1), Inches(6.5), Inches(0.5),
                      c["subtitle"], BODY_FONT, 16, MID)
    meta = []
    if c.get("author"):
        meta.append(c["author"])
    if c.get("date"):
        meta.append(c["date"])
    if meta:
        _add_text_box(slide, Inches(0.5), Inches(4.2), Inches(5), Inches(0.7),
                      "\n".join(meta), BODY_FONT, 12, MID)
    # Logo in upper-right
    logo = c.get("logo_path", LOGO_PATH)
    if logo and os.path.isfile(logo):
        slide.shapes.add_picture(logo, Inches(7.6), Inches(0.3), Inches(2.0))
    img = c.get("imagePath", "")
    if img and os.path.exists(img):
        slide.shapes.add_picture(img, Inches(5.4), Inches(0.3), Inches(4.3), Inches(5.0))


def build_in_brief(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "In Brief"))
    bullets = c.get("bullets", [])
    startY = Inches(1.2)
    rowH = Inches(0.95)
    gap = Inches(0.08)

    for i, b in enumerate(bullets):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, Inches(0.5), y, Inches(9), rowH, lt)
        _add_rect(slide, Inches(0.5), y, Inches(0.12), rowH, col)
        _add_oval(slide, Inches(0.8), Emu(y + (rowH - Inches(0.42)) // 2),
                  Inches(0.42), Inches(0.42), col)
        _add_text_box(slide, Inches(0.8), Emu(y + (rowH - Inches(0.42)) // 2),
                      Inches(0.42), Inches(0.42), str(i + 1),
                      TITLE_FONT, 13, WHITE, bold=True, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Inches(1.5), Emu(y + Inches(0.05)),
                      Inches(7.8), Emu(rowH - Inches(0.1)),
                      b, BODY_FONT, 13, DARK, valign=MSO_ANCHOR.MIDDLE)


def build_section_divider(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    if c.get("sectionNumber"):
        _add_text_box(slide, Inches(0.6), Inches(1.2), Inches(2), Inches(0.8),
                      f"0{c['sectionNumber']}", TITLE_FONT, 48, GREEN_MID, bold=True,
                      valign=MSO_ANCHOR.BOTTOM)
    _add_text_box(slide, Inches(0.6), Inches(2.0), Inches(8), Inches(1.2),
                  c.get("title", "Section"), TITLE_FONT, 36, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.6), Inches(3.3), Inches(8), Inches(0.6),
                      c["subtitle"], BODY_FONT, 14, WHITE)


def build_stat_callout(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Key Metric"))
    # Visual anchor circle behind stat
    sc = GREEN
    cs = Inches(2.2)
    cx = Emu((W - cs) // 2)
    cy = Emu(Inches(1.5))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, cs, cs)
    oval_fill = oval.fill
    oval_fill.solid()
    r, g, b = sc[0], sc[1], sc[2]
    light_r = min(255, r + int((255 - r) * 0.85))
    light_g = min(255, g + int((255 - g) * 0.85))
    light_b = min(255, b + int((255 - b) * 0.85))
    oval_fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
    oval.line.fill.background()
    _add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.8),
                  c.get("stat", "—"), TITLE_FONT, 96, GREEN, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if c.get("headline"):
        _add_text_box(slide, Inches(1.0), Inches(3.1), Inches(8), Inches(0.7),
                      c["headline"], BODY_FONT, 18, DARK, bold=True, align=PP_ALIGN.CENTER)
    if c.get("detail"):
        _add_text_box(slide, Inches(1.5), Inches(3.8), Inches(7), Inches(0.8),
                      c["detail"], BODY_FONT, 12, MID, align=PP_ALIGN.CENTER)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True, align=PP_ALIGN.CENTER)


def build_quote(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "In Their Words"))
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(1), Inches(1),
                  "\u201C", "Georgia", 72, GREEN_LIGHT, bold=True)
    _add_text_box(slide, Inches(1.5), Inches(1.6), Inches(7.0), Inches(2.0),
                  c.get("quote", ""), BODY_FONT, 18, DARK, italic=True,
                  valign=MSO_ANCHOR.MIDDLE, line_spacing=25)
    _add_rect(slide, Inches(1.5), Inches(3.8), Inches(1.5), Inches(0.05), GREEN)
    if c.get("attribution"):
        _add_text_box(slide, Inches(1.5), Inches(3.95), Inches(7), Inches(0.5),
                      c["attribution"], BODY_FONT, 12, MID)
    if c.get("context"):
        _add_text_box(slide, Inches(1.5), Inches(4.3), Inches(7), Inches(0.4),
                      c["context"], BODY_FONT, 10, MID, italic=True)


def build_comparison(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Comparison"))

    left_color = ORANGE
    right_color = GREEN

    left_label = c.get("leftLabel", "Before")
    right_label = c.get("rightLabel", "After")
    left_items = c.get("leftItems", [])
    right_items = c.get("rightItems", [])

    content_top = Inches(1.2)
    card_w = Inches(4.2)
    card_h = Inches(3.8)
    gap = Inches(0.4)
    header_h = Inches(0.5)

    # Left card
    left_x = Inches(0.5)
    _add_rect(slide, left_x, content_top, card_w, card_h, WHITE)
    _add_rect(slide, left_x, content_top, card_w, header_h, left_color)
    _add_text_box(slide, Emu(left_x + Inches(0.2)), content_top,
                  Emu(card_w - Inches(0.4)), header_h,
                  left_label, BODY_FONT, 14, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(left_items[:6]):
        iy = Emu(content_top + header_h + Inches(0.15) + i * Inches(0.5))
        text = item if isinstance(item, str) else str(item)
        _add_text_box(slide, Emu(left_x + Inches(0.25)), iy,
                      Emu(card_w - Inches(0.5)), Inches(0.45),
                      text, BODY_FONT, 13, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)

    # Right card
    right_x = Emu(left_x + card_w + gap)
    _add_rect(slide, right_x, content_top, card_w, card_h, WHITE)
    _add_rect(slide, right_x, content_top, card_w, header_h, right_color)
    _add_text_box(slide, Emu(right_x + Inches(0.2)), content_top,
                  Emu(card_w - Inches(0.4)), header_h,
                  right_label, BODY_FONT, 14, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(right_items[:6]):
        iy = Emu(content_top + header_h + Inches(0.15) + i * Inches(0.5))
        text = item if isinstance(item, str) else str(item)
        _add_text_box(slide, Emu(right_x + Inches(0.25)), iy,
                      Emu(card_w - Inches(0.5)), Inches(0.45),
                      text, BODY_FONT, 13, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)


def build_text_graph(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Title"))
    texts = c.get("text", [])
    if not isinstance(texts, list):
        texts = [texts]
    for i, t in enumerate(texts):
        _add_text_box(slide, Inches(0.5), Emu(Inches(1.2) + i * Inches(1.15)),
                      Inches(4.2), Inches(1.1), t, BODY_FONT, 12, DARK)
    _add_rect(slide, Inches(4.9), Inches(1.2), Inches(0.06), Inches(3.6), GREEN)

    from pptx.chart.data import CategoryChartData
    chart_data_raw = c.get("chartData", [{"name": "S1", "labels": ["A", "B", "C"], "values": [25, 45, 30]}])
    chart_data = CategoryChartData()
    if chart_data_raw:
        cd = chart_data_raw[0]
        chart_data.categories = cd.get("labels", ["A", "B", "C"])
        chart_data.add_series(cd.get("name", "Series 1"), cd.get("values", [25, 45, 30]))

    ct = XL_CHART_TYPE.COLUMN_CLUSTERED
    if c.get("chartType") == "line":
        ct = XL_CHART_TYPE.LINE
    elif c.get("chartType") == "pie":
        ct = XL_CHART_TYPE.PIE

    slide.shapes.add_chart(ct, Inches(5.15), Inches(1.0), Inches(4.5), Inches(4.0), chart_data)


def build_process_flow(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Process"))
    steps = c.get("steps", [])
    n = min(len(steps), 5)
    if n == 0: return

    labels = [s.get("title", f"Step {i+1}") for i, s in enumerate(steps[:n])]
    _draw_stepper_bar(slide, labels, n, Emu(CONTENT_TOP + Inches(0.05)))

    card_top = Emu(CONTENT_TOP + Inches(0.85))
    gap = Inches(0.2)

    if n <= 3:
        card_w = int((_CW - (n - 1) * gap) / n)
        card_h = Inches(3.0)
        for i, step in enumerate(steps[:n]):
            cx = Emu(_LM + i * (card_w + gap))
            _add_rounded_rect(slide, cx, card_top, card_w, card_h, OFF_WHITE)
            _add_rect(slide, cx, card_top, Inches(0.06), card_h, COLORS[i % len(COLORS)])
            _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(card_top + Inches(0.15)),
                          Emu(card_w - Inches(0.4)), Inches(0.5),
                          step.get("title", ""), BODY_FONT, 13, DARK, bold=True)
            if step.get("detail"):
                _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(card_top + Inches(0.65)),
                              Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.85)),
                              step["detail"], BODY_FONT, 10, MID, line_spacing=14)
    elif n == 4:
        card_w = int((_CW - gap) / 2)
        card_h = Inches(1.4)
        for i, step in enumerate(steps[:4]):
            col = i % 2
            row = i // 2
            cx = Emu(_LM + col * (card_w + gap))
            cy = Emu(card_top + row * (card_h + gap))
            _add_rounded_rect(slide, cx, cy, card_w, card_h, OFF_WHITE)
            _add_rect(slide, cx, cy, Inches(0.06), card_h, COLORS[i % len(COLORS)])
            _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.1)),
                          Emu(card_w - Inches(0.4)), Inches(0.4),
                          step.get("title", ""), BODY_FONT, 12, DARK, bold=True)
            if step.get("detail"):
                _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.5)),
                              Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.65)),
                              step["detail"], BODY_FONT, 10, MID, line_spacing=14)
    else:
        card_w = int((_CW - 2 * gap) / 3)
        card_h = Inches(1.35)
        for i, step in enumerate(steps[:5]):
            if i < 3:
                cx = Emu(_LM + i * (card_w + gap))
                cy = card_top
            else:
                offset = (_CW - 2 * card_w - gap) // 2
                cx = Emu(_LM + offset + (i - 3) * (card_w + gap))
                cy = Emu(card_top + card_h + gap)
            _add_rounded_rect(slide, cx, cy, card_w, card_h, OFF_WHITE)
            _add_rect(slide, cx, cy, Inches(0.06), card_h, COLORS[i % len(COLORS)])
            _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.1)),
                          Emu(card_w - Inches(0.4)), Inches(0.35),
                          step.get("title", ""), BODY_FONT, 12, DARK, bold=True)
            if step.get("detail"):
                _add_text_box(slide, Emu(cx + Inches(0.2)), Emu(cy + Inches(0.45)),
                              Emu(card_w - Inches(0.4)), Emu(card_h - Inches(0.55)),
                              step["detail"], BODY_FONT, 10, MID, line_spacing=14)


def build_process_flow_reveal(prs, c):
    """Generates N slides, one per step, with progressive stepper."""
    steps = c.get("steps", [])
    n = min(len(steps), 5)
    if n == 0: return

    labels = [s.get("title", f"Step {i+1}") for i, s in enumerate(steps[:n])]

    for step_idx in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _header(slide, c.get("title", "Process"))

        _draw_stepper_bar(slide, labels, step_idx + 1, Emu(CONTENT_TOP + Inches(0.05)))

        card_y = Emu(CONTENT_TOP + Inches(0.85))
        card_w = Inches(6.0)
        card_h = Inches(3.0)
        card_x = Emu(_LM + (_CW - card_w) // 2)

        step = steps[step_idx]
        _add_rounded_rect(slide, card_x, card_y, card_w, card_h, OFF_WHITE)
        _add_rect(slide, card_x, card_y, Inches(0.08), card_h, COLORS[step_idx % len(COLORS)])
        _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.2)),
                      Emu(card_w - Inches(0.5)), Inches(0.6),
                      step.get("title", ""), BODY_FONT, 16, DARK, bold=True)
        if step.get("detail"):
            _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.85)),
                          Emu(card_w - Inches(0.5)), Emu(card_h - Inches(1.1)),
                          step["detail"], BODY_FONT, 12, MID, line_spacing=16)


def build_matrix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Framework"))
    quads = c.get("quadrants", [{}, {}, {}, {}])
    qW = Inches(4.15)
    qH = Inches(1.65)
    gap = Inches(0.2)
    sX = Inches(0.85)
    sY = Inches(1.5)

    for i, q in enumerate(quads[:4]):
        col_idx = i % 2
        row = i // 2
        x = Emu(sX + col_idx * (qW + gap))
        y = Emu(sY + row * (qH + gap))
        col = COLORS[i]
        lt = LIGHTS[i]

        _add_rect(slide, x, y, qW, qH, lt)
        _add_rect(slide, x, y, qW, Inches(0.08), col)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.12)),
                      Emu(qW - Inches(0.3)), Inches(0.3),
                      q.get("label", ""), TITLE_FONT, 12, col, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.45)),
                      Emu(qW - Inches(0.3)), Emu(qH - Inches(0.6)),
                      q.get("detail", ""), BODY_FONT, 10, DARK)


def build_methods(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Approach"))
    for i, f in enumerate(c.get("fields", [])):
        y = Emu(Inches(1.3) + i * Inches(0.8))
        col = COLORS[i % len(COLORS)]
        _add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(0.7), col)
        _add_text_box(slide, Inches(0.75), y, Inches(2.0), Inches(0.7),
                      f.get("label", ""), BODY_FONT, 12, col, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Inches(2.8), y, Inches(6.7), Inches(0.7),
                      f.get("value", ""), BODY_FONT, 12, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_hypotheses(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Hypotheses"))
    for i, h in enumerate(c.get("hypotheses", [])):
        y = Emu(Inches(1.2) + i * Inches(0.75))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, Inches(0.5), y, Inches(8.2), Inches(0.65), lt)
        _add_rect(slide, Inches(0.5), y, Inches(0.1), Inches(0.65), col)
        _add_text_box(slide, Inches(0.75), y, Inches(0.5), Inches(0.65),
                      f"H{i + 1}", TITLE_FONT, 14, col, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Inches(1.3), y, Inches(6.2), Inches(0.65),
                      h.get("text", ""), BODY_FONT, 12, DARK,
                      valign=MSO_ANCHOR.MIDDLE)
        if h.get("status"):
            sc = GREEN if h["status"] == "Confirmed" else (ORANGE if h["status"] == "Rejected" else MID)
            _add_rect(slide, Inches(7.8), Emu(y + (Inches(0.65) - Inches(0.3)) // 2),
                      Inches(0.9), Inches(0.3), sc)
            _add_text_box(slide, Inches(7.8), Emu(y + (Inches(0.65) - Inches(0.3)) // 2),
                          Inches(0.9), Inches(0.3),
                          h["status"], BODY_FONT, 8, WHITE, bold=True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def build_wsn_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Key Finding"))
    labels = c.get("labels", ["What", "So What", "Now What"])
    cols = [
        (labels[0], GREEN, GREEN_LIGHT, c.get("what", {})),
        (labels[1], BLUE, BLUE_LIGHT, c.get("soWhat", {})),
        (labels[2], PURPLE, PURPLE_LIGHT, c.get("nowWhat", {})),
    ]
    colW = Inches(2.85)
    gap = Inches(0.2)
    startX = Inches(0.5)
    startY = Inches(1.2)
    cardH = Inches(3.8)

    for i, (label, color, light, data) in enumerate(cols):
        x = Emu(startX + i * (colW + gap))
        _add_rect(slide, x, startY, colW, cardH, light)
        _add_rect(slide, x, startY, colW, Inches(0.1), color)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(startY + Inches(0.15)),
                      Emu(colW - Inches(0.3)), Inches(0.35),
                      label, TITLE_FONT, 14, color, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(startY + Inches(0.55)),
                      Emu(colW - Inches(0.3)), Inches(0.8),
                      data.get("headline", ""), BODY_FONT, 12, DARK, bold=True)
        if data.get("detail"):
            _add_text_box(slide, Emu(x + Inches(0.15)), Emu(startY + Inches(1.4)),
                          Emu(colW - Inches(0.3)), Inches(1.7),
                          data["detail"], BODY_FONT, 10, MID)


def build_wsn_reveal(prs, c):
    """3 progressive slides: What -> So What -> Now What with running summary."""
    labels = c.get("labels", ["What", "So What", "Now What"])
    WSN_KEYS = ["what", "soWhat", "nowWhat"]
    WSN_COLORS = [GREEN, BLUE, PURPLE]
    WSN_LIGHTS = [GREEN_LIGHT, BLUE_LIGHT, PURPLE_LIGHT]

    for step in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_rect(slide, Inches(0), Inches(0), W, Inches(0.15), GREEN)
        _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
                      c.get("title", "Key Finding"), TITLE_FONT, 28, DARK, bold=True)

        # Stepper bar at top
        _draw_stepper_bar(slide, labels, step + 1, Emu(CONTENT_TOP + Inches(0.05)))

        # Featured card for current step
        data = c.get(WSN_KEYS[step], {})
        card_y = Emu(CONTENT_TOP + Inches(0.85))
        card_w = Inches(7.0)
        card_h = Inches(2.4)
        card_x = Emu(_LM + (_CW - card_w) // 2)

        _add_rounded_rect(slide, card_x, card_y, card_w, card_h, WSN_LIGHTS[step])
        _add_rect(slide, card_x, card_y, Inches(0.10), card_h, WSN_COLORS[step])
        _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.15)),
                      Emu(card_w - Inches(0.5)), Inches(0.6),
                      data.get("headline", ""), BODY_FONT, 15, DARK, bold=True,
                      line_spacing=18)
        if data.get("detail"):
            _add_text_box(slide, Emu(card_x + Inches(0.25)), Emu(card_y + Inches(0.8)),
                          Emu(card_w - Inches(0.5)), Emu(card_h - Inches(1.0)),
                          data["detail"], BODY_FONT, 12, MID, line_spacing=16)

        # Running summary (only for step >= 1)
        if step > 0:
            summary_y = Emu(card_y + card_h + Inches(0.15))
            # Thin divider
            _add_rect(slide, _LM, summary_y, Inches(3.0), Inches(0.02), LIGHT)
            # Previous step summaries
            for j in range(step):
                iy = Emu(summary_y + Inches(0.1) + j * Inches(0.28))
                prev_data = c.get(WSN_KEYS[j], {})
                summary_text = prev_data.get("summary", prev_data.get("headline", ""))
                # Small bullet square
                _add_rect(slide, _LM, Emu(iy + Inches(0.05)), Inches(0.08), Inches(0.08), LIGHT)
                _add_text_box(slide, Emu(_LM + Inches(0.18)), iy,
                              Emu(_CW - Inches(0.18)), Inches(0.25),
                              summary_text, BODY_FONT, 9, MID,
                              valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Findings & Recommendations"))
    items = c.get("items", [])
    startY = Inches(1.15)
    rowH = Inches(0.82)
    gap = Inches(0.08)

    for i, item in enumerate(items[:5]):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]

        _add_rect(slide, Inches(0.5), y, Inches(3.9), rowH, lt)
        _add_rect(slide, Inches(0.5), y, Inches(0.1), rowH, col)
        _add_text_box(slide, Inches(0.75), y, Inches(3.5), rowH,
                      item.get("finding", ""), BODY_FONT, 10.5, DARK,
                      valign=MSO_ANCHOR.MIDDLE)

        _add_oval(slide, Inches(4.62), Emu(y + (rowH - Inches(0.42)) // 2),
                  Inches(0.42), Inches(0.42), col)
        _add_text_box(slide, Inches(4.62), Emu(y + (rowH - Inches(0.42)) // 2),
                      Inches(0.42), Inches(0.42), "\u2192",
                      BODY_FONT, 14, WHITE, bold=True, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)

        _add_rect(slide, Inches(5.2), y, Inches(4.3), rowH, lt)
        _add_rect(slide, Inches(5.2), y, Inches(0.1), rowH, BLUE)
        _add_text_box(slide, Inches(5.45), y, Inches(3.9), rowH,
                      item.get("recommendation", ""), BODY_FONT, 10.5, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_findings_recs_dense(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Complete Findings"))
    items = c.get("items", [])
    startY = Inches(1.15)
    rowH = Inches(0.56)
    gap = Inches(0.08)

    for i, item in enumerate(items[:8]):
        y = Emu(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        bg = OFF_WHITE if i % 2 == 0 else WHITE

        _add_rect(slide, Inches(0.5), y, Inches(4.2), rowH, bg)
        _add_oval(slide, Inches(0.2), Emu(y + (rowH - Inches(0.3)) // 2),
                  Inches(0.3), Inches(0.3), col)
        _add_text_box(slide, Inches(0.2), Emu(y + (rowH - Inches(0.3)) // 2),
                      Inches(0.3), Inches(0.3), str(i + 1),
                      TITLE_FONT, 9, WHITE, bold=True, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Inches(0.55), y, Inches(4.1), rowH,
                      item.get("finding", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE)

        _add_rect(slide, Inches(4.85), y, Inches(4.65), rowH, bg)
        _add_rect(slide, Inches(4.85), y, Inches(0.06), rowH, BLUE)
        _add_text_box(slide, Inches(5.0), y, Inches(4.4), rowH,
                      item.get("recommendation", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE)


def build_open_questions(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Open Questions"))
    questions = c.get("questions", [])
    cardW = Inches(4.2)
    cardH = Inches(1.7)
    gX = Inches(0.3)
    gY = Inches(0.2)
    gridX = Inches(0.5)
    gridY = Inches(1.2)

    for i, question in enumerate(questions[:4]):
        col_idx = i % 2
        row = i // 2
        x = Emu(gridX + col_idx * (cardW + gX))
        y = Emu(gridY + row * (cardH + gY))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]

        _add_rect(slide, x, y, cardW, cardH, lt)
        _add_rect(slide, x, y, cardW, Inches(0.08), col)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.15)),
                      Inches(0.5), Inches(0.5), str(i + 1),
                      TITLE_FONT, 26, col, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.15)), Emu(y + Inches(0.7)),
                      Emu(cardW - Inches(0.3)), Inches(0.85),
                      question, BODY_FONT, 12, DARK)


def build_agenda(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Agenda"))
    items = c.get("items", [])
    for i, item in enumerate(items):
        y = Emu(Inches(1.2) + i * Inches(0.75))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]

        _add_rect(slide, Inches(0.5), y, Inches(9), Inches(0.65), lt)
        _add_rect(slide, Inches(0.5), y, Inches(0.1), Inches(0.65), col)
        _add_oval(slide, Inches(0.8), Emu(y + (Inches(0.65) - Inches(0.4)) // 2),
                  Inches(0.4), Inches(0.4), col)
        _add_text_box(slide, Inches(0.8), Emu(y + (Inches(0.65) - Inches(0.4)) // 2),
                      Inches(0.4), Inches(0.4), str(i + 1),
                      TITLE_FONT, 13, WHITE, bold=True, align=PP_ALIGN.CENTER,
                      valign=MSO_ANCHOR.MIDDLE)

        title_text = item if isinstance(item, str) else item.get("title", "")
        _add_text_box(slide, Inches(1.45), y, Inches(5.5), Inches(0.65),
                      title_text, BODY_FONT, 14, DARK, bold=True,
                      valign=MSO_ANCHOR.MIDDLE)
        if isinstance(item, dict) and item.get("detail"):
            _add_text_box(slide, Inches(7.0), y, Inches(2.3), Inches(0.65),
                          item["detail"], BODY_FONT, 10, MID,
                          align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def build_progressive_reveal(prs, c):
    takeaways = c.get("takeaways", [])
    for n in range(min(len(takeaways), 5)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_rect(slide, Inches(0), Inches(0), W, Inches(0.1), GREEN)
        _add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.6),
                      c.get("title", "Building the Picture"), TITLE_FONT, 24, DARK, bold=True)

        cur = takeaways[n]
        col = COLORS[n % len(COLORS)]
        lt = LIGHTS[n % len(LIGHTS)]

        _add_rect(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(2.4), lt)
        _add_rect(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.1), col)
        _add_text_box(slide, Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.7),
                      cur.get("headline", ""), BODY_FONT, 16, DARK, bold=True)
        if cur.get("detail"):
            _add_text_box(slide, Inches(0.7), Inches(2.05), Inches(8.6), Inches(1.2),
                          cur["detail"], BODY_FONT, 11, MID)

        # Thin subtle divider
        _add_rect(slide, _LM, Inches(3.7), Inches(3.0), Inches(0.02), LIGHT)
        # Running summary items (no header, all gray)
        for j in range(n + 1):
            ty = Emu(Inches(3.8) + j * Inches(0.28))
            _add_rect(slide, _LM, Emu(ty + Inches(0.04)), Inches(0.08), Inches(0.08), LIGHT)
            _add_text_box(slide, Emu(_LM + Inches(0.18)), ty,
                          Emu(_CW - Inches(0.18)), Inches(0.25),
                          takeaways[j].get("summary", takeaways[j].get("headline", "")),
                          BODY_FONT, 9, MID,
                          valign=MSO_ANCHOR.MIDDLE)


def build_closer(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    _add_text_box(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(1.2),
                  c.get("title", "Thank You"), TITLE_FONT, 44, WHITE, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
    _add_rect(slide, Inches(3.75), Inches(2.75), Inches(2.5), Inches(0.05), GREEN_MID)
    if c.get("subtitle"):
        _add_text_box(slide, Inches(0.5), Inches(2.95), Inches(9), Inches(0.5),
                      c["subtitle"], BODY_FONT, 16, WHITE, align=PP_ALIGN.CENTER)
    if c.get("contact"):
        _add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
                      c["contact"], BODY_FONT, 12, GREEN_LIGHT, align=PP_ALIGN.CENTER)


def build_timeline(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Timeline"))
    milestones = c.get("milestones", [])
    n = min(len(milestones), 6)
    if n == 0:
        return
    lineY = Inches(2.8)
    startX = 0.75
    endX = 9.25
    step = (endX - startX) / max(n - 1, 1)
    # Horizontal line
    _add_rect(slide, Inches(startX), lineY, Inches(endX - startX), Inches(0.04), LIGHT)
    STATUS_CLR = {"complete": GREEN, "current": GOLD, "upcoming": MID}
    for i, m in enumerate(milestones[:n]):
        cx = Inches(startX + i * step)
        col = STATUS_CLR.get(m.get("status", "upcoming"), MID)
        # Dot
        ds = Inches(0.22)
        _add_oval(slide, Emu(cx - ds // 2), Emu(lineY - ds // 2), ds, ds, col)
        # Date above
        _add_text_box(slide, Emu(cx - Inches(0.6)), Emu(lineY - Inches(0.7)),
                      Inches(1.2), Inches(0.35), m.get("date", ""),
                      BODY_FONT, 10, col, bold=True, align=PP_ALIGN.CENTER)
        # Title below
        _add_text_box(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.2)),
                      Inches(1.4), Inches(0.4), m.get("title", ""),
                      BODY_FONT, 11, DARK, bold=True, align=PP_ALIGN.CENTER)
        if m.get("detail"):
            _add_text_box(slide, Emu(cx - Inches(0.7)), Emu(lineY + Inches(0.6)),
                          Inches(1.4), Inches(0.7), m["detail"],
                          BODY_FONT, 9, MID, align=PP_ALIGN.CENTER, line_spacing=12)


def build_data_table(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Data"))
    headers = c.get("headers", [])
    rows = c.get("rows", [])
    hc = c.get("highlightCol", None)
    nCols = len(headers)
    nRows = min(len(rows), 10)
    if nCols == 0:
        return
    tW = 9.0
    tX = 0.5
    colW = tW / nCols
    rowH = min(0.42, 3.6 / max(nRows + 1, 1))
    hdrY = 1.2
    # Header row
    _add_rect(slide, Inches(tX), Inches(hdrY), Inches(tW), Inches(rowH), GREEN)
    for j, h in enumerate(headers):
        _add_text_box(slide, Inches(tX + j * colW + 0.08), Inches(hdrY),
                      Inches(colW - 0.16), Inches(rowH), h,
                      BODY_FONT, 10, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    # Data rows
    for i, row in enumerate(rows[:nRows]):
        ry = hdrY + (i + 1) * rowH
        bg = OFF_WHITE if i % 2 == 0 else WHITE
        _add_rect(slide, Inches(tX), Inches(ry), Inches(tW), Inches(rowH), bg)
        for j, cell in enumerate(row[:nCols]):
            is_hl = hc is not None and j == hc
            _add_text_box(slide, Inches(tX + j * colW + 0.08), Inches(ry),
                          Inches(colW - 0.16), Inches(rowH), str(cell),
                          BODY_FONT, 10, GREEN if is_hl else DARK,
                          bold=is_hl, valign=MSO_ANCHOR.MIDDLE)
    if c.get("note"):
        _add_text_box(slide, Inches(tX), Inches(hdrY + (nRows + 1) * rowH + 0.1),
                      Inches(tW), Inches(0.3), c["note"],
                      BODY_FONT, 8, MID, italic=True)


def build_multi_stat(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Key Metrics"))
    stats = c.get("stats", [])
    n = min(len(stats), 4)
    if n == 0:
        return
    totalW = 9.0
    gap = 0.25
    statW = (totalW - (n - 1) * gap) / n
    for i, s in enumerate(stats[:n]):
        x = Inches(0.5 + i * (statW + gap))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        cardH = Inches(3.6)
        _add_rect(slide, x, Inches(1.2), Inches(statW), cardH, lt)
        _add_rect(slide, x, Inches(1.2), Inches(statW), Inches(0.1), col)
        _add_text_box(slide, x, Inches(1.5), Inches(statW), Inches(1.4),
                      s.get("value", "\u2014"), TITLE_FONT, 48, col, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, x, Inches(2.9), Inches(statW), Inches(0.45),
                      s.get("label", ""), BODY_FONT, 13, DARK, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
        if s.get("detail"):
            _add_text_box(slide, x, Inches(3.4), Inches(statW), Inches(0.8),
                          s["detail"], BODY_FONT, 10, MID,
                          align=PP_ALIGN.CENTER, line_spacing=14)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True,
                      align=PP_ALIGN.CENTER)


def build_persona(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Persona"))
    # Left panel: name, archetype, traits
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(4.2), Inches(3.8), GREEN_LIGHT)
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.1), GREEN)
    _add_text_box(slide, Inches(0.7), Inches(1.4), Inches(3.8), Inches(0.5),
                  c.get("name", ""), TITLE_FONT, 22, GREEN, bold=True)
    _add_text_box(slide, Inches(0.7), Inches(1.9), Inches(3.8), Inches(0.35),
                  c.get("archetype", ""), BODY_FONT, 13, GOLD, bold=True)
    # Traits
    traits = c.get("traits", [])
    for i, t in enumerate(traits[:5]):
        ty = Emu(Inches(2.4) + i * Inches(0.42))
        col = COLORS[i % len(COLORS)]
        _add_oval(slide, Inches(0.75), Emu(ty + Inches(0.06)), Inches(0.14),
                  Inches(0.14), col)
        _add_text_box(slide, Inches(1.0), ty, Inches(3.5), Inches(0.35),
                      t, BODY_FONT, 11, DARK)
    # Right panel: strategy
    _add_rect(slide, Inches(4.9), Inches(1.2), Inches(4.6), Inches(3.8), BLUE_LIGHT)
    _add_rect(slide, Inches(4.9), Inches(1.2), Inches(4.6), Inches(0.1), BLUE)
    _add_text_box(slide, Inches(5.1), Inches(1.4), Inches(4.2), Inches(0.3),
                  "STRATEGY", BODY_FONT, 10, BLUE, bold=True)
    _add_text_box(slide, Inches(5.1), Inches(1.8), Inches(4.2), Inches(1.5),
                  c.get("strategy", ""), BODY_FONT, 12, DARK, line_spacing=17)
    if c.get("detail"):
        _add_text_box(slide, Inches(5.1), Inches(3.4), Inches(4.2), Inches(1.4),
                      c["detail"], BODY_FONT, 10, MID, line_spacing=14)


def build_risk_tradeoff(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Risk & Reward"))
    risks = c.get("risks", [])
    rewards = c.get("rewards", [])
    SEV = {"high": RGBColor(0xC2, 0x3B, 0x22), "medium": GOLD, "low": GREEN}
    # Left: Risks
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.45),
              RGBColor(0xC2, 0x3B, 0x22))
    _add_text_box(slide, Inches(0.7), Inches(1.2), Inches(3.8), Inches(0.45),
                  "RISKS", BODY_FONT, 12, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(risks[:5]):
        ry = Emu(Inches(1.75) + i * Inches(0.7))
        sc = SEV.get(r.get("severity", "medium"), GOLD)
        _add_rect(slide, Inches(0.5), ry, Inches(4.2), Inches(0.6), OFF_WHITE)
        _add_rect(slide, Inches(0.5), ry, Inches(0.1), Inches(0.6), sc)
        _add_text_box(slide, Inches(0.75), ry, Inches(3.8), Inches(0.3),
                      r.get("label", ""), BODY_FONT, 11, DARK, bold=True)
        if r.get("detail"):
            _add_text_box(slide, Inches(0.75), Emu(ry + Inches(0.28)),
                          Inches(3.8), Inches(0.3), r["detail"],
                          BODY_FONT, 9, MID)
    # Right: Rewards
    _add_rect(slide, Inches(5.3), Inches(1.2), Inches(4.2), Inches(0.45), GREEN)
    _add_text_box(slide, Inches(5.5), Inches(1.2), Inches(3.8), Inches(0.45),
                  "REWARDS", BODY_FONT, 12, WHITE, bold=True,
                  valign=MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(rewards[:5]):
        ry = Emu(Inches(1.75) + i * Inches(0.7))
        _add_rect(slide, Inches(5.3), ry, Inches(4.2), Inches(0.6), GREEN_LIGHT)
        _add_rect(slide, Inches(5.3), ry, Inches(0.1), Inches(0.6), GREEN)
        _add_text_box(slide, Inches(5.55), ry, Inches(3.8), Inches(0.3),
                      r.get("label", ""), BODY_FONT, 11, DARK, bold=True)
        if r.get("detail"):
            _add_text_box(slide, Inches(5.55), Emu(ry + Inches(0.28)),
                          Inches(3.8), Inches(0.3), r["detail"],
                          BODY_FONT, 9, MID)


def build_appendix(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # No header bar — muted style
    _add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.5),
                  c.get("title", "Appendix"), TITLE_FONT, 18, MID, bold=True)
    _add_rect(slide, Inches(0.5), Inches(0.75), Inches(1.2), Inches(0.04), LIGHT)
    sections = c.get("sections", [])
    y_cursor = 0.9
    for i, s in enumerate(sections):
        if y_cursor > 4.8:
            break
        col = COLORS[i % len(COLORS)]
        _add_rect(slide, Inches(0.5), Inches(y_cursor), Inches(0.08),
                  Inches(0.8), col)
        _add_text_box(slide, Inches(0.7), Inches(y_cursor), Inches(8.8),
                      Inches(0.25), s.get("label", ""),
                      BODY_FONT, 9, col, bold=True)
        _add_text_box(slide, Inches(0.7), Inches(y_cursor + 0.25), Inches(8.8),
                      Inches(0.7), s.get("content", ""),
                      BODY_FONT, 8, MID, line_spacing=11)
        y_cursor += 1.05


def build_before_after(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Transformation"))
    bef = c.get("before", {})
    aft = c.get("after", {})
    interv = c.get("intervention", "")
    # Before panel
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(3.3), Inches(3.6), ORANGE_LIGHT)
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(3.3), Inches(0.1), ORANGE)
    _add_text_box(slide, Inches(0.7), Inches(1.4), Inches(2.9), Inches(0.35),
                  bef.get("label", "Before"), TITLE_FONT, 14, ORANGE, bold=True)
    _add_text_box(slide, Inches(0.7), Inches(1.85), Inches(2.9), Inches(2.8),
                  bef.get("detail", ""), BODY_FONT, 11, DARK, line_spacing=16)
    # Intervention arrow zone
    _add_oval(slide, Inches(4.15), Inches(2.5), Inches(0.5), Inches(0.5), GOLD)
    _add_text_box(slide, Inches(4.15), Inches(2.5), Inches(0.5), Inches(0.5),
                  "\u2192", BODY_FONT, 18, WHITE, bold=True, align=PP_ALIGN.CENTER,
                  valign=MSO_ANCHOR.MIDDLE)
    _add_text_box(slide, Inches(3.95), Inches(3.1), Inches(0.9), Inches(1.4),
                  interv, BODY_FONT, 9, MID, align=PP_ALIGN.CENTER, line_spacing=12)
    # After panel
    _add_rect(slide, Inches(5.0), Inches(1.2), Inches(4.5), Inches(3.6), GREEN_LIGHT)
    _add_rect(slide, Inches(5.0), Inches(1.2), Inches(4.5), Inches(0.1), GREEN)
    _add_text_box(slide, Inches(5.2), Inches(1.4), Inches(4.1), Inches(0.35),
                  aft.get("label", "After"), TITLE_FONT, 14, GREEN, bold=True)
    _add_text_box(slide, Inches(5.2), Inches(1.85), Inches(4.1), Inches(2.8),
                  aft.get("detail", ""), BODY_FONT, 11, DARK, line_spacing=16)


def build_summary(prs, c):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Summary"))
    sections = c.get("sections", [])
    n = len(sections)
    if n == 0:
        return
    gap = 0.2
    colW = (9.0 - (n - 1) * gap) / n
    for i, sec in enumerate(sections[:4]):
        x = Inches(0.5 + i * (colW + gap))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, x, Inches(1.2), Inches(colW), Inches(3.8), lt)
        _add_rect(slide, x, Inches(1.2), Inches(colW), Inches(0.1), col)
        _add_text_box(slide, Emu(x + Inches(0.12)), Inches(1.35),
                      Inches(colW - 0.24), Inches(0.35),
                      sec.get("heading", ""), TITLE_FONT, 12, col, bold=True)
        points = sec.get("points", [])
        for j, p in enumerate(points[:5]):
            py = Emu(Inches(1.8) + j * Inches(0.55))
            _add_oval(slide, Emu(x + Inches(0.12)), Emu(py + Inches(0.04)),
                      Inches(0.12), Inches(0.12), col)
            _add_text_box(slide, Emu(x + Inches(0.3)), py,
                          Inches(colW - 0.45), Inches(0.5),
                          p, BODY_FONT, 10, DARK, line_spacing=13)


def build_quote_full(prs, c):
    """Full-bleed dark bg, centered quote. No header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    # Top color strip
    _add_rect(slide, Inches(0), Inches(0), W, Inches(0.1), GREEN)
    # Gold line above quote
    _add_rect(slide, Inches(3.5), Inches(1.0), Inches(3.0), Inches(0.05), GOLD)
    _add_text_box(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(2.5),
                  c.get("quote", ""), TITLE_FONT, 26, WHITE, italic=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=36)
    # Gold line below quote
    _add_rect(slide, Inches(3.5), Inches(4.0), Inches(3.0), Inches(0.05), GOLD)
    if c.get("attribution"):
        _add_text_box(slide, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.4),
                      c["attribution"], BODY_FONT, 13, GREEN_MID, bold=True,
                      align=PP_ALIGN.CENTER)
    if c.get("context"):
        _add_text_box(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.35),
                      c["context"], BODY_FONT, 10, MID, italic=True,
                      align=PP_ALIGN.CENTER)


def build_stat_hero(prs, c):
    """One hero stat + supporting stats below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Key Metric"))
    hero = c.get("hero", {})
    _add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.5),
                  hero.get("value", "\u2014"), TITLE_FONT, 72, GREEN, bold=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if hero.get("label"):
        _add_text_box(slide, Inches(1.5), Inches(2.7), Inches(7.0), Inches(0.45),
                      hero["label"], BODY_FONT, 16, DARK, bold=True,
                      align=PP_ALIGN.CENTER)
    _add_rect(slide, Inches(2.0), Inches(3.3), Inches(6.0), Inches(0.04), LIGHT)
    supporting = c.get("supporting", [])
    n = min(len(supporting), 4)
    if n > 0:
        sw = 7.0 / n
        for i, s in enumerate(supporting[:n]):
            x = Inches(1.5 + i * sw)
            col = COLORS[i % len(COLORS)]
            lt = LIGHTS[i % len(LIGHTS)]
            _add_rect(slide, x, Inches(3.5), Inches(sw - 0.15), Inches(1.2), lt)
            _add_rect(slide, x, Inches(3.5), Inches(sw - 0.15), Inches(0.08), col)
            _add_text_box(slide, x, Inches(3.65), Inches(sw - 0.15), Inches(0.55),
                          s.get("value", ""), TITLE_FONT, 24, col, bold=True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _add_text_box(slide, x, Inches(4.2), Inches(sw - 0.15), Inches(0.35),
                          s.get("label", ""), BODY_FONT, 10, MID,
                          align=PP_ALIGN.CENTER)
    if c.get("source"):
        _add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.3),
                      c["source"], BODY_FONT, 9, MID, italic=True,
                      align=PP_ALIGN.CENTER)


def build_in_brief_featured(prs, c):
    """One featured insight at top, supporting points below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "In Brief"))
    featured = c.get("featured", "")
    supporting = c.get("supporting", [])
    # Featured block
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.3), GREEN_LIGHT)
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.1), GREEN)
    _add_text_box(slide, Inches(0.7), Inches(1.4), Inches(8.6), Inches(1.0),
                  featured, BODY_FONT, 16, DARK, bold=True, line_spacing=22,
                  valign=MSO_ANCHOR.MIDDLE)
    # Supporting points
    for i, s in enumerate(supporting[:4]):
        y = Emu(Inches(2.7) + i * Inches(0.65))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, Inches(0.5), y, Inches(9), Inches(0.55), lt)
        _add_rect(slide, Inches(0.5), y, Inches(0.1), Inches(0.55), col)
        _add_text_box(slide, Inches(0.75), y, Inches(8.5), Inches(0.55),
                      s, BODY_FONT, 12, DARK, valign=MSO_ANCHOR.MIDDLE)


def build_persona_duo(prs, c):
    """Two personas side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Archetype Comparison"))
    personas = c.get("personas", [{}, {}])

    for idx, p in enumerate(personas[:2]):
        x = Inches(0.5) if idx == 0 else Inches(5.3)
        w = Inches(4.2)
        col = COLORS[idx % len(COLORS)]
        lt = LIGHTS[idx % len(LIGHTS)]
        _add_rect(slide, x, Inches(1.2), w, Inches(3.8), lt)
        _add_rect(slide, x, Inches(1.2), w, Inches(0.1), col)
        _add_text_box(slide, Emu(x + Inches(0.15)), Inches(1.4), Emu(w - Inches(0.3)),
                      Inches(0.4), p.get("name", ""), TITLE_FONT, 18, col, bold=True)
        _add_text_box(slide, Emu(x + Inches(0.15)), Inches(1.8), Emu(w - Inches(0.3)),
                      Inches(0.3), p.get("archetype", ""), BODY_FONT, 11, GOLD, bold=True)
        traits = p.get("traits", [])
        for j, t in enumerate(traits[:4]):
            ty = Emu(Inches(2.2) + j * Inches(0.35))
            jcol = COLORS[(idx * 2 + j) % len(COLORS)]
            _add_oval(slide, Emu(x + Inches(0.15)), Emu(ty + Inches(0.06)),
                      Inches(0.12), Inches(0.12), jcol)
            _add_text_box(slide, Emu(x + Inches(0.35)), ty,
                          Emu(w - Inches(0.5)), Inches(0.3),
                          t, BODY_FONT, 10, DARK)
        if p.get("strategy"):
            _add_rect(slide, Emu(x + Inches(0.15)), Inches(3.7),
                      Emu(w - Inches(0.3)), Inches(0.04), col)
            _add_text_box(slide, Emu(x + Inches(0.15)), Inches(3.8),
                          Emu(w - Inches(0.3)), Inches(1.0),
                          p["strategy"], BODY_FONT, 10, MID, line_spacing=14)
    # Center divider
    _add_rect(slide, Inches(4.85), Inches(1.4), Inches(0.14), Inches(3.4), GREEN)


def build_process_flow_vertical(prs, c):
    """Vertical steps with down arrows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Process"))
    steps = c.get("steps", [])
    count = min(len(steps), 3)
    if count == 0:
        return
    avail = 3.8
    arrowH = 0.35
    stepH = (avail - (count - 1) * arrowH) / count
    for i, step in enumerate(steps[:count]):
        y = Inches(1.2 + i * (stepH + arrowH))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, Inches(0.5), y, Inches(9), Inches(stepH), lt)
        _add_rect(slide, Inches(0.5), y, Inches(9), Inches(0.1), col)
        _add_oval(slide, Inches(0.7), Emu(y + Inches(0.15)), Inches(0.4),
                  Inches(0.4), col)
        _add_text_box(slide, Inches(0.7), Emu(y + Inches(0.15)), Inches(0.4),
                      Inches(0.4), str(i + 1), TITLE_FONT, 13, WHITE, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text_box(slide, Inches(1.3), Emu(y + Inches(0.15)), Inches(3.0),
                      Inches(0.4), step.get("title", ""), BODY_FONT, 12, DARK,
                      bold=True, valign=MSO_ANCHOR.MIDDLE)
        if step.get("detail"):
            _add_text_box(slide, Inches(1.3), Emu(y + Inches(0.6)),
                          Inches(7.8), Inches(stepH - 0.7),
                          step["detail"], BODY_FONT, 10, MID, line_spacing=14)
        if i < count - 1:
            _add_text_box(slide, Inches(0.75), Emu(y + Inches(stepH + 0.02)),
                          Inches(0.4), Inches(0.3), "\u2193", BODY_FONT, 18,
                          col, bold=True, align=PP_ALIGN.CENTER)


def build_text_cards(prs, c):
    """Points in a 2x2 or 2x3 card grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Key Points"))
    items = c.get("items", [])
    n = min(len(items), 6)
    if n == 0:
        return
    cols = 2 if n <= 4 else 3
    rows_count = (n + cols - 1) // cols
    gap = 0.2
    cardW = (9.0 - (cols - 1) * gap) / cols
    cardH = (3.8 - (rows_count - 1) * gap) / rows_count
    for i, item in enumerate(items[:n]):
        col_idx = i % cols
        row = i // cols
        x = Inches(0.5 + col_idx * (cardW + gap))
        y = Emu(Inches(1.2) + row * (Inches(cardH) + Inches(gap)))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, x, y, Inches(cardW), Inches(cardH), lt)
        _add_rect(slide, x, y, Inches(cardW), Inches(0.08), col)
        title_text = item if isinstance(item, str) else item.get("title", "")
        detail_text = "" if isinstance(item, str) else item.get("detail", "")
        if detail_text:
            _add_text_box(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.15)),
                          Inches(cardW - 0.24), Inches(0.35),
                          title_text, BODY_FONT, 12, DARK, bold=True)
            _add_text_box(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.55)),
                          Inches(cardW - 0.24), Inches(cardH - 0.65),
                          detail_text, BODY_FONT, 10, MID, line_spacing=14)
        else:
            _add_text_box(slide, Emu(x + Inches(0.12)), Emu(y + Inches(0.15)),
                          Inches(cardW - 0.24), Inches(cardH - 0.25),
                          title_text, BODY_FONT, 12, DARK,
                          valign=MSO_ANCHOR.MIDDLE, line_spacing=17)


def build_text_columns(prs, c):
    """2-3 columns of flowing text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Overview"))
    columns = c.get("columns", [])
    n = min(len(columns), 3)
    if n == 0:
        return
    gap = 0.2
    colW = (9.0 - (n - 1) * gap) / n
    for i, col_data in enumerate(columns[:n]):
        x = Inches(0.5 + i * (colW + gap))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        heading = col_data if isinstance(col_data, str) else col_data.get("heading", "")
        body = "" if isinstance(col_data, str) else col_data.get("body", "")
        if isinstance(col_data, str):
            body = col_data
            heading = ""
        if heading:
            _add_rect(slide, x, Inches(1.2), Inches(colW), Inches(0.45), col)
            _add_text_box(slide, Emu(x + Inches(0.1)), Inches(1.2),
                          Inches(colW - 0.2), Inches(0.45),
                          heading, BODY_FONT, 12, WHITE, bold=True,
                          valign=MSO_ANCHOR.MIDDLE)
            _add_rect(slide, x, Inches(1.65), Inches(colW), Inches(3.15), lt)
            _add_text_box(slide, Emu(x + Inches(0.1)), Inches(1.75),
                          Inches(colW - 0.2), Inches(3.0),
                          body, BODY_FONT, 11, DARK, line_spacing=16)
        else:
            _add_rect(slide, x, Inches(1.2), Inches(colW), Inches(3.6), lt)
            _add_rect(slide, x, Inches(1.2), Inches(colW), Inches(0.08), col)
            _add_text_box(slide, Emu(x + Inches(0.1)), Inches(1.35),
                          Inches(colW - 0.2), Inches(3.4),
                          body, BODY_FONT, 11, DARK, line_spacing=16)


def build_text_narrative(prs, c):
    """Big lede + body paragraphs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Overview"))
    lede = c.get("lede", "")
    body = c.get("body", "")
    if isinstance(body, list):
        body = "\n\n".join(body)
    # Lede block with colored left accent
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), GREEN_LIGHT)
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(0.12), Inches(1.2), GREEN)
    _add_text_box(slide, Inches(0.8), Inches(1.25), Inches(8.5), Inches(1.1),
                  lede, TITLE_FONT, 16, DARK, bold=True, line_spacing=23,
                  valign=MSO_ANCHOR.MIDDLE)
    # Separator
    _add_rect(slide, Inches(0.5), Inches(2.55), Inches(9), Inches(0.04), LIGHT)
    # Body text
    _add_text_box(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(2.5),
                  body, BODY_FONT, 12, MID, line_spacing=18)


def build_text_nested(prs, c):
    """Colored blocks + children text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Detail"))
    items = c.get("items", [])
    n = min(len(items), 4)
    if n == 0:
        return
    avail = 3.8
    gap = 0.15
    cardH = (avail - (n - 1) * gap) / n
    for i, item in enumerate(items[:n]):
        y = Emu(Inches(1.2) + i * (Inches(cardH) + Inches(gap)))
        text = item if isinstance(item, str) else item.get("text", "")
        children = [] if isinstance(item, str) else item.get("children", [])
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        # Colored label block on left
        _add_rect(slide, Inches(0.5), y, Inches(2.2), Inches(cardH), col)
        _add_text_box(slide, Inches(0.62), y, Inches(1.96), Inches(cardH),
                      text, BODY_FONT, 11, WHITE, bold=True,
                      valign=MSO_ANCHOR.MIDDLE, line_spacing=15)
        # Children on right
        child_texts = []
        for ch in children:
            ch_text = ch if isinstance(ch, str) else ch.get("text", "")
            grandchildren = [] if isinstance(ch, str) else ch.get("children", [])
            child_texts.append(ch_text)
            for gc in grandchildren:
                gc_text = gc if isinstance(gc, str) else gc.get("text", "")
                child_texts.append("  \u2022 " + gc_text)
        body = "\n".join(child_texts)
        _add_rect(slide, Inches(2.8), y, Inches(6.7), Inches(cardH), lt)
        _add_text_box(slide, Inches(2.95), y, Inches(6.4), Inches(cardH),
                      body, BODY_FONT, 10, DARK, valign=MSO_ANCHOR.MIDDLE,
                      line_spacing=14)


def build_text_split(prs, c):
    """Left: big message. Right: evidence points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Key Point"))
    headline = c.get("headline", "")
    detail = c.get("detail", "")
    points = c.get("points", [])
    # Left panel
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(4.2), Inches(3.8), GREEN_LIGHT)
    _add_rect(slide, Inches(0.5), Inches(1.2), Inches(0.12), Inches(3.8), GREEN)
    _add_text_box(slide, Inches(0.8), Inches(1.4), Inches(3.7), Inches(1.8),
                  headline, TITLE_FONT, 18, DARK, bold=True, line_spacing=26,
                  valign=MSO_ANCHOR.MIDDLE)
    if detail:
        _add_text_box(slide, Inches(0.8), Inches(3.3), Inches(3.7), Inches(1.4),
                      detail, BODY_FONT, 11, MID, line_spacing=16)
    # Divider
    _add_rect(slide, Inches(4.85), Inches(1.4), Inches(0.06), Inches(3.4), LIGHT)
    # Right: evidence points
    for i, p in enumerate(points[:6]):
        py = Emu(Inches(1.3) + i * Inches(0.6))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        _add_rect(slide, Inches(5.1), py, Inches(4.4), Inches(0.5), lt)
        _add_rect(slide, Inches(5.1), py, Inches(0.08), Inches(0.5), col)
        _add_text_box(slide, Inches(5.3), py, Inches(4.1), Inches(0.5),
                      p, BODY_FONT, 11, DARK, valign=MSO_ANCHOR.MIDDLE,
                      line_spacing=15)


def build_text_annotated(prs, c):
    """Colored label rectangles, text on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, c.get("title", "Analysis"))
    items = c.get("items", [])
    n = min(len(items), 5)
    if n == 0:
        return
    avail = 3.8
    gap = 0.12
    rowH = min((avail - (n - 1) * gap) / n, 0.85)
    total = n * rowH + (n - 1) * gap
    startY = 1.2 + (avail - total) * 0.2
    for i, item in enumerate(items[:n]):
        y = Inches(startY + i * (rowH + gap))
        col = COLORS[i % len(COLORS)]
        lt = LIGHTS[i % len(LIGHTS)]
        # Label rectangle
        _add_rect(slide, Inches(0.5), y, Inches(1.8), Inches(rowH), col)
        _add_text_box(slide, Inches(0.5), y, Inches(1.8), Inches(rowH),
                      item.get("label", ""), BODY_FONT, 10, WHITE, bold=True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Text on right
        _add_rect(slide, Inches(2.4), y, Inches(7.1), Inches(rowH), lt)
        _add_text_box(slide, Inches(2.55), y, Inches(6.8), Inches(rowH),
                      item.get("text", ""), BODY_FONT, 11, DARK,
                      valign=MSO_ANCHOR.MIDDLE, line_spacing=16)


# ============================================================
# DISPATCH
# ============================================================
BUILDERS = {
    "title": build_title,
    "in_brief": build_in_brief,
    "section_divider": build_section_divider,
    "stat_callout": build_stat_callout,
    "quote": build_quote,
    "comparison": build_comparison,
    "text_graph": build_text_graph,
    "process_flow": build_process_flow,
    "matrix": build_matrix,
    "methods": build_methods,
    "hypotheses": build_hypotheses,
    "wsn_dense": build_wsn_dense,
    "wsn_reveal": build_wsn_reveal,
    "findings_recs": build_findings_recs,
    "findings_recs_dense": build_findings_recs_dense,
    "open_questions": build_open_questions,
    "agenda": build_agenda,
    "progressive_reveal": build_progressive_reveal,
    "closer": build_closer,
    "timeline": build_timeline,
    "data_table": build_data_table,
    "multi_stat": build_multi_stat,
    "persona": build_persona,
    "risk_tradeoff": build_risk_tradeoff,
    "appendix": build_appendix,
    "before_after": build_before_after,
    "summary": build_summary,
    "quote_full": build_quote_full,
    "stat_hero": build_stat_hero,
    "in_brief_featured": build_in_brief_featured,
    "persona_duo": build_persona_duo,
    "process_flow_vertical": build_process_flow_vertical,
    "process_flow_reveal": build_process_flow_reveal,
    "text_cards": build_text_cards,
    "text_columns": build_text_columns,
    "text_narrative": build_text_narrative,
    "text_nested": build_text_nested,
    "text_split": build_text_split,
    "text_annotated": build_text_annotated,
}


def build_deck(slide_configs, output_path):
    """Build a complete deck from a list of (slide_type, data_dict) tuples."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_type, data in slide_configs:
        if slide_type == "skip":
            continue
        builder = BUILDERS.get(slide_type)
        if builder:
            builder(prs, data)

    prs.save(output_path)
    return output_path
